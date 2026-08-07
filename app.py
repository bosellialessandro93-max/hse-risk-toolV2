from io import BytesIO
from datetime import datetime, date
import json
import html
import urllib.parse
import urllib.request
import urllib.error

import altair as alt
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Pt, Inches as DocxInches


# =========================================================
# CONFIGURAZIONE STREAMLIT
# =========================================================
st.set_page_config(
    page_title="HSE Risk Platform · V2.2",
    page_icon="🦺",
    layout="wide"
)



# =========================================================
# DESIGN SYSTEM V2.2 · RESPONSIBLE UX/UI
# =========================================================
st.markdown(
    """
    <style>
    :root {
        --hse-electric: #0AA9F4;
        --hse-electric-strong: #008ED6;
        --hse-navy: #071A2B;
        --hse-navy-2: #0D2C40;
        --hse-ink: #102A3A;
        --hse-muted: #617786;
        --hse-bg: #F5F8FB;
        --hse-surface: #FFFFFF;
        --hse-soft: #EAF7FD;
        --hse-border: #DCE6ED;
        --hse-success: #168A5B;
        --hse-warning: #B66A05;
        --hse-danger: #C53B3B;
    }

    html, body, [class*="css"] {
        font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }

    .stApp {
        background: var(--hse-bg);
        color: var(--hse-ink);
    }

    .block-container {
        max-width: 1460px;
        padding-top: 1.35rem;
        padding-bottom: 3.5rem;
    }

    /* Sidebar */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #061725 0%, #0A2638 100%);
        border-right: 1px solid rgba(255,255,255,.08);
    }
    [data-testid="stSidebar"] * { color: #EDF7FC; }
    [data-testid="stSidebar"] label { color: #EAF6FC !important; }
    [data-testid="stSidebar"] [data-baseweb="select"] * { color: var(--hse-ink) !important; }
    [data-testid="stSidebar"] input { color: var(--hse-ink) !important; }
    [data-testid="stSidebar"] div[role="radiogroup"] > label {
        border-radius: 10px;
        padding: .48rem .65rem;
        margin: .12rem 0;
        transition: background .15s ease;
    }
    [data-testid="stSidebar"] div[role="radiogroup"] > label:hover {
        background: rgba(10,169,244,.12);
    }

    /* Header / hero */
    .hse-hero {
        padding: 1.7rem 1.9rem;
        border-radius: 18px;
        background: linear-gradient(115deg, #071A2B 0%, #0D3046 70%, #0A5A82 100%);
        box-shadow: 0 14px 36px rgba(13,44,64,.13);
        margin-bottom: 1.15rem;
        position: relative;
        overflow: hidden;
        border: 1px solid rgba(255,255,255,.08);
    }
    .hse-hero:after {
        content: "";
        position: absolute;
        width: 210px; height: 210px;
        border-radius: 50%;
        right: -70px; top: -115px;
        background: rgba(10,169,244,.22);
    }
    .hse-eyebrow {
        color: #62D4FF;
        font-weight: 800;
        letter-spacing: .13em;
        font-size: .73rem;
        text-transform: uppercase;
        margin-bottom: .42rem;
    }
    .hse-hero h1 {
        color: white;
        margin: 0;
        font-size: clamp(1.75rem, 2.5vw, 2.35rem);
        letter-spacing: -.025em;
    }
    .hse-hero p {
        color: #DCECF5;
        max-width: 920px;
        margin: .65rem 0 0;
        font-size: 1rem;
        line-height: 1.55;
    }

    /* Reusable cards */
    .hse-card, .hse-panel {
        background: var(--hse-surface);
        border: 1px solid var(--hse-border);
        border-radius: 16px;
        padding: 1.15rem 1.25rem;
        box-shadow: 0 6px 20px rgba(15,51,72,.045);
    }
    .hse-card { min-height: 145px; }
    .hse-card h3, .hse-panel h3 { margin-top: 0; color: var(--hse-ink); }
    .hse-card p, .hse-panel p { color: var(--hse-muted); line-height: 1.5; }
    .hse-accent { color: var(--hse-electric-strong); font-weight: 800; }

    .hse-kicker {
        color: var(--hse-muted);
        font-size: .78rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: .08em;
    }
    .hse-section-title {
        color: var(--hse-ink);
        font-weight: 850;
        margin-top: .4rem;
        letter-spacing: -.015em;
    }
    .hse-badge, .hse-chip {
        display: inline-block;
        padding: .28rem .68rem;
        border-radius: 999px;
        background: #E8F7FE;
        color: #087BAF;
        border: 1px solid #BFE8F8;
        font-weight: 750;
        font-size: .8rem;
    }
    .hse-chip-neutral {
        background: #F1F5F8;
        color: #526A79;
        border-color: #DDE6EC;
    }

    /* Steps */
    .hse-stepper {
        display: grid;
        grid-template-columns: repeat(5, minmax(0, 1fr));
        gap: .5rem;
        margin: .85rem 0 1.15rem;
    }
    .hse-step {
        background: #FFFFFF;
        border: 1px solid var(--hse-border);
        border-radius: 12px;
        padding: .72rem .8rem;
        color: #526B7A;
        font-size: .83rem;
        font-weight: 700;
    }
    .hse-step span {
        display:block;
        color: var(--hse-electric-strong);
        font-size: .72rem;
        letter-spacing:.08em;
        margin-bottom:.18rem;
    }

    /* Metrics */
    div[data-testid="stMetric"] {
        background: #FFFFFF;
        border: 1px solid var(--hse-border);
        border-radius: 14px;
        padding: .9rem 1rem;
        box-shadow: 0 5px 16px rgba(15,51,72,.04);
    }
    div[data-testid="stMetric"] label { color: var(--hse-muted) !important; }

    /* Inputs */
    div[data-baseweb="input"], div[data-baseweb="select"] > div,
    div[data-testid="stNumberInput"] > div > div {
        border-radius: 10px !important;
    }
    div[data-testid="stExpander"] {
        background: #FFFFFF;
        border: 1px solid var(--hse-border) !important;
        border-radius: 14px !important;
        overflow: hidden;
    }

    /* Buttons */
    .stButton > button,
    .stDownloadButton > button {
        border-radius: 10px !important;
        min-height: 43px;
        font-weight: 750 !important;
        border: 1px solid #BFD8E6 !important;
        transition: all .14s ease-in-out;
        box-shadow: none;
    }
    .stButton > button:hover,
    .stDownloadButton > button:hover {
        transform: translateY(-1px);
        border-color: #75CFF2 !important;
        box-shadow: 0 6px 16px rgba(10,120,175,.10);
    }
    .stButton > button[kind="primary"] {
        background: linear-gradient(100deg, #007EBE 0%, #0AA9F4 100%) !important;
        color: white !important;
        border: none !important;
    }
    .stDownloadButton > button {
        background: #FFFFFF !important;
        color: #087BAF !important;
    }

    /* Tabs */
    div[data-baseweb="tab-list"] {
        gap: .25rem;
        background: #EDF4F8;
        padding: .3rem;
        border-radius: 12px;
        border: 1px solid #DDE8EE;
    }
    button[data-baseweb="tab"] {
        border-radius: 9px;
        padding-left: .8rem !important;
        padding-right: .8rem !important;
    }

    /* Risk result */
    .risk-result {
        background: #FFFFFF;
        border: 1px solid var(--hse-border);
        border-radius: 18px;
        padding: 1.35rem 1.45rem;
        box-shadow: 0 8px 24px rgba(15,51,72,.055);
        margin: .5rem 0 1rem;
    }
    .risk-result-grid {
        display: grid;
        grid-template-columns: 170px 1fr;
        gap: 1.3rem;
        align-items: center;
    }
    .risk-number {
        font-size: 3.55rem;
        line-height: 1;
        font-weight: 900;
        letter-spacing: -.05em;
        color: var(--risk-color, #0A3551);
    }
    .risk-denom { color: #7A8E9A; font-size: 1rem; font-weight: 650; }
    .risk-level {
        display:inline-block;
        margin-top:.48rem;
        padding:.3rem .72rem;
        border-radius:999px;
        font-size:.78rem;
        font-weight:850;
        letter-spacing:.05em;
        color: var(--risk-color, #0A3551);
        background: color-mix(in srgb, var(--risk-color, #0A3551) 10%, white);
        border:1px solid color-mix(in srgb, var(--risk-color, #0A3551) 28%, white);
    }
    .risk-bar {
        height: 10px;
        border-radius: 999px;
        background: linear-gradient(90deg,
            #24956A 0% 20%,
            #8AAE37 20% 40%,
            #D29B27 40% 60%,
            #D46A39 60% 80%,
            #C53B3B 80% 100%);
        margin: .75rem 0 .45rem;
        position:relative;
    }
    .risk-marker {
        position:absolute;
        left: var(--risk-position, 0%);
        top:50%; transform:translate(-50%, -50%);
        width:18px; height:18px;
        border-radius:50%;
        background:#FFFFFF;
        border:4px solid var(--risk-color, #0A3551);
        box-shadow:0 2px 7px rgba(0,0,0,.18);
    }

    .driver-list {
        margin: .35rem 0 0;
        padding:0;
        list-style:none;
    }
    .driver-list li {
        display:flex;
        justify-content:space-between;
        gap:1rem;
        padding:.52rem 0;
        border-bottom:1px solid #EDF2F5;
        color:#334F60;
    }
    .driver-list li:last-child { border-bottom:none; }

    /* Timeline / attention cards */
    .attention-item, .timeline-item {
        background:#FFFFFF;
        border:1px solid var(--hse-border);
        border-left:4px solid var(--status-color, #0AA9F4);
        border-radius:12px;
        padding:.82rem .95rem;
        margin:.5rem 0;
    }
    .attention-title { font-weight:800; color:var(--hse-ink); }
    .attention-meta { color:var(--hse-muted); font-size:.86rem; margin-top:.2rem; }

    .trust-strip {
        display:flex;
        gap:.55rem;
        flex-wrap:wrap;
        margin:.55rem 0 1.25rem;
    }
    .trust-item {
        background:#FFFFFF;
        border:1px solid var(--hse-border);
        color:#496272;
        border-radius:999px;
        padding:.34rem .7rem;
        font-size:.79rem;
        font-weight:700;
    }

    .responsible-note {
        border-left:4px solid #0AA9F4;
        background:#EEF8FC;
        padding:.8rem 1rem;
        border-radius:10px;
        color:#35586B;
        font-size:.9rem;
        margin:.7rem 0;
    }

    @media (max-width: 900px) {
        .hse-stepper { grid-template-columns: 1fr 1fr; }
        .risk-result-grid { grid-template-columns: 1fr; }
        .block-container { padding-left: .8rem; padding-right: .8rem; }
    }
    </style>
    """,
    unsafe_allow_html=True
)

# =========================================================
# ARCHIVIAZIONE SUPABASE (FREE TIER READY)
# =========================================================
def get_supabase_config():
    """Legge le credenziali da Streamlit Secrets senza esporle nel codice."""
    try:
        url = st.secrets.get("SUPABASE_URL", "")
        key = st.secrets.get("SUPABASE_KEY", "")
    except Exception:
        url, key = "", ""
    return str(url).rstrip("/"), str(key)


def database_ready():
    url, key = get_supabase_config()
    return bool(url and key)


def get_supabase_secret_key():
    """Chiave privilegiata opzionale, usata solo lato server per le funzioni Admin."""
    try:
        return str(st.secrets.get("SUPABASE_SECRET_KEY", ""))
    except Exception:
        return ""


def _supabase_request(method, path, payload=None, query=None, privileged=False):
    url, public_key = get_supabase_config()
    key = get_supabase_secret_key() if privileged else public_key
    if not url or not key:
        if privileged:
            raise RuntimeError("Funzioni Admin non configurate: aggiungi SUPABASE_SECRET_KEY nei Secrets.")
        raise RuntimeError("Database non configurato: aggiungi SUPABASE_URL e SUPABASE_KEY nei Secrets.")

    endpoint = f"{url}/rest/v1/{path}"
    if query:
        endpoint += "?" + urllib.parse.urlencode(query, safe="*,.()")

    body = None
    if payload is not None:
        body = json.dumps(payload, ensure_ascii=False, default=str).encode("utf-8")

    req = urllib.request.Request(endpoint, data=body, method=method)
    req.add_header("apikey", key)

    # Supabase distingue tra le nuove API key opache (sb_publishable_... /
    # sb_secret_...) e le vecchie chiavi JWT (anon / service_role).
    # Le nuove chiavi sb_* NON devono essere inviate come Bearer token:
    # basta l'header apikey. Per compatibilità con eventuali chiavi legacy
    # JWT continuiamo invece ad aggiungere Authorization: Bearer <key>.
    if not str(key).startswith("sb_"):
        req.add_header("Authorization", f"Bearer {key}")

    req.add_header("Content-Type", "application/json")
    req.add_header("Accept", "application/json")
    if method == "POST":
        req.add_header("Prefer", "return=representation")
    elif method == "DELETE":
        req.add_header("Prefer", "return=representation")

    try:
        with urllib.request.urlopen(req, timeout=15) as response:
            raw = response.read().decode("utf-8")
            return json.loads(raw) if raw else []
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(f"Errore database ({exc.code}): {detail}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"Connessione database non disponibile: {exc.reason}") from exc


def build_storage_payload(report, created_by=""):
    data = report["data"]
    return {
        "project_name": data.get("nome", ""),
        "phase": data.get("fase", ""),
        "risk": round(float(report["risk"]), 2),
        "level": report["level"],
        "created_by": created_by or data.get("created_by", ""),
        "payload": {
            "data": data,
            "technical_data": report["technical_data"],
            "summary": report["summary"],
            "root_cause": report["root_cause"],
            "drivers": report["driver_df"].to_dict(orient="records"),
            "actions": report["actions_df"].to_dict(orient="records"),
            "action_plan": report["action_plan_df"].to_dict(orient="records"),
            "scorecard": report["scorecard_df"].to_dict(orient="records"),
            "scorecard_overall": report["scorecard_overall"],
        }
    }


def save_assessment(report, created_by=""):
    return _supabase_request(
        "POST",
        "risk_assessments",
        payload=build_storage_payload(report, created_by)
    )


def fetch_assessments(limit=5000):
    if not database_ready():
        return []
    return _supabase_request(
        "GET",
        "risk_assessments",
        query={
            "select": "id,created_at,project_name,phase,risk,level,created_by,payload",
            "order": "created_at.desc",
            "limit": str(limit),
        }
    )


def delete_assessment(assessment_id):
    """Elimina una valutazione usando esclusivamente la secret key server-side."""
    return _supabase_request(
        "DELETE",
        "risk_assessments",
        query={"id": f"eq.{int(assessment_id)}"},
        privileged=True,
    )


def get_admin_password():
    try:
        return str(st.secrets.get("ADMIN_PASSWORD", ""))
    except Exception:
        return ""


def admin_ready():
    return bool(get_admin_password() and get_supabase_secret_key())


def assessment_code(row):
    """Codice leggibile derivato da cantiere, data e ID senza modificare lo schema DB."""
    site = str(row.get("Cantiere", "HSE")).upper().replace(" ", "-")
    site = ''.join(ch for ch in site if ch.isalnum() or ch == '-')[:18] or "HSE"
    dt = row.get("Data")
    try:
        stamp = pd.Timestamp(dt).strftime("%Y%m%d")
    except Exception:
        stamp = datetime.now().strftime("%Y%m%d")
    try:
        ident = int(row.get("ID"))
    except Exception:
        ident = 0
    return f"{site}-{stamp}-{ident:04d}"


def enrich_assessments_dataframe(df):
    if df.empty:
        return df
    out = df.copy()
    out["Codice"] = out.apply(lambda row: assessment_code(row), axis=1)
    return out


def assessments_dataframe(rows):
    records = []
    for row in rows:
        payload = row.get("payload") or {}
        data = payload.get("data") or {}
        technical = payload.get("technical_data") or {}
        records.append({
            "ID": row.get("id"),
            "Data": row.get("created_at"),
            "Cantiere": row.get("project_name") or data.get("nome", ""),
            "Fase": row.get("phase") or data.get("fase", ""),
            "Risk Index": row.get("risk"),
            "Livello": row.get("level"),
            "Compilato da": row.get("created_by", ""),
            "Ispezioni": data.get("inspections", 0),
            "NC": data.get("num_nc", 0),
            "Stop Work": data.get("stopworks", 0),
            "Criticità aperte": data.get("crit_open", 0),
            "Appaltatori": data.get("app", 0),
            "Subappaltatori": data.get("sub", 0),
            "HSE imprese": data.get("company_hse", 0),
            "Interferenza": data.get("interference_level", 0),
            "Sensibilizzazioni": data.get("awareness_count", 0),
            "Indice detection": technical.get("nc_detection_ratio", 0),
            "Attività": ", ".join(data.get("activities", [])),
        })
    df = pd.DataFrame(records)
    if not df.empty:
        df["Data"] = pd.to_datetime(df["Data"], errors="coerce", utc=True).dt.tz_convert(None)
        df["Risk Index"] = pd.to_numeric(df["Risk Index"], errors="coerce")
    return df


def render_db_notice():
    st.info(
        "Il database non è ancora collegato. Il calcolo resta disponibile, "
        "mentre storico, dashboard ed export centralizzati richiedono Supabase."
    )


def risk_level_color(level):
    return {
        "Basso": "#24956A",
        "Medio basso": "#7B9E30",
        "Medio": "#C58B18",
        "Alto": "#D36132",
        "Critico": "#C53B3B",
    }.get(str(level), "#0A7FB9")


def risk_level_rank(level):
    return {"Basso": 1, "Medio basso": 2, "Medio": 3, "Alto": 4, "Critico": 5}.get(str(level), 0)


def safe_html(value):
    return html.escape(str(value if value is not None else ""))


def navigate_to(page_name):
    """Callback sicura: eseguita prima del rerun del widget di navigazione."""
    st.session_state["main_navigation"] = page_name


def render_hero(title, subtitle, eyebrow="HSE RISK PLATFORM · V2.2"):
    st.markdown(
        f"""
        <div class="hse-hero">
            <div class="hse-eyebrow">{safe_html(eyebrow)}</div>
            <h1>{safe_html(title)}</h1>
            <p>{safe_html(subtitle)}</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_trust_strip(database_status=None):
    db_text = database_status or ("Database connected" if database_ready() else "Database not connected")
    st.markdown(
        f"""
        <div class="trust-strip">
            <span class="trust-item">✓ Deterministic risk model</span>
            <span class="trust-item">✓ Explainable drivers</span>
            <span class="trust-item">✓ Traceable assessments</span>
            <span class="trust-item">✓ {safe_html(db_text)}</span>
            <span class="trust-item">Model V2.2</span>
        </div>
        """,
        unsafe_allow_html=True,
    )


def calculate_data_completeness(data):
    """Misura prudenziale di completezza, non una certificazione della qualità del dato."""
    score = 100
    notes = []
    if not str(data.get("nome", "")).strip():
        score -= 25
        notes.append("nome cantiere non valorizzato")
    if not str(data.get("created_by", "")).strip():
        score -= 10
        notes.append("compilatore non indicato")
    if not data.get("activities"):
        score -= 20
        notes.append("nessuna attività selezionata")
    num_nc = int(data.get("num_nc", 0) or 0)
    detailed_nc = len(data.get("nc_items") or [])
    if num_nc > 0 and detailed_nc < num_nc:
        missing_ratio = (num_nc - detailed_nc) / max(num_nc, 1)
        deduction = min(20, round(20 * missing_ratio))
        score -= deduction
        notes.append(f"{num_nc - detailed_nc} NC non dettagliate")
    awareness_count = int(data.get("awareness_count", 0) or 0)
    if awareness_count > 0 and not data.get("awareness_types"):
        score -= 15
        notes.append("sensibilizzazioni senza tipologia")
    score = max(0, min(100, int(score)))
    if score >= 85:
        label = "Buona"
    elif score >= 65:
        label = "Parziale"
    else:
        label = "Limitata"
    return score, label, notes


def _latest_with_delta(df):
    if df.empty:
        return df
    ordered = df.sort_values(["Cantiere", "Data"]).copy()
    ordered["Risk precedente"] = ordered.groupby("Cantiere")["Risk Index"].shift(1)
    ordered["Delta"] = ordered["Risk Index"] - ordered["Risk precedente"]
    return ordered.groupby("Cantiere", as_index=False).tail(1)


def _previous_site_risk(site_name):
    if not database_ready() or not str(site_name).strip():
        return None
    try:
        df, _ = _load_archive_df()
        subset = df[df["Cantiere"].astype(str).str.casefold() == str(site_name).strip().casefold()]
        if subset.empty:
            return None
        return float(subset.sort_values("Data").iloc[-1]["Risk Index"])
    except Exception:
        return None


def render_home():
    render_hero(
        "Operational HSE risk, made visible.",
        "Un unico punto di accesso per valutare il rischio, leggere i trend, consultare lo storico e costruire dataset di reporting.",
        eyebrow="HSE RISK PLATFORM · OVERVIEW"
    )

    df = pd.DataFrame()
    db_error = None
    if database_ready():
        try:
            rows = fetch_assessments(limit=10000)
            df = enrich_assessments_dataframe(assessments_dataframe(rows))
        except Exception as exc:
            db_error = str(exc)

    render_trust_strip("Database connected" if database_ready() and not db_error else "Database unavailable" if db_error else "Database not connected")

    cta1, cta2, cta3 = st.columns([1.25, 1, 1])
    with cta1:
        st.button(
            "＋ Nuova valutazione",
            type="primary",
            use_container_width=True,
            on_click=navigate_to,
            args=("🧮 Risk Assessment",),
            key="home_new_assessment",
        )
    with cta2:
        st.button(
            "Apri Dashboard",
            use_container_width=True,
            on_click=navigate_to,
            args=("📊 Dashboard & Storico",),
            key="home_dashboard",
        )
    with cta3:
        st.button(
            "Estrai dati",
            use_container_width=True,
            on_click=navigate_to,
            args=("📥 Estrazione dati",),
            key="home_export",
        )

    st.markdown("### Portfolio overview")
    if db_error:
        st.warning(f"Database configurato ma non raggiungibile: {db_error}")
    if not df.empty:
        latest = _latest_with_delta(df)
        k1, k2, k3, k4 = st.columns(4)
        k1.metric("Cantieri monitorati", int(df["Cantiere"].nunique()))
        k2.metric("Valutazioni", len(df))
        k3.metric("Risk medio portfolio", f"{latest['Risk Index'].mean():.1f}")
        k4.metric("Siti Alto / Critico", int(latest["Livello"].isin(["Alto", "Critico"]).sum()))

        attention = latest.sort_values(["Risk Index", "Data"], ascending=[False, False]).head(4)
        st.markdown("### Attention required")
        if attention.empty:
            st.info("Nessun dato disponibile.")
        else:
            for _, row in attention.iterrows():
                color = risk_level_color(row["Livello"])
                delta = row.get("Delta")
                if pd.isna(delta):
                    delta_text = "prima valutazione"
                elif delta > 0:
                    delta_text = f"↑ +{delta:.1f} vs precedente"
                elif delta < 0:
                    delta_text = f"↓ {delta:.1f} vs precedente"
                else:
                    delta_text = "→ invariato"
                when = pd.Timestamp(row["Data"]).strftime("%d/%m/%Y") if pd.notna(row["Data"]) else "—"
                st.markdown(
                    f"""
                    <div class="attention-item" style="--status-color:{color}">
                        <div class="attention-title">{safe_html(row['Cantiere'])} · Risk {float(row['Risk Index']):.0f} · {safe_html(row['Livello'])}</div>
                        <div class="attention-meta">{when} · {safe_html(delta_text)} · fase {safe_html(row['Fase'])}</div>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
    else:
        k1, k2, k3, k4 = st.columns(4)
        k1.metric("Cantieri monitorati", 0)
        k2.metric("Valutazioni", 0)
        k3.metric("Risk medio portfolio", "—")
        k4.metric("Siti Alto / Critico", 0)
        if not database_ready():
            render_db_notice()
        else:
            st.info("Il database è connesso. Salva la prima valutazione per popolare la dashboard.")

    st.markdown("### Workspace")
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown(
            '<div class="hse-card"><div class="hse-kicker">ASSESS</div><h3>Risk Assessment</h3><p>Workflow guidato, risultato spiegabile, driver, action plan e salvataggio controllato.</p><span class="hse-badge">Decision support</span></div>',
            unsafe_allow_html=True,
        )
    with c2:
        st.markdown(
            '<div class="hse-card"><div class="hse-kicker">MONITOR</div><h3>Dashboard & History</h3><p>Trend, confronto tra cantieri, ultima valutazione e dettaglio storico tracciabile.</p><span class="hse-badge">Portfolio view</span></div>',
            unsafe_allow_html=True,
        )
    with c3:
        st.markdown(
            '<div class="hse-card"><div class="hse-kicker">REPORT</div><h3>Data Export</h3><p>Filtri combinabili, anteprima del perimetro e download Excel/CSV per reporting.</p><span class="hse-badge">Controlled export</span></div>',
            unsafe_allow_html=True,
        )

    st.markdown(
        '<div class="responsible-note"><b>Responsible use.</b> Il Risk Index supporta la prioritizzazione e la lettura dei segnali HSE. Non sostituisce la valutazione dei rischi, le procedure applicabili, il giudizio professionale o le responsabilità previste dalla normativa.</div>',
        unsafe_allow_html=True,
    )


def _load_archive_df():
    rows = fetch_assessments()
    df = assessments_dataframe(rows)
    return enrich_assessments_dataframe(df), rows


def _payload_for_id(rows, assessment_id):
    for row in rows:
        if str(row.get("id")) == str(assessment_id):
            return row.get("payload") or {}
    return {}


def _filter_by_period(df, period):
    if df.empty or not isinstance(period, (tuple, list)) or len(period) != 2:
        return df
    start, end = period
    return df[(df["Data"].dt.date >= start) & (df["Data"].dt.date <= end)]


def render_assessment_detail(payload, archive_row=None):
    data = payload.get("data") or {}
    technical = payload.get("technical_data") or {}
    drivers = pd.DataFrame(payload.get("drivers") or [])
    actions = pd.DataFrame(payload.get("actions") or [])
    action_plan = pd.DataFrame(payload.get("action_plan") or [])
    scorecard = pd.DataFrame(payload.get("scorecard") or [])

    if archive_row is not None:
        st.markdown(f"### {archive_row.get('Codice', '')} · {archive_row.get('Cantiere', '')}")
        a, b, c, d = st.columns(4)
        a.metric("Risk Index", f"{float(archive_row.get('Risk Index', 0)):.1f} / 100")
        b.metric("Livello", str(archive_row.get("Livello", "")))
        c.metric("Fase", str(archive_row.get("Fase", "")))
        d.metric("Compilato da", str(archive_row.get("Compilato da", "")) or "—")

    t1, t2, t3, t4 = st.tabs(["Input", "Risk drivers", "Actions", "Decision output"])
    with t1:
        input_rows = [
            ("Ispezioni", data.get("inspections", 0)),
            ("NC", data.get("num_nc", 0)),
            ("Stop Work", data.get("stopworks", 0)),
            ("Criticità aperte", data.get("crit_open", 0)),
            ("Criticità chiuse nei tempi", data.get("crit_ontime", 0)),
            ("Criticità chiuse in ritardo", data.get("crit_late", 0)),
            ("Appaltatori", data.get("app", 0)),
            ("Subappaltatori", data.get("sub", 0)),
            ("HSE imprese", data.get("company_hse", 0)),
            ("Interferenza", data.get("interference_level", 0)),
            ("Sensibilizzazioni", data.get("awareness_count", 0)),
            ("Attività", ", ".join(data.get("activities", []))),
            ("Tipi sensibilizzazione", ", ".join(data.get("awareness_types", []))),
            ("Detection", technical.get("nc_detection_ratio", 0)),
        ]
        st.dataframe(pd.DataFrame(input_rows, columns=["Indicatore", "Valore"]), use_container_width=True, hide_index=True)
        nc_items = pd.DataFrame(data.get("nc_items") or [])
        if not nc_items.empty:
            st.markdown("#### Non conformità")
            st.dataframe(nc_items, use_container_width=True, hide_index=True)
    with t2:
        if not drivers.empty:
            st.altair_chart(create_driver_chart(drivers[["Driver", "Valore"]]), use_container_width=True)
            st.dataframe(drivers, use_container_width=True, hide_index=True)
        else:
            st.info("Driver non disponibili per questa valutazione.")
    with t3:
        if not action_plan.empty:
            st.markdown("#### Action Plan")
            st.dataframe(action_plan, use_container_width=True, hide_index=True)
        if not actions.empty:
            st.markdown("#### Recommended actions")
            st.dataframe(actions, use_container_width=True, hide_index=True)
        if not scorecard.empty:
            st.markdown("#### HSE Scorecard")
            st.dataframe(scorecard, use_container_width=True, hide_index=True)
    with t4:
        if payload.get("summary"):
            st.markdown("#### Executive Summary")
            st.markdown(payload.get("summary"))
        if payload.get("root_cause"):
            with st.expander("Root Cause Analysis", expanded=False):
                st.markdown(payload.get("root_cause"))


def render_dashboard():
    render_hero(
        "Portfolio risk overview",
        "Individua dove intervenire, verifica l'andamento dei cantieri e apri il dettaglio tracciabile delle valutazioni storiche.",
        eyebrow="DASHBOARD & HISTORY · V2.2"
    )
    if not database_ready():
        render_db_notice()
        return
    try:
        df, rows = _load_archive_df()
    except Exception as exc:
        st.error(str(exc))
        return
    if df.empty:
        st.info("Nessuna valutazione ancora archiviata.")
        return

    last_update = df["Data"].max()
    if pd.notna(last_update):
        st.caption(f"Data updated · {pd.Timestamp(last_update).strftime('%d/%m/%Y %H:%M')} · Model V2.2")

    min_date = df["Data"].min().date()
    max_date = df["Data"].max().date()
    with st.expander("Filters", expanded=True):
        f1, f2, f3 = st.columns([1.25, 1, 1])
        with f1:
            period = st.date_input("Periodo", value=(min_date, max_date), min_value=min_date, max_value=max_date, key="dash_period")
        with f2:
            site_options = ["Tutti"] + sorted(df["Cantiere"].dropna().astype(str).unique().tolist())
            selected_site = st.selectbox("Cantiere", site_options, key="dash_site")
        with f3:
            level_options = ["Tutti", "Basso", "Medio basso", "Medio", "Alto", "Critico"]
            selected_level = st.selectbox("Risk level", level_options, key="dash_level")

    view = _filter_by_period(df, period)
    if selected_site != "Tutti":
        view = view[view["Cantiere"] == selected_site]
    if selected_level != "Tutti":
        view = view[view["Livello"] == selected_level]

    if view.empty:
        st.warning("Nessuna valutazione corrisponde ai filtri selezionati.")
        return

    latest = _latest_with_delta(view)
    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Cantieri", int(view["Cantiere"].nunique()))
    k2.metric("Risk medio attuale", f"{latest['Risk Index'].mean():.1f}")
    k3.metric("Siti Alto / Critico", int(latest["Livello"].isin(["Alto", "Critico"]).sum()))
    k4.metric("Valutazioni nel periodo", len(view))

    st.markdown("### Attention required")
    attention = latest.sort_values(["Risk Index", "Data"], ascending=[False, False]).head(5)
    for _, row in attention.iterrows():
        color = risk_level_color(row["Livello"])
        delta = row.get("Delta")
        if pd.isna(delta):
            delta_text = "prima valutazione nel perimetro"
        elif delta > 0:
            delta_text = f"↑ +{delta:.1f}"
        elif delta < 0:
            delta_text = f"↓ {delta:.1f}"
        else:
            delta_text = "→ 0.0"
        st.markdown(
            f"""
            <div class="attention-item" style="--status-color:{color}">
                <div class="attention-title">{safe_html(row['Cantiere'])} · Risk {float(row['Risk Index']):.0f} · {safe_html(row['Livello'])}</div>
                <div class="attention-meta">Trend {safe_html(delta_text)} · NC {int(row.get('NC',0) or 0)} · Interferenza {int(row.get('Interferenza',0) or 0)}/10 · {safe_html(row.get('Fase',''))}</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    left, right = st.columns([1.65, 1])
    with left:
        st.markdown("### Risk evolution")
        trend = (
            alt.Chart(view.dropna(subset=["Data", "Risk Index"]))
            .mark_line(point=True, strokeWidth=2.7)
            .encode(
                x=alt.X("Data:T", title="Data"),
                y=alt.Y("Risk Index:Q", scale=alt.Scale(domain=[0, 100]), title="Risk Index"),
                color=alt.Color("Cantiere:N", legend=alt.Legend(title="Cantiere")),
                tooltip=["Codice:N", "Cantiere:N", alt.Tooltip("Data:T"), alt.Tooltip("Risk Index:Q", format=".1f"), "Livello:N"]
            ).properties(height=360)
        )
        st.altair_chart(trend, use_container_width=True)
    with right:
        st.markdown("### Current risk distribution")
        level_counts = latest.groupby("Livello", as_index=False).size().rename(columns={"size": "Cantieri"})
        level_chart = (
            alt.Chart(level_counts)
            .mark_bar(cornerRadiusTopLeft=5, cornerRadiusTopRight=5)
            .encode(
                x=alt.X("Livello:N", sort=["Basso", "Medio basso", "Medio", "Alto", "Critico"], title=None),
                y=alt.Y("Cantieri:Q", title="Cantieri"),
                color=alt.Color(
                    "Livello:N",
                    scale=alt.Scale(
                        domain=["Basso", "Medio basso", "Medio", "Alto", "Critico"],
                        range=["#24956A", "#7B9E30", "#C58B18", "#D36132", "#C53B3B"],
                    ),
                    legend=None,
                ),
                tooltip=["Livello:N", "Cantieri:Q"]
            ).properties(height=360)
        )
        st.altair_chart(level_chart, use_container_width=True)

    st.markdown("### Operational signal")
    indicator = st.selectbox(
        "Indicatore da analizzare",
        ["NC", "Ispezioni", "Interferenza", "Criticità aperte", "Sensibilizzazioni", "Indice detection"],
        key="dash_indicator",
    )
    indicator_chart = (
        alt.Chart(view.dropna(subset=["Data"]))
        .mark_line(point=True, strokeWidth=2.25)
        .encode(
            x=alt.X("Data:T", title="Data"),
            y=alt.Y(f"{indicator}:Q", title=indicator),
            color=alt.Color("Cantiere:N", legend=alt.Legend(title="Cantiere")),
            tooltip=["Codice:N", "Cantiere:N", alt.Tooltip("Data:T"), alt.Tooltip(f"{indicator}:Q", format=".2f")]
        ).properties(height=285)
    )
    st.altair_chart(indicator_chart, use_container_width=True)

    st.markdown("### Current project status")
    current_table = latest[["Codice", "Cantiere", "Data", "Risk Index", "Livello", "Delta", "NC", "Ispezioni", "Interferenza", "Compilato da"]].copy()
    current_table = current_table.sort_values("Risk Index", ascending=False)
    st.dataframe(current_table, use_container_width=True, hide_index=True)

    st.markdown("### Assessment history")
    history = view.sort_values("Data", ascending=False)
    h1, h2 = st.columns([1, 1.7])
    with h1:
        timeline_site = st.selectbox("Timeline cantiere", sorted(history["Cantiere"].dropna().astype(str).unique()), key="timeline_site")
        site_history = history[history["Cantiere"] == timeline_site].sort_values("Data", ascending=False).head(8)
        for _, row in site_history.iterrows():
            color = risk_level_color(row["Livello"])
            when = pd.Timestamp(row["Data"]).strftime("%d %b %Y") if pd.notna(row["Data"]) else "—"
            st.markdown(
                f"""
                <div class="timeline-item" style="--status-color:{color}">
                    <div class="attention-title">{when} · Risk {float(row['Risk Index']):.0f}</div>
                    <div class="attention-meta">{safe_html(row['Livello'])} · {safe_html(row['Fase'])} · NC {int(row.get('NC',0) or 0)}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )
    with h2:
        st.dataframe(history[["Codice", "Cantiere", "Data", "Fase", "Risk Index", "Livello", "Compilato da"]], use_container_width=True, hide_index=True)
        options = {f"{row['Codice']} · {row['Cantiere']} · Risk {float(row['Risk Index']):.1f}": row for _, row in history.iterrows()}
        selected_label = st.selectbox("Open assessment detail", ["— Seleziona —"] + list(options.keys()), key="history_detail")
        if selected_label != "— Seleziona —":
            row = options[selected_label]
            payload = _payload_for_id(rows, row["ID"])
            render_assessment_detail(payload, row.to_dict())


def render_extraction():
    render_hero(
        "Build your dataset",
        "Definisci il perimetro, verifica quali record saranno inclusi e scarica solo i dati necessari al reporting.",
        eyebrow="DATA EXPORT · V2.2"
    )
    if not database_ready():
        render_db_notice()
        return
    try:
        df, _rows = _load_archive_df()
    except Exception as exc:
        st.error(str(exc))
        return
    if df.empty:
        st.info("Nessun dato disponibile per l'estrazione.")
        return

    min_date = df["Data"].min().date()
    max_date = df["Data"].max().date()
    with st.expander("Filters · lascia vuoto per includere tutti i valori", expanded=True):
        f1, f2, f3 = st.columns(3)
        with f1:
            period = st.date_input("Periodo", value=(min_date, max_date), min_value=min_date, max_value=max_date, key="extract_period")
        with f2:
            sites = st.multiselect("Cantiere", sorted(df["Cantiere"].dropna().astype(str).unique()), placeholder="Tutti i cantieri", key="extract_sites")
        with f3:
            phases = st.multiselect("Fase", sorted(df["Fase"].dropna().astype(str).unique()), placeholder="Tutte le fasi", key="extract_phases")

        f4, f5, f6 = st.columns(3)
        with f4:
            levels = st.multiselect("Risk level", ["Basso", "Medio basso", "Medio", "Alto", "Critico"], placeholder="Tutti i livelli", key="extract_levels")
        with f5:
            authors = st.multiselect("Compilato da", sorted([x for x in df["Compilato da"].dropna().astype(str).unique() if x.strip()]), placeholder="Tutti", key="extract_authors")
        with f6:
            risk_range = st.slider("Risk Index", 0, 100, (0, 100), key="extract_risk")

        f7, f8 = st.columns(2)
        with f7:
            activity_text = st.text_input("Attività contiene", placeholder="es. scavi, elettrici", key="extract_activity")
        with f8:
            min_nc = st.number_input("NC minime", min_value=0, value=0, step=1, key="extract_min_nc")

    filtered = _filter_by_period(df.copy(), period)
    if sites:
        filtered = filtered[filtered["Cantiere"].isin(sites)]
    if phases:
        filtered = filtered[filtered["Fase"].isin(phases)]
    if levels:
        filtered = filtered[filtered["Livello"].isin(levels)]
    if authors:
        filtered = filtered[filtered["Compilato da"].isin(authors)]
    filtered = filtered[filtered["Risk Index"].between(risk_range[0], risk_range[1])]
    filtered = filtered[pd.to_numeric(filtered["NC"], errors="coerce").fillna(0) >= min_nc]
    if activity_text.strip():
        filtered = filtered[filtered["Attività"].fillna("").str.contains(activity_text.strip(), case=False, regex=False)]

    st.markdown("### Selection summary")
    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Valutazioni incluse", len(filtered))
    k2.metric("Cantieri inclusi", int(filtered["Cantiere"].nunique()) if not filtered.empty else 0)
    k3.metric("Risk medio", f"{filtered['Risk Index'].mean():.1f}" if not filtered.empty else "—")
    k4.metric("Alto / Critico", int(filtered["Livello"].isin(["Alto", "Critico"]).sum()) if not filtered.empty else 0)

    st.markdown("### Preview")
    st.caption("L'anteprima riflette esattamente il perimetro che verrà esportato.")
    st.dataframe(filtered.sort_values("Data", ascending=False), use_container_width=True, hide_index=True)
    if filtered.empty:
        st.warning("Nessun record corrisponde ai filtri selezionati.")
        return

    export_df = filtered.copy()
    export_df["Data"] = export_df["Data"].dt.strftime("%Y-%m-%d %H:%M:%S")
    csv_data = export_df.to_csv(index=False).encode("utf-8-sig")
    excel_buffer = BytesIO()
    with pd.ExcelWriter(excel_buffer, engine="openpyxl") as writer:
        export_df.to_excel(writer, sheet_name="Valutazioni", index=False)
        summary = pd.DataFrame([
            {"Indicatore": "Valutazioni", "Valore": len(filtered)},
            {"Indicatore": "Cantieri", "Valore": int(filtered["Cantiere"].nunique())},
            {"Indicatore": "Risk medio", "Valore": round(float(filtered["Risk Index"].mean()), 2)},
            {"Indicatore": "Risk massimo", "Valore": round(float(filtered["Risk Index"].max()), 2)},
            {"Indicatore": "Alto/Critico", "Valore": int(filtered["Livello"].isin(["Alto", "Critico"]).sum())},
            {"Indicatore": "Generato il", "Valore": datetime.now().strftime("%d/%m/%Y %H:%M")},
            {"Indicatore": "Model version", "Valore": "V2.2"},
        ])
        summary.to_excel(writer, sheet_name="Sintesi", index=False)
    excel_buffer.seek(0)

    st.markdown("### Export")
    e1, e2 = st.columns(2)
    with e1:
        st.download_button("Download CSV", csv_data, "HSE_Risk_Export.csv", "text/csv", use_container_width=True)
    with e2:
        st.download_button("Download Excel", excel_buffer, "HSE_Risk_Export.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)


def render_methodology():
    render_hero(
        "Model governance & methodology",
        "Come leggere il Risk Index, quali dati lo influenzano e quali sono i limiti di utilizzo della piattaforma.",
        eyebrow="TRANSPARENCY · MODEL V2.2"
    )
    render_trust_strip()

    g1, g2 = st.columns([1, 1.2])
    with g1:
        st.markdown("### Risk classification")
        st.dataframe(pd.DataFrame([
            {"Intervallo": "0–20", "Livello": "Basso"},
            {"Intervallo": "21–40", "Livello": "Medio basso"},
            {"Intervallo": "41–60", "Livello": "Medio"},
            {"Intervallo": "61–80", "Livello": "Alto"},
            {"Intervallo": "81–100", "Livello": "Critico"},
        ]), hide_index=True, use_container_width=True)
    with g2:
        st.markdown("### What drives the score")
        st.markdown(
            "Il modello combina **non conformità, criticità, attività, interferenze, complessità organizzativa, "
            "copertura HSE, detection, rapporto ispezioni/NC e correlazioni critiche**. Sensibilizzazioni e Stop Work "
            "agiscono come fattori protettivi entro limiti definiti. La fase del cantiere non genera un malus diretto, "
            "ma viene utilizzata per verificare la coerenza delle sensibilizzazioni."
        )

    st.markdown("### Responsible platform principles")
    p1, p2, p3 = st.columns(3)
    with p1:
        st.markdown('<div class="hse-card"><h3>Explainability</h3><p>Il punteggio viene accompagnato dai driver positivi e negativi, così l’utente può comprenderne la formazione.</p></div>', unsafe_allow_html=True)
    with p2:
        st.markdown('<div class="hse-card"><h3>Traceability</h3><p>Ogni valutazione salvata conserva data, cantiere, compilatore, input principali e output del modello.</p></div>', unsafe_allow_html=True)
    with p3:
        st.markdown('<div class="hse-card"><h3>Human oversight</h3><p>La piattaforma supporta la decisione, ma non automatizza autorizzazioni, Stop Work o giudizi professionali.</p></div>', unsafe_allow_html=True)

    st.markdown("### Model limitations")
    st.info(
        "Il Risk Index è uno strumento di supporto decisionale. Non sostituisce la valutazione dei rischi prevista dalla normativa, "
        "il PSC/POS/DUVRI o altra documentazione applicabile, le procedure aziendali, le verifiche sul campo e il giudizio professionale. "
        "Il risultato dipende dalla correttezza e dalla completezza dei dati inseriti."
    )


def render_admin():
    render_hero(
        "Archive management",
        "Area riservata per la manutenzione dello storico. Le operazioni distruttive sono separate dall'esperienza utente ordinaria.",
        eyebrow="ADMIN · V2.2"
    )
    if not database_ready():
        render_db_notice()
        return
    if not admin_ready():
        st.warning(
            "Area Admin non ancora attiva. Per abilitarla aggiungi nei Secrets di Streamlit "
            "`ADMIN_PASSWORD` e `SUPABASE_SECRET_KEY`."
        )
        st.info("Finché questi Secrets non sono configurati, nessun utente può cancellare valutazioni dall'app.")
        return

    if "admin_authenticated" not in st.session_state:
        st.session_state.admin_authenticated = False

    if not st.session_state.admin_authenticated:
        st.markdown('<div class="responsible-note"><b>Restricted area.</b> L’accesso Admin abilita operazioni irreversibili sullo storico. Utilizzare solo per correzioni o manutenzione autorizzata.</div>', unsafe_allow_html=True)
        password = st.text_input("Password amministratore", type="password", key="admin_password_input")
        if st.button("Accedi all'area Admin", type="primary", use_container_width=True):
            if password and password == get_admin_password():
                st.session_state.admin_authenticated = True
                st.rerun()
            else:
                st.error("Password amministratore non corretta.")
        return

    top1, top2 = st.columns([4, 1])
    with top1:
        st.success("Accesso amministratore attivo.")
    with top2:
        if st.button("Esci", use_container_width=True):
            st.session_state.admin_authenticated = False
            st.rerun()

    try:
        df, rows = _load_archive_df()
    except Exception as exc:
        st.error(str(exc))
        return
    if df.empty:
        st.info("Archivio vuoto.")
        return

    st.markdown("### Archived assessments")
    st.dataframe(df[["Codice", "Cantiere", "Data", "Fase", "Risk Index", "Livello", "Compilato da"]].sort_values("Data", ascending=False), use_container_width=True, hide_index=True)

    labels = {f"{row['Codice']} · {row['Cantiere']} · {row['Data']:%d/%m/%Y %H:%M}": row for _, row in df.sort_values("Data", ascending=False).iterrows()}
    selected = st.selectbox("Seleziona la valutazione da gestire", list(labels.keys()), key="admin_selected")
    row = labels[selected]
    payload = _payload_for_id(rows, row["ID"])
    with st.expander("Visualizza dettaglio prima dell'eliminazione", expanded=False):
        render_assessment_detail(payload, row.to_dict())

    st.markdown("### Delete assessment")
    st.warning("Operazione definitiva: il record verrà rimosso da Supabase e sparirà immediatamente da Dashboard, Storico ed Estrazione.")
    confirm = st.checkbox(f"Confermo di voler eliminare {row['Codice']}", key="admin_confirm_delete")
    typed = st.text_input("Scrivi ELIMINA per confermare", key="admin_delete_text")
    if st.button("Elimina definitivamente", type="primary", use_container_width=True, disabled=not (confirm and typed.strip().upper() == "ELIMINA")):
        try:
            delete_assessment(row["ID"])
            st.success(f"Valutazione {row['Codice']} eliminata correttamente.")
            st.session_state.pop("admin_confirm_delete", None)
            st.session_state.pop("admin_delete_text", None)
            st.rerun()
        except Exception as exc:
            st.error(f"Eliminazione non riuscita: {exc}")


# =========================================================
# PARAMETRI DEL MODELLO
# =========================================================

# Gravità NC ridotte a tre livelli
NC_SEVERITY_WEIGHTS = {
    "bassa": 3,
    "media": 8,
    "alta": 16
}

NC_THEME_WEIGHTS = {
    "elettrico": 16,
    "scavi": 16,
    "quota": 14,
    "sollevamento": 14,
    "spazi confinati": 18,
    "incendio": 14,
    "viabilità": 10,
    "dpi": 7,
    "ordine e pulizia": 6,
    "documentale": 5,
    "ambientale": 8,
    "altro": 6
}

ACTIVITY_WEIGHTS = {
    "elettrici": 18,
    "scavi": 18,
    "lavori in quota": 16,
    "sollevamenti": 16,
    "spazi confinati": 20,
    "hot work": 15,
    "movimentazione mezzi": 12,
    "civili": 8,
    "meccanici": 10,
    "altro": 6
}

AWARENESS_BONUS = {
    "rischio elettrico": 8,
    "scavi e sottoservizi": 8,
    "lavori in quota": 7,
    "sollevamenti": 7,
    "spazi confinati": 9,
    "hot work / incendio": 7,
    "viabilità e mezzi": 6,
    "dpi": 4,
    "ordine e pulizia": 4,
    "toolbox generale": 3,
    "altro": 2
}

LINKS_ACTIVITY_NC = {
    "elettrici": ["elettrico"],
    "scavi": ["scavi"],
    "lavori in quota": ["quota"],
    "sollevamenti": ["sollevamento"],
    "spazi confinati": ["spazi confinati"],
    "hot work": ["incendio"],
    "movimentazione mezzi": ["viabilità"]
}

LINKS_AWARENESS_ACTIVITY = {
    "rischio elettrico": ["elettrici"],
    "scavi e sottoservizi": ["scavi"],
    "lavori in quota": ["lavori in quota"],
    "sollevamenti": ["sollevamenti"],
    "spazi confinati": ["spazi confinati"],
    "hot work / incendio": ["hot work"],
    "viabilità e mezzi": ["movimentazione mezzi"],
    "dpi": ["civili", "meccanici", "elettrici", "scavi"],
    "ordine e pulizia": ["civili", "meccanici"]
}


# La fase non genera più un malus diretto sul Risk Index.
# Viene però utilizzata per verificare la coerenza delle
# sensibilizzazioni con il momento operativo del cantiere.
LINKS_PHASE_AWARENESS = {
    "cantierizzazione": [
        "scavi e sottoservizi",
        "viabilità e mezzi",
        "ordine e pulizia",
        "toolbox generale"
    ],
    "costruzione": [
        "lavori in quota",
        "sollevamenti",
        "dpi",
        "ordine e pulizia",
        "toolbox generale"
    ],
    "commissioning": [
        "rischio elettrico",
        "hot work / incendio",
        "dpi",
        "toolbox generale"
    ],
    "punch list": [
        "rischio elettrico",
        "lavori in quota",
        "ordine e pulizia",
        "dpi",
        "toolbox generale"
    ]
}

STOP_WORK_BONUS_PER_EVENT = 1.5
STOP_WORK_BONUS_CAP = 6.0
PHASE_AWARENESS_BONUS_PER_MATCH = 2
PHASE_AWARENESS_BONUS_CAP = 6


# =========================================================
# FUNZIONI GENERALI
# =========================================================
def normalize_list(values):
    return [
        value.strip().lower()
        for value in values
        if value and value.strip()
    ]


def risk_level(risk):
    if risk <= 20:
        return "Basso"
    if risk <= 40:
        return "Medio basso"
    if risk <= 60:
        return "Medio"
    if risk <= 80:
        return "Alto"
    return "Critico"


def safe_filename(value):
    return (
        value.strip()
        .replace(" ", "_")
        .replace("/", "_")
        .replace("\\", "_")
    )


# =========================================================
# INTERFERENZA ATTIVITÀ
# =========================================================
def calculate_interference_score(
    interference_level,
    activities
):
    """
    Il livello di interferenza è impostato dall'utilizzatore da 1 a 10.

    Il risultato è corretto in funzione di:
    - pericolosità dell'attività più critica;
    - numero delle attività contemporanee.

    Livello 1 = impatto minimo.
    Livello 10 = impatto massimo.
    """

    if not activities:
        return 0, 0, 1

    highest_activity_weight = max(
        ACTIVITY_WEIGHTS.get(activity, 6)
        for activity in activities
    )

    # Fattore compreso indicativamente tra 0,30 e 1
    activity_criticality_factor = (
        highest_activity_weight /
        max(ACTIVITY_WEIGHTS.values())
    )

    # Aumento fino al 25% in presenza di più attività
    activity_count_factor = min(
        1 + max(0, len(activities) - 1) * 0.08,
        1.25
    )

    normalized_level = (interference_level - 1) / 9

    interference_score = (
        normalized_level
        * 100
        * activity_criticality_factor
        * activity_count_factor
    )

    interference_score = min(
        round(interference_score, 1),
        100
    )

    return (
        interference_score,
        round(activity_criticality_factor, 2),
        round(activity_count_factor, 2)
    )


# =========================================================
# COPERTURA HSE IMPRESE
# =========================================================
def calculate_hse_coverage(
    contractors,
    company_hse
):
    """
    Confronta il numero di HSE delle imprese con gli appaltatori.

    - HSE = appaltatori: nessun bonus/malus.
    - HSE < appaltatori: malus fino a +30.
    - HSE > appaltatori: bonus fino a -15.
    """

    if contractors <= 0:
        return 0, "Non valutabile"

    ratio = company_hse / contractors

    if company_hse == contractors:
        return 0, "Copertura equilibrata"

    if company_hse < contractors:
        shortage_ratio = (
            contractors - company_hse
        ) / contractors

        malus = min(
            shortage_ratio * 30,
            30
        )

        return round(malus, 1), "Copertura insufficiente"

    surplus_ratio = (
        company_hse - contractors
    ) / contractors

    bonus = min(
        surplus_ratio * 15,
        15
    )

    return -round(bonus, 1), "Copertura superiore al minimo"


# =========================================================
# CALCOLO RISCHIO
# =========================================================
def calculate_risk(data):
    details = []
    explanations = []

    inspections = data["inspections"]
    num_nc = data["num_nc"]
    nc_items = data["nc_items"]
    activities = data["activities"]
    awareness_types = data["awareness_types"]
    awareness_count = data["awareness_count"]

    # -----------------------------------------------------
    # 1. NON CONFORMITÀ
    # -----------------------------------------------------
    nc_score = 0

    for item in nc_items:
        theme = item["theme"]
        severity = item["severity"]

        theme_score = NC_THEME_WEIGHTS.get(theme, 6)
        severity_score = NC_SEVERITY_WEIGHTS.get(severity, 8)

        item_score = theme_score + severity_score
        nc_score += item_score

        explanations.append(
            f"NC su '{theme}' con gravità '{severity}': "
            f"+{item_score} punti grezzi."
        )

    # Le NC non dettagliate ricevono un peso medio prudenziale
    unclassified_nc = max(
        num_nc - len(nc_items),
        0
    )

    if unclassified_nc > 0:
        unclassified_score = min(
            unclassified_nc * 8,
            40
        )

        nc_score += unclassified_score

        explanations.append(
            f"{unclassified_nc} NC non dettagliate: "
            f"+{unclassified_score} punti prudenziali."
        )

    nc_score = min(nc_score, 100)

    # -----------------------------------------------------
    # 2. INDICE DI DETECTION
    # Rapporto tra NC e ispezioni
    # -----------------------------------------------------
    nc_detection_ratio = (
        num_nc / max(inspections, 1)
    )

    if nc_detection_ratio >= 1:
        detection_penalty = 25
        explanations.append(
            "Indice di detection molto elevato: "
            "almeno una NC per ispezione, +25 punti."
        )

    elif nc_detection_ratio >= 0.5:
        detection_penalty = 15
        explanations.append(
            "Indice di detection significativo: "
            "almeno una NC ogni due ispezioni, +15 punti."
        )

    elif nc_detection_ratio >= 0.25:
        detection_penalty = 8
        explanations.append(
            "Indice di detection moderato: "
            "almeno una NC ogni quattro ispezioni, +8 punti."
        )

    else:
        detection_penalty = 0
        explanations.append(
            "Indice di detection contenuto: "
            "nessun malus aggiuntivo."
        )

    # -----------------------------------------------------
    # 3. BONUS/MALUS RAPPORTO ISPEZIONI
    # -----------------------------------------------------
    if inspections == 0 and num_nc > 0:
        inspection_ratio_adjustment = 30

        explanations.append(
            "NC presenti senza ispezioni registrate: "
            "malus rapporto ispezioni +30."
        )

    elif inspections < 3 and num_nc >= 3:
        inspection_ratio_adjustment = 25

        explanations.append(
            "Numero di ispezioni insufficiente rispetto alle NC: "
            "malus rapporto ispezioni +25."
        )

    elif inspections >= 20 and num_nc <= 2:
        inspection_ratio_adjustment = -10

        explanations.append(
            "Elevato presidio ispettivo con poche NC: "
            "bonus rapporto ispezioni -10."
        )

    elif inspections >= 10 and num_nc <= 2:
        inspection_ratio_adjustment = -5

        explanations.append(
            "Buon presidio ispettivo con poche NC: "
            "bonus rapporto ispezioni -5."
        )

    else:
        inspection_ratio_adjustment = 0

        explanations.append(
            "Rapporto ispezioni/NC nella fascia neutra: "
            "nessun bonus o malus."
        )

    # -----------------------------------------------------
    # 4. STOP WORK - SOLO BONUS, CON IMPATTO LIMITATO
    # -----------------------------------------------------
    stop_work_bonus = min(
        data["stopworks"] * STOP_WORK_BONUS_PER_EVENT,
        STOP_WORK_BONUS_CAP
    )

    if data["stopworks"] > 0:
        explanations.append(
            f"{data['stopworks']} Stop Work registrati: "
            f"bonus cultura della sicurezza -{stop_work_bonus:.1f}. "
            "Il beneficio è limitato per evitare che compensi "
            "eccessivamente altri fattori di rischio."
        )
    else:
        explanations.append(
            "Nessuno Stop Work registrato: nessun bonus applicato."
        )

    # -----------------------------------------------------
    # 5. CRITICITÀ
    # -----------------------------------------------------
    crit_score = min(
        data["crit_open"] * 12
        + data["crit_late"] * 10
        + data["crit_ontime"] * 3,
        100
    )

    if data["crit_open"] > 0:
        explanations.append(
            f"{data['crit_open']} criticità aperte: "
            "incremento del rischio."
        )

    if data["crit_late"] > 0:
        explanations.append(
            f"{data['crit_late']} criticità risolte in ritardo: "
            "peggioramento del follow-up."
        )

    if data["crit_ontime"] > 0:
        explanations.append(
            f"{data['crit_ontime']} criticità risolte nei tempi: "
            "mantengono un peso residuo limitato."
        )

    # -----------------------------------------------------
    # 6. COMPLESSITÀ ORGANIZZATIVA
    # -----------------------------------------------------
    complexity_score = min(
        data["app"] * 5
        + data["sub"] * 8,
        70
    )

    if data["sub"] >= 5:
        complexity_score += 15

        explanations.append(
            "Numero elevato di subappaltatori: "
            "+15 punti grezzi aggiuntivi."
        )

    complexity_score = min(
        complexity_score,
        100
    )

    # -----------------------------------------------------
    # 7. COPERTURA HSE IMPRESE
    # -----------------------------------------------------
    hse_coverage_adjustment, hse_coverage_status = (
        calculate_hse_coverage(
            data["app"],
            data["company_hse"]
        )
    )

    if hse_coverage_adjustment > 0:
        explanations.append(
            f"Copertura HSE insufficiente: "
            f"{data['company_hse']} HSE per "
            f"{data['app']} appaltatori, "
            f"malus +{hse_coverage_adjustment}."
        )

    elif hse_coverage_adjustment < 0:
        explanations.append(
            f"Copertura HSE superiore al rapporto 1:1: "
            f"{data['company_hse']} HSE per "
            f"{data['app']} appaltatori, "
            f"bonus {hse_coverage_adjustment}."
        )

    else:
        explanations.append(
            "Copertura HSE equilibrata rispetto agli appaltatori: "
            "nessun bonus o malus."
        )

    # -----------------------------------------------------
    # 8. TIPOLOGIA E NUMERO ATTIVITÀ
    # -----------------------------------------------------
    activity_score = 0

    for activity in activities:
        score = ACTIVITY_WEIGHTS.get(
            activity,
            6
        )

        activity_score += score

        explanations.append(
            f"Attività '{activity}': "
            f"+{score} punti grezzi."
        )

    # Mantiene la logica già presente relativa al numero
    # di attività contemporanee.
    if len(activities) >= 3:
        activity_score += 20

        explanations.append(
            "Tre o più tipologie di attività contemporanee: "
            "+20 punti grezzi."
        )

    if len(activities) >= 5:
        activity_score += 15

        explanations.append(
            "Cinque o più tipologie di attività contemporanee: "
            "+15 punti grezzi aggiuntivi."
        )

    activity_score = min(
        activity_score,
        100
    )

    # -----------------------------------------------------
    # 9. LIVELLO INTERFERENZA ATTIVITÀ
    # -----------------------------------------------------
    (
        interference_score,
        activity_criticality_factor,
        activity_count_factor
    ) = calculate_interference_score(
        data["interference_level"],
        activities
    )

    if activities:
        explanations.append(
            f"Livello di interferenza "
            f"{data['interference_level']}/10, "
            f"corretto per pericolosità e numero attività: "
            f"{interference_score} punti grezzi."
        )
    else:
        explanations.append(
            "Nessuna attività selezionata: "
            "interferenza non valorizzata."
        )

    # -----------------------------------------------------
    # 10. CORRELAZIONE ATTIVITÀ / NC
    # -----------------------------------------------------
    correlation_penalty = 0

    nc_themes = [
        item["theme"]
        for item in nc_items
    ]

    for activity in activities:
        expected_nc_themes = LINKS_ACTIVITY_NC.get(
            activity,
            []
        )

        for theme in expected_nc_themes:
            if theme in nc_themes:
                correlation_penalty += 15

                explanations.append(
                    f"Correlazione critica tra attività "
                    f"'{activity}' e NC su '{theme}': "
                    "+15 punti."
                )

    correlation_penalty = min(
        correlation_penalty,
        45
    )

    # -----------------------------------------------------
    # 11. FASE DEL CANTIERE
    # -----------------------------------------------------
    # La fase resta un dato operativo e viene utilizzata per
    # valutare la coerenza delle sensibilizzazioni. Non genera
    # più alcun malus diretto sul Risk Index.
    phase_score = 0

    explanations.append(
        f"Fase '{data['fase']}': nessun impatto diretto "
        "sul Risk Index; utilizzata per la verifica di coerenza "
        "delle sensibilizzazioni."
    )

    # -----------------------------------------------------
    # 12. BONUS SENSIBILIZZAZIONI
    # -----------------------------------------------------
    awareness_bonus = 0

    for awareness in awareness_types:
        bonus = AWARENESS_BONUS.get(
            awareness,
            2
        )

        awareness_bonus += bonus

        explanations.append(
            f"Sensibilizzazione '{awareness}': "
            f"bonus -{bonus}."
        )

    if awareness_count >= 5:
        awareness_bonus += 5

        explanations.append(
            "Almeno cinque sensibilizzazioni: "
            "bonus aggiuntivo -5."
        )

    if awareness_count >= 10:
        awareness_bonus += 5

        explanations.append(
            "Almeno dieci sensibilizzazioni: "
            "bonus continuità -5."
        )

    for awareness in awareness_types:
        covered_activities = LINKS_AWARENESS_ACTIVITY.get(
            awareness,
            []
        )

        for activity in activities:
            if activity in covered_activities:
                awareness_bonus += 5

                explanations.append(
                    f"Sensibilizzazione '{awareness}' "
                    f"coerente con attività '{activity}': "
                    "bonus -5."
                )

    phase_awareness_matches = [
        awareness
        for awareness in awareness_types
        if awareness in LINKS_PHASE_AWARENESS.get(
            data["fase"],
            []
        )
    ]

    phase_awareness_bonus = min(
        len(phase_awareness_matches)
        * PHASE_AWARENESS_BONUS_PER_MATCH,
        PHASE_AWARENESS_BONUS_CAP
    )

    if phase_awareness_bonus > 0:
        awareness_bonus += phase_awareness_bonus

        explanations.append(
            f"Sensibilizzazioni coerenti con la fase "
            f"'{data['fase']}': "
            f"{', '.join(phase_awareness_matches)}. "
            f"Bonus aggiuntivo -{phase_awareness_bonus}."
        )
    else:
        explanations.append(
            f"Nessuna sensibilizzazione specificamente coerente "
            f"con la fase '{data['fase']}': "
            "nessun bonus di coerenza fase."
        )

    awareness_bonus = min(
        awareness_bonus,
        35
    )

    # -----------------------------------------------------
    # FORMULA FINALE
    # -----------------------------------------------------
    risk = (
        nc_score * 0.20
        + crit_score * 0.15
        + activity_score * 0.18
        + complexity_score * 0.10
        + interference_score * 0.12
        + detection_penalty
        + inspection_ratio_adjustment
        + correlation_penalty
        + hse_coverage_adjustment
        - awareness_bonus
        - stop_work_bonus
    )

    risk = max(
        0,
        min(100, risk)
    )

    level = risk_level(risk)

    # Nel grafico vengono mostrati gli impatti effettivi
    # già moltiplicati per il peso della formula.
    details.append((
        "NC",
        round(nc_score * 0.20, 1)
    ))

    details.append((
        "Criticità",
        round(crit_score * 0.15, 1)
    ))

    details.append((
        "Attività",
        round(activity_score * 0.18, 1)
    ))

    details.append((
        "Interferenza attività",
        round(interference_score * 0.12, 1)
    ))

    details.append((
        "Complessità organizzativa",
        round(complexity_score * 0.10, 1)
    ))

    details.append((
        "Copertura HSE imprese",
        round(hse_coverage_adjustment, 1)
    ))

    details.append((
        "Bonus Stop Work",
        round(-stop_work_bonus, 1)
    ))

    details.append((
        "Indice di detection",
        round(detection_penalty, 1)
    ))

    details.append((
        "Bonus/Malus rapporto ispezioni",
        round(inspection_ratio_adjustment, 1)
    ))

    details.append((
        "Correlazioni critiche",
        round(correlation_penalty, 1)
    ))

    details.append((
        "Bonus sensibilizzazioni",
        round(-awareness_bonus, 1)
    ))

    technical_data = {
        "nc_detection_ratio": round(nc_detection_ratio, 2),
        "interference_score": interference_score,
        "activity_criticality_factor": activity_criticality_factor,
        "activity_count_factor": activity_count_factor,
        "hse_coverage_adjustment": hse_coverage_adjustment,
        "hse_coverage_status": hse_coverage_status,
        "stop_work_bonus": round(stop_work_bonus, 1),
        "phase_direct_impact": 0,
        "phase_awareness_matches": phase_awareness_matches,
        "phase_awareness_bonus": phase_awareness_bonus
    }

    return (
        risk,
        level,
        details,
        explanations,
        technical_data
    )


# =========================================================
# AZIONI E ALERT
# =========================================================
def generate_base_actions(
    data,
    risk,
    level,
    technical_data
):
    actions = []

    nc_themes = [
        item["theme"]
        for item in data["nc_items"]
    ]

    high_nc = [
        item
        for item in data["nc_items"]
        if item["severity"] == "alta"
    ]

    if risk >= 80:
        actions.append((
            "MUST HAVE",
            "🔴 Rischio critico: convocare immediatamente "
            "il coordinamento operativo e validare le misure "
            "prima della prosecuzione."
        ))

    if high_nc:
        actions.append((
            "MUST HAVE",
            f"🔴 Presenti {len(high_nc)} NC ad alta gravità. "
            "Definire responsabile, scadenza, azione correttiva "
            "e verifica di efficacia."
        ))

    if technical_data["nc_detection_ratio"] >= 1:
        actions.append((
            "MUST HAVE",
            "🔴 Indice di detection molto elevato: viene rilevata "
            "almeno una NC per ogni ispezione. Analizzare le cause "
            "ricorrenti e rafforzare la prevenzione."
        ))

    elif technical_data["nc_detection_ratio"] >= 0.5:
        actions.append((
            "MUST HAVE",
            "🟠 Indice di detection significativo: aumentare "
            "il controllo preventivo e verificare la ripetitività "
            "delle NC."
        ))

    if data["interference_level"] >= 8:
        actions.append((
            "MUST HAVE",
            f"🔴 Interferenza attività molto elevata "
            f"({data['interference_level']}/10): valutare "
            "separazione temporale o spaziale delle lavorazioni."
        ))

    elif data["interference_level"] >= 6:
        actions.append((
            "MUST HAVE",
            f"🟠 Interferenza attività elevata "
            f"({data['interference_level']}/10): predisporre "
            "coordinamento giornaliero e matrice interferenze."
        ))

    if data["company_hse"] < data["app"]:
        actions.append((
            "MUST HAVE",
            f"🟠 Copertura HSE insufficiente: "
            f"{data['company_hse']} HSE rispetto a "
            f"{data['app']} appaltatori. Rafforzare il presidio "
            "o formalizzare una diversa organizzazione di controllo."
        ))

    elif data["company_hse"] > data["app"]:
        actions.append((
            "NICE TO HAVE",
            "🟢 Copertura HSE superiore al rapporto 1:1. "
            "Mantenere il presidio e orientarlo sulle attività "
            "a maggiore interferenza."
        ))

    for activity in data["activities"]:
        if activity == "elettrici":
            actions.append((
                "MUST HAVE",
                "⚡ Verificare sezionamenti, LOTO, autorizzazioni, "
                "DPI e competenze PES/PAV/PEI."
            ))

        elif activity == "scavi":
            actions.append((
                "MUST HAVE",
                "🚧 Verificare sottoservizi, stabilità dei fronti, "
                "accessi, delimitazioni e condizioni dello scavo."
            ))

        elif activity == "lavori in quota":
            actions.append((
                "MUST HAVE",
                "🪜 Verificare parapetti, ancoraggi, PLE, "
                "imbracature e piano di emergenza."
            ))

        elif activity == "sollevamenti":
            actions.append((
                "MUST HAVE",
                "🏗️ Verificare piano di sollevamento, portate, "
                "accessori, interferenze e segregazione dell'area."
            ))

        elif activity == "spazi confinati":
            actions.append((
                "MUST HAVE",
                "☠️ Verificare autorizzazione, atmosfera, "
                "recupero di emergenza e presidio esterno."
            ))

        elif activity == "hot work":
            actions.append((
                "MUST HAVE",
                "🔥 Verificare permesso hot work, estintori, "
                "materiali combustibili e fire watch."
            ))

    if len(data["activities"]) >= 3:
        actions.append((
            "MUST HAVE",
            "🔴 Predisporre matrice interferenziale giornaliera "
            "e briefing pre-job tra le imprese."
        ))

    if "elettrico" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC elettriche: sospendere le lavorazioni "
            "non conformi fino al ripristino delle condizioni sicure."
        ))

    if "scavi" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC sugli scavi: rivalutare stabilità, accessi, "
            "delimitazioni e sottoservizi."
        ))

    if "quota" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC sui lavori in quota: verificare immediatamente "
            "protezioni collettive e sistemi anticaduta."
        ))

    if data["awareness_count"] == 0:
        actions.append((
            "MUST HAVE",
            "🟡 Nessuna sensibilizzazione registrata: "
            "pianificare toolbox mirati prima delle attività critiche."
        ))

    missing_awareness = []

    if (
        "elettrici" in data["activities"]
        and "rischio elettrico" not in data["awareness_types"]
    ):
        missing_awareness.append("rischio elettrico")

    if (
        "scavi" in data["activities"]
        and "scavi e sottoservizi" not in data["awareness_types"]
    ):
        missing_awareness.append("scavi e sottoservizi")

    if (
        "lavori in quota" in data["activities"]
        and "lavori in quota" not in data["awareness_types"]
    ):
        missing_awareness.append("lavori in quota")

    if (
        "sollevamenti" in data["activities"]
        and "sollevamenti" not in data["awareness_types"]
    ):
        missing_awareness.append("sollevamenti")

    if missing_awareness:
        actions.append((
            "MUST HAVE",
            "🟠 Mancano sensibilizzazioni mirate su: "
            + ", ".join(missing_awareness)
            + "."
        ))

    if data["crit_open"] > 0:
        actions.append((
            "MUST HAVE",
            f"🟠 Criticità aperte: {data['crit_open']}. "
            "Definire priorità, owner e data di chiusura."
        ))

    if data["crit_late"] > 0:
        actions.append((
            "MUST HAVE",
            f"🟠 Criticità chiuse in ritardo: "
            f"{data['crit_late']}. Analizzare le cause del ritardo."
        ))

    if level in ["Basso", "Medio basso"]:
        actions.append((
            "NICE TO HAVE",
            "🟢 Mantenere il trend positivo con ispezioni mirate, "
            "tracciamento NC e toolbox sulle attività giornaliere."
        ))

    if len(actions) < 4:
        actions.append((
            "IDEA",
            "Mantenere il presidio operativo e aggiornare "
            "la valutazione al variare delle attività."
        ))

    return actions


# =========================================================
# EXECUTIVE SUMMARY
# =========================================================
def local_executive_summary(
    data,
    risk,
    level,
    driver_df,
    technical_data
):
    malus_drivers = (
        driver_df[driver_df["Valore"] > 0]
        .sort_values("Valore", ascending=False)
        .head(3)
    )

    if malus_drivers.empty:
        top_text = "nessun driver negativo dominante"
    else:
        top_text = ", ".join(
            f"{row['Driver']} ({round(row['Valore'], 1)})"
            for _, row in malus_drivers.iterrows()
        )

    activities = (
        ", ".join(data["activities"])
        if data["activities"]
        else "nessuna attività critica selezionata"
    )

    awareness = (
        ", ".join(data["awareness_types"])
        if data["awareness_types"]
        else "nessuna sensibilizzazione mirata"
    )

    if level in ["Alto", "Critico"]:
        decision = (
            "Si raccomanda un coordinamento operativo immediato, "
            "con verifica delle attività critiche, delle interferenze, "
            "della copertura HSE e delle NC prioritarie."
        )

    elif level == "Medio":
        decision = (
            "Si raccomanda di rafforzare il presidio operativo, "
            "ridurre le interferenze e programmare azioni mirate."
        )

    else:
        decision = (
            "Il livello risulta sotto controllo, ma deve essere "
            "mantenuto il monitoraggio al variare delle attività."
        )

    return f"""
Il cantiere **{data['nome']}** presenta un Risk Index pari a
**{round(risk)} / 100**, classificato come **{level}**.

I principali driver negativi sono: **{top_text}**.

Le attività in corso sono: **{activities}**.

Il livello di interferenza dichiarato è **{data['interference_level']}/10**,
con un impatto calcolato pari a
**{technical_data['interference_score']} punti grezzi**.

La copertura HSE è pari a **{data['company_hse']} HSE**
rispetto a **{data['app']} appaltatori**:
**{technical_data['hse_coverage_status']}**.

L'indice di detection è pari a
**{technical_data['nc_detection_ratio']} NC per ispezione**.

Le sensibilizzazioni registrate sono: **{awareness}**.

{decision}
"""


# =========================================================
# ROOT CAUSE
# =========================================================
def local_root_cause(
    data,
    risk,
    level,
    technical_data,
    driver_df=None
):
    categories = {
        "Non conformità": [],
        "Attività": [],
        "Interferenze": [],
        "Organizzazione": [],
        "Controlli": [],
        "Cultura HSE": []
    }

    high_nc = [item for item in data["nc_items"] if item["severity"] == "alta"]
    nc_themes = [item["theme"] for item in data["nc_items"]]

    if data["num_nc"] > 0:
        categories["Non conformità"].append(
            f"Sono presenti {data['num_nc']} NC, di cui {len(high_nc)} ad alta gravità."
        )
    else:
        categories["Non conformità"].append("Non risultano NC registrate nel periodo analizzato.")

    if technical_data["nc_detection_ratio"] >= 0.5:
        categories["Non conformità"].append(
            f"L'indice di detection è significativo ({technical_data['nc_detection_ratio']} NC/ispezione)."
        )

    if data["activities"]:
        categories["Attività"].append(
            "Attività presenti: " + ", ".join(data["activities"]) + "."
        )
    if len(data["activities"]) >= 3:
        categories["Attività"].append(
            "La contemporaneità di almeno tre tipologie di lavoro aumenta la complessità operativa."
        )

    categories["Interferenze"].append(
        f"Il livello dichiarato è {data['interference_level']}/10, con impatto tecnico pari a {technical_data['interference_score']} punti grezzi."
    )
    if data["interference_level"] >= 6:
        categories["Interferenze"].append(
            "È necessario rafforzare la separazione spaziale o temporale e il coordinamento giornaliero."
        )

    categories["Organizzazione"].append(
        f"Sono presenti {data['app']} appaltatori, {data['sub']} subappaltatori e {data['company_hse']} HSE delle imprese."
    )
    categories["Organizzazione"].append(
        f"Valutazione della copertura HSE: {technical_data['hse_coverage_status']}."
    )
    if data["sub"] >= 5:
        categories["Organizzazione"].append(
            "L'elevato numero di subappaltatori richiede una governance più strutturata."
        )

    if data["crit_open"] > 0:
        categories["Controlli"].append(
            f"Sono ancora aperte {data['crit_open']} criticità."
        )
    if data["crit_late"] > 0:
        categories["Controlli"].append(
            f"Sono state chiuse in ritardo {data['crit_late']} criticità, segnale di follow-up da rafforzare."
        )
    if data["inspections"] == 0:
        categories["Controlli"].append("Non risultano ispezioni registrate.")
    else:
        categories["Controlli"].append(f"Sono state registrate {data['inspections']} ispezioni HSE.")

    if data["awareness_count"] > 0:
        categories["Cultura HSE"].append(
            f"Sono state eseguite {data['awareness_count']} sensibilizzazioni."
        )
    else:
        categories["Cultura HSE"].append(
            "Non risultano sensibilizzazioni registrate: manca un elemento preventivo rilevante."
        )
    if data["stopworks"] > 0:
        categories["Cultura HSE"].append(
            f"Sono stati registrati {data['stopworks']} Stop Work, considerati un segnale positivo di partecipazione."
        )

    contribution_map = {
        "Non conformità": ["NC", "Indice di detection", "Correlazioni critiche"],
        "Attività": ["Attività"],
        "Interferenze": ["Interferenza attività"],
        "Organizzazione": ["Complessità organizzativa", "Copertura HSE imprese"],
        "Controlli": ["Criticità", "Bonus/Malus rapporto ispezioni"],
        "Cultura HSE": ["Bonus Stop Work", "Bonus sensibilizzazioni"],
    }

    contributions = {}
    if driver_df is not None:
        for category, drivers in contribution_map.items():
            value = driver_df[driver_df["Driver"].isin(drivers)]["Valore"].sum()
            contributions[category] = round(float(value), 1)
    else:
        contributions = {category: 0.0 for category in categories}

    positive = {key: max(value, 0) for key, value in contributions.items()}
    total_positive = sum(positive.values())
    percentages = {
        key: round((value / total_positive) * 100) if total_positive > 0 else 0
        for key, value in positive.items()
    }

    blocks = []
    for category, items in categories.items():
        impact = contributions.get(category, 0)
        percentage = percentages.get(category, 0)
        blocks.append(f"### {category}\n\n**Contributo netto:** {impact:+.1f} punti | **Quota dei driver negativi:** {percentage}%")
        blocks.extend(f"- {item}" for item in items)

    dominant = max(positive, key=positive.get) if positive else "nessuna categoria"
    conclusion = (
        f"### Conclusione complessiva\n\nIl Risk Index è **{round(risk)}/100 ({level})**. "
        f"La categoria che contribuisce maggiormente ai driver negativi è **{dominant}**. "
        "Le priorità devono essere definite sulla base dei contributi più elevati, senza trascurare "
        "eventuali NC ad alta gravità o condizioni non controllate anche quando il loro peso numerico è limitato."
    )

    return "\n\n".join(blocks + [conclusion])


# =========================================================
# TOOLBOX TALK
# =========================================================
def local_toolbox_talk(
    data,
    technical_data
):
    activities = data["activities"]

    nc_themes = [
        item["theme"]
        for item in data["nc_items"]
    ]

    title_parts = []

    if activities:
        title_parts.extend(activities[:2])

    if nc_themes:
        title_parts.extend(nc_themes[:2])

    title = (
        " / ".join(title_parts)
        if title_parts
        else "Sicurezza operativa giornaliera"
    )

    rules = [
        "Non iniziare attività senza autorizzazione e briefing pre-job.",
        "Verificare le interferenze con le altre imprese prima di iniziare.",
        "Segnalare immediatamente variazioni del programma di lavoro.",
        "Mantenere l'area ordinata, accessibile e segregata.",
        "Utilizzare i DPI previsti per il rischio specifico.",
        "Applicare lo Stop Work in caso di condizione non controllata."
    ]

    if data["interference_level"] >= 6:
        rules.append(
            "Il livello di interferenza è elevato: rispettare "
            "aree, orari e sequenze operative concordate."
        )

    if data["company_hse"] < data["app"]:
        rules.append(
            "In assenza di presidio HSE dedicato, il preposto deve "
            "segnalare immediatamente ogni variazione o criticità."
        )

    if "elettrici" in activities or "elettrico" in nc_themes:
        rules.append(
            "Per le attività elettriche verificare sezionamento, "
            "assenza tensione e procedura LOTO."
        )

    if "scavi" in activities or "scavi" in nc_themes:
        rules.append(
            "Per gli scavi verificare sottoservizi, stabilità "
            "dei fronti e accessi sicuri."
        )

    if "lavori in quota" in activities or "quota" in nc_themes:
        rules.append(
            "Per i lavori in quota verificare le protezioni "
            "anticaduta prima dell'accesso."
        )

    if "sollevamenti" in activities:
        rules.append(
            "Nessuno deve sostare sotto carichi sospesi "
            "o nelle aree di manovra."
        )

    rules_text = "\n".join(
        f"- {rule}"
        for rule in rules
    )

    return f"""
### Titolo

**Toolbox Talk - {title}**

### Dati operativi

- Fase: **{data['fase']}**
- Interferenza: **{data['interference_level']}/10**
- HSE imprese: **{data['company_hse']}**
- Appaltatori: **{data['app']}**
- Indice di detection:
  **{technical_data['nc_detection_ratio']} NC/ispezione**

### Obiettivo

Allineare la squadra sui rischi della giornata,
sulle interferenze e sulle condizioni minime di sicurezza.

### Regole operative

{rules_text}

### Comportamenti vietati

- Iniziare lavorazioni fuori programma senza coordinamento.
- Rimuovere segregazioni o protezioni.
- Accedere alle aree di un'altra impresa senza autorizzazione.
- Proseguire in presenza di un rischio non controllato.

### Domande finali

1. Quali attività interferiscono con la nostra?
2. Chi coordina le lavorazioni di oggi?
3. Qual è il rischio principale?
4. Chi deve essere avvisato in caso di variazione?
5. Quando deve essere applicato lo Stop Work?
"""


# =========================================================
# ACTION PLAN
# =========================================================
def local_action_plan(
    data,
    risk,
    level,
    technical_data
):
    rows = []

    if risk >= 80:
        rows.append([
            "Coordinamento HSE immediato",
            "Critica",
            "Site Manager / HSE",
            "Immediata",
            "Verbale e piano azioni"
        ])

    if technical_data["nc_detection_ratio"] >= 0.5:
        rows.append([
            "Analizzare indice di detection e ricorrenza NC",
            "Alta",
            "HSE Manager",
            "48 ore",
            "Analisi NC per tema e causa"
        ])

    if data["interference_level"] >= 6:
        rows.append([
            "Formalizzare matrice interferenze giornaliera",
            "Alta",
            "Construction Manager",
            "Prima delle attività",
            "Matrice e briefing firmato"
        ])

    if data["interference_level"] >= 8:
        rows.append([
            "Valutare separazione temporale o spaziale delle attività",
            "Critica",
            "Site Manager",
            "Immediata",
            "Programma lavori aggiornato"
        ])

    if data["company_hse"] < data["app"]:
        rows.append([
            "Rafforzare il presidio HSE delle imprese",
            "Alta",
            "Responsabili appaltatori",
            "48 ore",
            "Organigramma e presenze HSE"
        ])

    if data["crit_open"] > 0:
        rows.append([
            "Chiudere le criticità aperte prioritarie",
            "Alta",
            "HSE / Responsabile impresa",
            "24-48 ore",
            "Registro criticità aggiornato"
        ])

    if data["crit_late"] > 0:
        rows.append([
            "Analizzare i ritardi nella chiusura",
            "Media",
            "HSE Manager",
            "7 giorni",
            "Analisi cause e azioni correttive"
        ])

    if "elettrici" in data["activities"]:
        rows.append([
            "Verificare attività elettriche e LOTO",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist e autorizzazioni"
        ])

    if "scavi" in data["activities"]:
        rows.append([
            "Verificare sicurezza degli scavi",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist scavi"
        ])

    if "lavori in quota" in data["activities"]:
        rows.append([
            "Verificare lavori in quota",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist quota e PLE"
        ])

    if data["awareness_count"] == 0:
        rows.append([
            "Eseguire toolbox mirato",
            "Alta",
            "HSE / Preposto",
            "Prima dell'attività",
            "Registro presenze"
        ])

    if not rows:
        rows.append([
            "Mantenere il monitoraggio HSE",
            "Media",
            "HSE",
            "7 giorni",
            "Report ispezione"
        ])

    return pd.DataFrame(
        rows,
        columns=[
            "Azione",
            "Priorità",
            "Owner suggerito",
            "Scadenza",
            "Evidenza richiesta"
        ]
    )


# =========================================================
# RIDUZIONE RISCHIO
# =========================================================
def local_reduce_risk(
    data,
    risk,
    technical_data
):
    if risk > 80:
        target = 80
    elif risk > 60:
        target = 60
    elif risk > 40:
        target = 40
    else:
        target = 20

    levers = []

    if technical_data["nc_detection_ratio"] >= 0.25:
        levers.append(
            "Ridurre la frequenza delle NC attraverso analisi "
            "delle cause ricorrenti e azioni preventive."
        )

    if data["interference_level"] >= 6:
        levers.append(
            "Ridurre il livello di interferenza separando "
            "le attività per area, orario o sequenza."
        )

    if data["company_hse"] < data["app"]:
        levers.append(
            "Aumentare la copertura HSE delle imprese "
            "fino ad almeno un rapporto 1:1 con gli appaltatori."
        )

    if data["crit_open"] > 0:
        levers.append(
            "Chiudere le criticità aperte più rilevanti."
        )

    if data["awareness_count"] == 0:
        levers.append(
            "Eseguire toolbox mirati sulle attività critiche."
        )

    if data["num_nc"] > 0:
        levers.append(
            "Chiudere le NC ad alta gravità con verifica di efficacia."
        )

    if data["sub"] >= 5:
        levers.append(
            "Rafforzare il coordinamento dei subappaltatori."
        )

    if not levers:
        levers.append(
            "Mantenere controlli periodici e aggiornare "
            "il rischio al cambio delle attività."
        )

    levers_text = "\n".join(
        f"- {lever}"
        for lever in levers
    )

    return f"""
### Obiettivo suggerito

Portare il Risk Index sotto **{target}**.

### Leve principali

{levers_text}

### Strategia consigliata

Agire prioritariamente su:

- indice di detection;
- interferenza tra attività;
- copertura HSE delle imprese;
- NC ad alta gravità;
- criticità aperte;
- sensibilizzazioni mirate.
"""


# =========================================================
# HSE ADVISOR LOCALE
# =========================================================
def local_hse_advisor(
    question,
    data,
    risk,
    level,
    actions_df,
    technical_data
):
    query = question.lower()

    if (
        "interfer" in query
        or "contemporan" in query
    ):
        return f"""
Il livello di interferenza dichiarato è
**{data['interference_level']}/10**.

L'impatto tecnico calcolato è pari a
**{technical_data['interference_score']} punti grezzi**.

Il valore è corretto in funzione della tipologia
e del numero delle attività in corso.

Le azioni prioritarie sono:

1. predisporre una matrice interferenziale;
2. separare le attività per area o orario;
3. effettuare un briefing tra imprese;
4. aggiornare il programma giornaliero;
5. definire chi coordina ogni interferenza.
"""

    if (
        "hse" in query
        or "appaltator" in query
        or "copertura" in query
    ):
        return f"""
Sono presenti **{data['company_hse']} HSE delle imprese**
rispetto a **{data['app']} appaltatori**.

La valutazione della copertura è:
**{technical_data['hse_coverage_status']}**.

Il correttivo applicato all'indice è pari a
**{technical_data['hse_coverage_adjustment']} punti**.

Un valore positivo è un malus.
Un valore negativo è un bonus.
"""

    if (
        "detection" in query
        or "ispezion" in query
        or "rapporto nc" in query
    ):
        return f"""
L'indice di detection è pari a
**{technical_data['nc_detection_ratio']} NC per ispezione**.

Il parametro rappresenta il rapporto tra il numero
totale delle NC e il numero delle ispezioni.

Più il rapporto cresce, maggiore è il malus,
perché aumenta la frequenza con cui vengono
individuate anomalie durante i controlli.
"""

    if (
        "24" in query
        or "domani" in query
        or "subito" in query
    ):
        top_actions = actions_df.head(5)["Azione"].tolist()

        actions_text = "\n".join(
            f"{index + 1}. {action}"
            for index, action in enumerate(top_actions)
        )

        return f"""
Nelle prossime 24 ore suggerisco:

{actions_text}
"""

    if (
        "principale" in query
        or "rischio" in query
    ):
        top_actions = actions_df.head(3)["Azione"].tolist()

        actions_text = "\n".join(
            f"- {action}"
            for action in top_actions
        )

        return f"""
Il rischio principale deriva dalla combinazione di:

- attività critiche;
- livello di interferenza {data['interference_level']}/10;
- indice di detection
  {technical_data['nc_detection_ratio']};
- copertura HSE
  {data['company_hse']} su {data['app']} appaltatori;
- criticità aperte e NC.

Azioni prioritarie:

{actions_text}
"""

    if (
        "ridurre" in query
        or "abbassare" in query
    ):
        return local_reduce_risk(
            data,
            risk,
            technical_data
        )

    if (
        "bonus stop" in query
        or "stop work" in query
        or "stopwork" in query
    ):
        return f"""
Sono stati registrati **{data['stopworks']} Stop Work**.

Nel modello gli Stop Work sono considerati esclusivamente
come un segnale positivo di cultura della sicurezza:
il bonus applicato è **-{technical_data['stop_work_bonus']} punti**.

Il beneficio è volutamente limitato a un massimo di
**-{STOP_WORK_BONUS_CAP} punti**, così non può compensare
NC, criticità o interferenze rilevanti.
"""

    if (
        "fermare" in query
        or "stop" in query
    ):
        if level in ["Alto", "Critico"]:
            return (
                "Con livello Alto o Critico valuterei la sospensione "
                "selettiva delle attività non controllate, in particolare "
                "in presenza di NC ad alta gravità o interferenze elevate."
            )

        return (
            "Non emerge automaticamente la necessità di fermare "
            "le attività, ma deve essere applicato lo Stop Work "
            "in presenza di una condizione non controllata."
        )

    return f"""
Il report presenta un livello **{level}**
con Risk Index **{round(risk)}/100**.

Per interpretarlo considera prioritariamente:

- NC e relativa gravità;
- indice di detection;
- rapporto tra ispezioni e NC;
- attività e livello di interferenza;
- copertura HSE delle imprese;
- criticità aperte;
- sensibilizzazioni mirate.

Prova a chiedere, ad esempio:

- Qual è il rischio principale?
- Come posso ridurre il rischio?
- La copertura HSE è sufficiente?
- Come viene calcolata l'interferenza?
- Cosa devo fare nelle prossime 24 ore?
"""


# =========================================================
# GRAFICO DRIVER BONUS / MALUS
# =========================================================
def create_driver_chart(driver_df):
    chart_df = driver_df.copy()

    chart_df["Tipo"] = chart_df["Valore"].apply(
        lambda value:
        "Malus"
        if value > 0
        else "Bonus"
        if value < 0
        else "Neutro"
    )

    chart_df["Etichetta"] = chart_df["Valore"].apply(
        lambda value: f"{value:+.1f}"
    )

    bars = (
        alt.Chart(chart_df)
        .mark_bar(cornerRadiusEnd=4)
        .encode(
            x=alt.X(
                "Valore:Q",
                title="Impatto sull'indice di rischio"
            ),
            y=alt.Y(
                "Driver:N",
                sort="-x",
                title=None
            ),
            color=alt.Color(
                "Tipo:N",
                scale=alt.Scale(
                    domain=[
                        "Malus",
                        "Bonus",
                        "Neutro"
                    ],
                    range=[
                        "#D62728",
                        "#2CA02C",
                        "#A0A0A0"
                    ]
                ),
                legend=alt.Legend(
                    title="Impatto"
                )
            ),
            tooltip=[
                alt.Tooltip(
                    "Driver:N",
                    title="Driver"
                ),
                alt.Tooltip(
                    "Valore:Q",
                    title="Impatto",
                    format=".1f"
                ),
                alt.Tooltip(
                    "Tipo:N",
                    title="Tipo"
                )
            ]
        )
    )

    zero_line = (
        alt.Chart(
            pd.DataFrame({"zero": [0]})
        )
        .mark_rule(
            color="#444444",
            strokeWidth=1
        )
        .encode(
            x="zero:Q"
        )
    )

    labels = (
        alt.Chart(chart_df)
        .mark_text(
            align="left",
            baseline="middle",
            dx=5
        )
        .encode(
            x="Valore:Q",
            y=alt.Y(
                "Driver:N",
                sort="-x"
            ),
            text="Etichetta:N"
        )
    )

    return (
        bars + zero_line + labels
    ).properties(
        height=420
    )



# =========================================================
# HSE SCORECARD
# =========================================================
def clamp(value, minimum=0.0, maximum=10.0):
    return max(minimum, min(maximum, value))


def generate_hse_scorecard(data, risk, technical_data):
    """Genera una scorecard sintetica su scala 0-10."""
    organization = 10.0
    organization -= min(data["sub"] * 0.35, 2.0)
    organization -= max(technical_data["hse_coverage_adjustment"], 0) / 6
    organization += min(max(-technical_data["hse_coverage_adjustment"], 0) / 8, 1.0)

    planning = 10.0
    planning -= max(data["interference_level"] - 3, 0) * 0.55
    planning -= 1.2 if len(data["activities"]) >= 3 else 0
    planning -= 0.8 if data["crit_late"] > 0 else 0

    operational = 10.0
    operational -= min(technical_data["nc_detection_ratio"] * 3.2, 4.0)
    operational -= min(data["crit_open"] * 0.45, 2.5)
    operational -= min(sum(1 for item in data["nc_items"] if item["severity"] == "alta") * 0.8, 2.0)

    culture = 6.0
    culture += min(data["awareness_count"] * 0.25, 2.0)
    culture += min(data["stopworks"] * 0.30, 1.5)
    culture += min(len(data["awareness_types"]) * 0.18, 0.8)
    culture -= 1.0 if data["awareness_count"] == 0 else 0

    coordination = 10.0
    coordination -= max(data["interference_level"] - 2, 0) * 0.55
    coordination -= min(max(len(data["activities"]) - 2, 0) * 0.45, 1.8)
    coordination -= 0.8 if data["sub"] >= 5 else 0

    rows = [
        ("Organizzazione", organization, "Adeguatezza della struttura, complessità imprese e copertura HSE."),
        ("Pianificazione", planning, "Programmazione delle attività, gestione delle interferenze e puntualità del follow-up."),
        ("Controllo operativo", operational, "Esito dei controlli, frequenza delle NC e criticità ancora aperte."),
        ("Cultura HSE", culture, "Sensibilizzazioni, partecipazione e utilizzo responsabile dello Stop Work."),
        ("Coordinamento", coordination, "Gestione della contemporaneità, delle imprese e delle lavorazioni interferenti."),
    ]

    scorecard_df = pd.DataFrame(rows, columns=["Area", "Punteggio", "Motivazione"])
    scorecard_df["Punteggio"] = scorecard_df["Punteggio"].apply(lambda value: round(clamp(value), 1))
    scorecard_df["Stelle"] = scorecard_df["Punteggio"].apply(
        lambda value: "★" * max(1, min(5, round(value / 2))) + "☆" * (5 - max(1, min(5, round(value / 2))))
    )
    overall = round(scorecard_df["Punteggio"].mean(), 1)
    return scorecard_df, overall


# =========================================================
# HSE TOOLKIT - CHECKLIST WORD
# =========================================================
CHECKLIST_LIBRARY = {
    "Controlli generali pre-attività": [
        "L'attività è inserita nel programma lavori / sinottico aggiornato?",
        "È stato eseguito il coordinamento con le imprese coinvolte e interferenti?",
        "Sono stati individuati il preposto e i referenti operativi dell'attività?",
        "La documentazione operativa è disponibile, aggiornata e coerente con l'attività?",
        "Sono state verificate formazione, informazione, addestramento e idoneità dei lavoratori?",
        "Sono disponibili attrezzature, mezzi e DPI idonei e verificati?",
        "L'area di lavoro è stata delimitata, segnalata e resa accessibile in sicurezza?",
        "Sono state valutate le condizioni ambientali e meteorologiche?",
        "Sono state definite le modalità di comunicazione e gestione delle emergenze?",
        "È stato eseguito il briefing / pre-job check prima dell'avvio?",
    ],
    "Sollevamenti": [
        "È stato predisposto e condiviso il piano di sollevamento?",
        "Sono stati verificati peso, baricentro e punti di presa del carico?",
        "La gru o l'apparecchio di sollevamento è idoneo e correttamente posizionato?",
        "Sono disponibili verifiche, controlli e documentazione dell'attrezzatura?",
        "Gruista, imbracatore e segnalatore possiedono formazione e abilitazioni richieste?",
        "Gli accessori di sollevamento sono identificati, integri e adeguati alla portata?",
        "L'area di manovra e la zona sottostante sono segregate?",
        "Sono state verificate interferenze con linee elettriche, strutture e altre attività?",
        "Sono definite modalità di comunicazione univoche tra gli operatori?",
        "Sono rispettati i limiti meteo e di vento previsti?",
    ],
    "Lavori elettrici": [
        "Sono definiti ruoli, responsabilità e autorizzazioni per il lavoro elettrico?",
        "Sono verificate le competenze PES, PAV, PEI o equivalenti richieste?",
        "È disponibile il permesso / piano di lavoro applicabile?",
        "Sono state eseguite sezionamento, blocco, identificazione e prevenzione delle richiusure?",
        "È stata verificata l'assenza di tensione con strumento idoneo?",
        "Sono state applicate eventuali terre e cortocircuiti previste?",
        "Sono definite e rispettate le distanze di sicurezza dalle parti attive?",
        "DPI, attrezzi isolati e strumenti sono idonei e verificati?",
        "Sono gestite le interferenze con impianti o stalli adiacenti in esercizio?",
        "Sono definite le procedure di emergenza e soccorso elettrico?",
    ],
    "Scavi": [
        "Sono stati individuati e tracciati i sottoservizi presenti?",
        "È disponibile l'autorizzazione allo scavo e la documentazione prevista?",
        "Sono state valutate stabilità dei fronti, profondità e caratteristiche del terreno?",
        "Sono previste armature, pendenze o sistemi di protezione adeguati?",
        "Sono presenti accessi e uscite sicure dallo scavo?",
        "Materiali e mezzi sono mantenuti a distanza sicura dal ciglio?",
        "Lo scavo è delimitato e segnalato anche nelle ore notturne?",
        "Sono gestite acqua, cedimenti, atmosfere pericolose e condizioni meteo?",
        "Sono impediti accessi e soste nel raggio operativo dei mezzi?",
        "È prevista una verifica prima di ogni turno e dopo eventi che possano alterare lo scavo?",
    ],
    "Lavori in quota": [
        "È stata privilegiata una protezione collettiva rispetto ai DPI anticaduta?",
        "Ponteggi, parapetti, trabattelli o opere provvisionali sono conformi e verificati?",
        "Sono individuati e verificati i punti di ancoraggio?",
        "Imbracature, cordini e dispositivi anticaduta sono integri e idonei?",
        "I lavoratori hanno formazione e addestramento specifici?",
        "Sono gestiti caduta materiali e segregazione dell'area sottostante?",
        "Scale e accessi sono utilizzati esclusivamente nelle condizioni consentite?",
        "Sono considerate condizioni meteo e vento?",
        "È disponibile un piano di recupero e soccorso dell'operatore sospeso?",
        "Sono state verificate interferenze con impianti elettrici e altre lavorazioni?",
    ],
    "PLE": [
        "La PLE è idonea alla quota, allo sbraccio e alle condizioni del terreno?",
        "Sono disponibili verifiche periodiche, controlli e manuale d'uso?",
        "L'operatore è abilitato e autorizzato all'utilizzo?",
        "È stata eseguita la verifica pre-uso?",
        "Stabilizzatori e piastre sono correttamente posizionati?",
        "Sono rispettati i limiti di portata e numero di occupanti?",
        "Sono gestite distanze da linee elettriche e ostacoli?",
        "L'area sottostante e il raggio di manovra sono segregati?",
        "È previsto l'utilizzo dei DPI anticaduta indicati dal costruttore?",
        "È noto il sistema di discesa di emergenza e presente personale addestrato a terra?",
    ],
    "Hot Work": [
        "È stato emesso il permesso di lavoro a caldo?",
        "Sono stati rimossi o protetti i materiali combustibili?",
        "Sono disponibili estintori idonei e facilmente raggiungibili?",
        "È stato nominato il fire watch quando necessario?",
        "Bombole, tubazioni e riduttori sono integri e correttamente posizionati?",
        "È garantita adeguata ventilazione?",
        "Sono controllate scintille, scorie e propagazione del calore verso aree adiacenti?",
        "Sono gestite atmosfere potenzialmente esplosive o infiammabili?",
        "I lavoratori utilizzano DPI idonei a calore, radiazioni e fumi?",
        "È previsto il controllo dell'area al termine dell'attività?",
    ],
    "Spazi confinati": [
        "L'ambiente è stato classificato e autorizzato come spazio confinato / sospetto di inquinamento?",
        "È disponibile una procedura specifica e un permesso di ingresso?",
        "Sono stati isolati impianti, energie e possibili ingressi di sostanze?",
        "Sono state eseguite misure atmosferiche prima e durante l'accesso?",
        "Sono garantite ventilazione e condizioni respirabili?",
        "Gli addetti possiedono formazione e addestramento specifici?",
        "È presente un sorvegliante esterno dedicato?",
        "Sono disponibili comunicazioni efficaci tra interno ed esterno?",
        "È predisposto un piano di emergenza con attrezzature di recupero?",
        "È impedito l'accesso a personale non autorizzato?",
    ],
    "Movimentazione mezzi": [
        "È definito e segnalato il piano di viabilità del cantiere?",
        "Sono separati, ove possibile, percorsi pedonali e carrabili?",
        "Mezzi e operatori sono autorizzati e in possesso dei requisiti previsti?",
        "Sono stati eseguiti i controlli pre-uso dei mezzi?",
        "Sono funzionanti segnalatori acustici, luminosi e dispositivi di visibilità?",
        "Sono individuati punti ciechi e necessità di moviere?",
        "Velocità e sensi di marcia sono definiti e rispettati?",
        "Le aree di carico, scarico e manovra sono segregate?",
        "Sono gestite interferenze con scavi, sollevamenti e lavorazioni a terra?",
        "Le condizioni del fondo e delle piste sono adeguate?",
    ],
    "Attività civili": [
        "Le aree di lavoro sono ordinate, pulite e prive di ostacoli?",
        "Sono verificati accessi, percorsi e protezioni dei bordi?",
        "Attrezzature manuali ed elettriche sono idonee e controllate?",
        "Sono gestite polveri, rumore e proiezione di materiali?",
        "Sono previste modalità sicure per taglio, perforazione e demolizioni minori?",
        "Sono stoccati correttamente materiali e prodotti?",
        "Sono controllate le interferenze con impianti esistenti e sottoservizi?",
        "Sono disponibili i DPI specifici per le attività?",
    ],
    "Attività meccaniche": [
        "Macchine e componenti sono stabilizzati prima delle lavorazioni?",
        "Sono isolate e scaricate le energie residue?",
        "Attrezzi e attrezzature sono adatti e in buono stato?",
        "Sono gestiti schiacciamento, cesoiamento e intrappolamento?",
        "Sono previsti sistemi sicuri per movimentazione e posizionamento dei componenti?",
        "Sono controllati fluidi in pressione e temperature residue?",
        "Sono definite le modalità di prova e rimessa in servizio?",
        "L'area è segregata durante prove, avviamenti e movimentazioni?",
    ],
}

ACTIVITY_TO_CHECKLIST = {
    "elettrici": "Lavori elettrici",
    "scavi": "Scavi",
    "lavori in quota": "Lavori in quota",
    "sollevamenti": "Sollevamenti",
    "spazi confinati": "Spazi confinati",
    "hot work": "Hot Work",
    "movimentazione mezzi": "Movimentazione mezzi",
    "civili": "Attività civili",
    "meccanici": "Attività meccaniche",
}


def generate_checklist_word(data, selected_sections):
    document = Document()
    section = document.sections[0]
    section.top_margin = DocxInches(0.65)
    section.bottom_margin = DocxInches(0.65)
    section.left_margin = DocxInches(0.7)
    section.right_margin = DocxInches(0.7)

    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(9)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CHECKLIST HSE PRE-ATTIVITÀ")
    run.bold = True
    run.font.size = Pt(16)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Verifica preventiva da compilare a cura dell'HSE di cantiere").italic = True

    info_table = document.add_table(rows=3, cols=4)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    info_table.style = "Table Grid"
    info_values = [
        ("Cantiere", data.get("nome", ""), "Data", datetime.now().strftime("%d/%m/%Y")),
        ("Fase", data.get("fase", ""), "Impresa", ""),
        ("Compilata da", "", "Attività / area", ""),
    ]
    for row, values in zip(info_table.rows, info_values):
        for index, value in enumerate(values):
            row.cells[index].text = str(value)
            row.cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            if index in (0, 2):
                row.cells[index].paragraphs[0].runs[0].bold = True

    document.add_paragraph()
    legend = document.add_paragraph()
    legend.add_run("Modalità di compilazione: ").bold = True
    legend.add_run("barrare SÌ, NO oppure N.A. e riportare nelle note le azioni necessarie.")

    for section_name in selected_sections:
        questions = CHECKLIST_LIBRARY.get(section_name, [])
        if not questions:
            continue

        heading = document.add_paragraph()
        heading.paragraph_format.space_before = Pt(10)
        heading.paragraph_format.space_after = Pt(4)
        run = heading.add_run(section_name.upper())
        run.bold = True
        run.font.size = Pt(11)

        table = document.add_table(rows=1, cols=6)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = "Table Grid"
        headers = ["N.", "Verifica", "SÌ", "NO", "N.A.", "Note / azioni"]
        for index, header in enumerate(headers):
            cell = table.rows[0].cells[index]
            cell.text = header
            cell.paragraphs[0].runs[0].bold = True
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        for number, question in enumerate(questions, start=1):
            cells = table.add_row().cells
            cells[0].text = str(number)
            cells[1].text = question
            cells[2].text = "☐"
            cells[3].text = "☐"
            cells[4].text = "☐"
            cells[5].text = ""
            for cell in cells:
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        widths = [0.35, 4.55, 0.45, 0.45, 0.5, 1.45]
        for row in table.rows:
            for index, width in enumerate(widths):
                row.cells[index].width = DocxInches(width)

    document.add_paragraph()
    notes = document.add_paragraph()
    notes.add_run("NOTE GENERALI / AZIONI DA INTRAPRENDERE").bold = True
    for _ in range(4):
        document.add_paragraph("________________________________________________________________________________")

    signatures = document.add_table(rows=2, cols=3)
    signatures.alignment = WD_TABLE_ALIGNMENT.CENTER
    signatures.style = "Table Grid"
    for index, label in enumerate(["HSE di cantiere", "Preposto", "Responsabile impresa"]):
        signatures.cell(0, index).text = label
        signatures.cell(0, index).paragraphs[0].runs[0].bold = True
        signatures.cell(1, index).text = "\nFirma: __________________________\n"

    footer = document.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("Checklist HSE generata tramite HSE Risk Platform")

    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer


# =========================================================
# POWERPOINT
# =========================================================
def generate_ppt(
    data,
    risk,
    level,
    driver_df,
    action_plan_df,
    summary,
    root_cause,
    toolbox,
    technical_data,
    scorecard_df,
    scorecard_overall
):
    prs = Presentation()

    slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    slide.shapes.title.text = "HSE Risk Report"

    slide.placeholders[1].text = (
        f"Cantiere: {data['nome']}\n"
        f"Data: {datetime.now().strftime('%d/%m/%Y')}"
    )

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Executive Summary"

    slide.placeholders[1].text = (
        summary
        .replace("**", "")
        .replace("###", "")
    )[:1700]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Sintesi del rischio"

    slide.placeholders[1].text = (
        f"Risk Index: {round(risk)} / 100\n"
        f"Livello: {level}\n\n"
        f"Fase: {data['fase']}\n"
        f"Indice di detection: "
        f"{technical_data['nc_detection_ratio']}\n"
        f"Interferenza: {data['interference_level']}/10\n"
        f"HSE imprese: {data['company_hse']}\n"
        f"Appaltatori: {data['app']}\n"
        f"Copertura HSE: "
        f"{technical_data['hse_coverage_status']}\n"
        f"Attività: "
        f"{', '.join(data['activities']) if data['activities'] else 'Nessuna'}"
    )

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = "Driver del rischio"

    table = slide.shapes.add_table(
        len(driver_df) + 1,
        3,
        Inches(0.6),
        Inches(1.3),
        Inches(8.5),
        Inches(5)
    ).table

    table.cell(0, 0).text = "Driver"
    table.cell(0, 1).text = "Impatto"
    table.cell(0, 2).text = "Tipo"

    for index, row in driver_df.iterrows():
        value = float(row["Valore"])

        impact_type = (
            "Malus"
            if value > 0
            else "Bonus"
            if value < 0
            else "Neutro"
        )

        table.cell(
            index + 1,
            0
        ).text = str(row["Driver"])

        table.cell(
            index + 1,
            1
        ).text = f"{value:+.1f}"

        table.cell(
            index + 1,
            2
        ).text = impact_type

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Action Plan"

    action_text = ""

    for _, row in action_plan_df.head(8).iterrows():
        action_text += (
            f"{row['Priorità']} - "
            f"{row['Azione']} - "
            f"{row['Scadenza']}\n"
        )

    slide.placeholders[1].text = action_text[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Root Cause Analysis"

    slide.placeholders[1].text = (
        root_cause
        .replace("###", "")
        .replace("**", "")
    )[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Toolbox Talk"

    slide.placeholders[1].text = (
        toolbox
        .replace("###", "")
        .replace("**", "")
    )[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = "HSE Scorecard"

    score_table = slide.shapes.add_table(
        len(scorecard_df) + 2,
        4,
        Inches(0.6),
        Inches(1.25),
        Inches(8.6),
        Inches(4.8)
    ).table

    score_table.cell(0, 0).text = "Area"
    score_table.cell(0, 1).text = "Stelle"
    score_table.cell(0, 2).text = "Punteggio"
    score_table.cell(0, 3).text = "Motivazione"

    for index, row in scorecard_df.iterrows():
        score_table.cell(index + 1, 0).text = str(row["Area"])
        score_table.cell(index + 1, 1).text = str(row["Stelle"])
        score_table.cell(index + 1, 2).text = f"{float(row['Punteggio']):.1f} / 10"
        score_table.cell(index + 1, 3).text = str(row["Motivazione"])

    last_row = len(scorecard_df) + 1
    score_table.cell(last_row, 0).text = "VALUTAZIONE COMPLESSIVA"
    score_table.cell(last_row, 1).text = ""
    score_table.cell(last_row, 2).text = f"{scorecard_overall:.1f} / 10"
    score_table.cell(last_row, 3).text = "Media delle cinque aree della scorecard."

    return prs


# =========================================================
# EXCEL
# =========================================================
def generate_excel(
    data,
    risk,
    level,
    driver_df,
    actions_df,
    action_plan_df,
    explanations_df,
    summary,
    root_cause,
    toolbox,
    reduce_risk,
    technical_data,
    scorecard_df,
    scorecard_overall
):
    buffer = BytesIO()

    summary_df = pd.DataFrame([{
        "Data": datetime.now().strftime("%d/%m/%Y %H:%M"),
        "Cantiere": data["nome"],
        "Risk Index": round(risk, 1),
        "Livello": level,
        "Fase": data["fase"],
        "Ispezioni": data["inspections"],
        "NC": data["num_nc"],
        "Indice di detection":
            technical_data["nc_detection_ratio"],
        "Interferenza attività":
            data["interference_level"],
        "Interference score":
            technical_data["interference_score"],
        "Appaltatori": data["app"],
        "Subappaltatori": data["sub"],
        "HSE imprese": data["company_hse"],
        "Copertura HSE":
            technical_data["hse_coverage_status"],
        "Correttivo copertura HSE":
            technical_data["hse_coverage_adjustment"],
        "Stop Work": data["stopworks"],
        "Criticità aperte": data["crit_open"],
        "Criticità in tempo": data["crit_ontime"],
        "Criticità in ritardo": data["crit_late"],
        "Sensibilizzazioni":
            data["awareness_count"],
        "Tipologie attività":
            ", ".join(data["activities"]),
        "Tipologie sensibilizzazioni":
            ", ".join(data["awareness_types"])
    }])

    local_outputs_df = pd.DataFrame([
        {
            "Sezione": "Executive Summary",
            "Testo": summary
        },
        {
            "Sezione": "Root Cause Analysis",
            "Testo": root_cause
        },
        {
            "Sezione": "Toolbox Talk",
            "Testo": toolbox
        },
        {
            "Sezione": "Come ridurre il rischio",
            "Testo": reduce_risk
        }
    ])

    nc_df = pd.DataFrame(
        data["nc_items"]
    )

    with pd.ExcelWriter(
        buffer,
        engine="openpyxl"
    ) as writer:
        summary_df.to_excel(
            writer,
            sheet_name="Sintesi",
            index=False
        )

        driver_df.to_excel(
            writer,
            sheet_name="Driver rischio",
            index=False
        )

        actions_df.to_excel(
            writer,
            sheet_name="Azioni base",
            index=False
        )

        action_plan_df.to_excel(
            writer,
            sheet_name="Action Plan",
            index=False
        )

        explanations_df.to_excel(
            writer,
            sheet_name="Motivazione",
            index=False
        )

        local_outputs_df.to_excel(
            writer,
            sheet_name="Output locali",
            index=False
        )

        nc_df.to_excel(
            writer,
            sheet_name="NC dettaglio",
            index=False
        )

        scorecard_export_df = scorecard_df.copy()
        scorecard_export_df.loc[len(scorecard_export_df)] = {
            "Area": "Valutazione complessiva",
            "Punteggio": scorecard_overall,
            "Motivazione": "Media delle cinque aree della scorecard",
            "Stelle": ""
        }
        scorecard_export_df.to_excel(
            writer,
            sheet_name="HSE Scorecard",
            index=False
        )

    buffer.seek(0)

    return buffer


# =========================================================
# ESECUZIONE CALCOLO E SESSION STATE
# =========================================================
def run_calculation(data):
    (
        risk,
        level,
        details,
        explanations,
        technical_data
    ) = calculate_risk(data)

    base_actions = generate_base_actions(
        data,
        risk,
        level,
        technical_data
    )

    driver_df = pd.DataFrame(
        details,
        columns=[
            "Driver",
            "Valore"
        ]
    )

    actions_df = pd.DataFrame(
        base_actions,
        columns=[
            "Priorità",
            "Azione"
        ]
    )

    explanations_df = pd.DataFrame(
        explanations,
        columns=[
            "Motivazione"
        ]
    )

    summary = local_executive_summary(
        data,
        risk,
        level,
        driver_df,
        technical_data
    )

    root_cause = local_root_cause(
        data,
        risk,
        level,
        technical_data,
        driver_df
    )

    toolbox = local_toolbox_talk(
        data,
        technical_data
    )

    action_plan_df = local_action_plan(
        data,
        risk,
        level,
        technical_data
    )

    reduce_risk = local_reduce_risk(
        data,
        risk,
        technical_data
    )

    scorecard_df, scorecard_overall = generate_hse_scorecard(
        data,
        risk,
        technical_data
    )

    st.session_state.report = {
        "data": data,
        "risk": risk,
        "level": level,
        "driver_df": driver_df,
        "actions_df": actions_df,
        "explanations_df": explanations_df,
        "summary": summary,
        "root_cause": root_cause,
        "toolbox": toolbox,
        "action_plan_df": action_plan_df,
        "reduce_risk": reduce_risk,
        "technical_data": technical_data,
        "scorecard_df": scorecard_df,
        "scorecard_overall": scorecard_overall
    }

    st.session_state.calcolo_effettuato = True
    st.session_state.advisor_answer = ""


if "calcolo_effettuato" not in st.session_state:
    st.session_state.calcolo_effettuato = False

if "advisor_answer" not in st.session_state:
    st.session_state.advisor_answer = ""



# =========================================================
# NAVIGAZIONE PRINCIPALE V2.2
# =========================================================
if st.session_state.get("main_navigation") == "🏠 Home":
    st.session_state["main_navigation"] = "🏠 Overview"
with st.sidebar:
    st.markdown("## HSE RISK")
    st.markdown("**PLATFORM**")
    st.caption("Operational risk intelligence")
    st.write("")
    page = st.radio(
        "Navigazione",
        ["🏠 Overview", "🧮 Risk Assessment", "📊 Dashboard & Storico", "📥 Estrazione dati", "ℹ️ Metodologia", "🔐 Admin"],
        label_visibility="collapsed",
        key="main_navigation"
    )
    st.divider()
    st.caption("Model V2.2 · Deterministic engine")
    st.caption("Supabase archive · Server-side Admin")

if page == "🏠 Overview":
    render_home()
    st.stop()
elif page == "📊 Dashboard & Storico":
    render_dashboard()
    st.stop()
elif page == "📥 Estrazione dati":
    render_extraction()
    st.stop()
elif page == "ℹ️ Metodologia":
    render_methodology()
    st.stop()
elif page == "🔐 Admin":
    render_admin()
    st.stop()


# =========================================================
# RISK ASSESSMENT · GUIDED WORKFLOW
# =========================================================
render_hero(
    "New Risk Assessment",
    "Inserisci i dati operativi, verifica la completezza dell'assessment e calcola un Risk Index spiegabile prima dell'eventuale salvataggio.",
    eyebrow="ASSESSMENT ENGINE · MODEL V2.2"
)
render_trust_strip()

st.markdown(
    """
    <div class="hse-stepper">
        <div class="hse-step"><span>01</span>Project</div>
        <div class="hse-step"><span>02</span>HSE Performance</div>
        <div class="hse-step"><span>03</span>Organisation</div>
        <div class="hse-step"><span>04</span>Activities</div>
        <div class="hse-step"><span>05</span>Review</div>
    </div>
    """,
    unsafe_allow_html=True,
)

step1, step2, step3, step4, step5 = st.tabs([
    "01 · Project",
    "02 · HSE Performance",
    "03 · Organisation",
    "04 · Activities",
    "05 · Review & Calculate",
])

with step1:
    st.markdown("### Project context")
    c1, c2, c3 = st.columns([1.35, 1, 1])
    with c1:
        nome = st.text_input(
            "Nome cantiere *",
            value="",
            placeholder="es. Udine BESS",
            key="assessment_site",
            help="Identifica il sito o il progetto. Il nome sarà utilizzato nello storico e nei report.",
        )
    with c2:
        fase = st.selectbox(
            "Fase *",
            ["cantierizzazione", "costruzione", "commissioning", "punch list"],
            key="assessment_phase",
            help="La fase contestualizza gli output e la coerenza delle sensibilizzazioni; non genera un malus diretto.",
        )
    with c3:
        created_by = st.text_input(
            "Compilato da",
            value="",
            placeholder="Nome e cognome",
            key="assessment_author",
            help="Il nominativo viene salvato nello storico per garantire tracciabilità.",
        )
    st.markdown(
        '<div class="responsible-note"><b>Traceability.</b> Usa un nome cantiere coerente nel tempo e indica il compilatore quando possibile: migliora la leggibilità dello storico senza modificare il Risk Index.</div>',
        unsafe_allow_html=True,
    )

with step2:
    st.markdown("### Controls & non-conformities")
    a1, a2, a3 = st.columns(3)
    with a1:
        inspections = st.number_input(
            "Numero ispezioni",
            min_value=0,
            max_value=1000,
            value=10,
            key="assessment_inspections",
            help="Numero di ispezioni HSE eseguite nel periodo analizzato.",
        )
    with a2:
        num_nc = st.number_input(
            "Numero totale NC",
            min_value=0,
            max_value=300,
            value=0,
            key="assessment_num_nc",
            help="Numero complessivo di non conformità rilevate nel periodo.",
        )
    with a3:
        stopworks = st.number_input(
            "Stop Work",
            min_value=0,
            max_value=100,
            value=0,
            key="assessment_stopworks",
            help="Segnale positivo di cultura della sicurezza con bonus limitato nel modello.",
        )

    nc_items = []
    nc_count_detail = st.number_input(
        "NC da dettagliare",
        min_value=0,
        max_value=20,
        value=0,
        key="assessment_nc_detail_count",
        help="Dettaglia tema e gravità delle NC per migliorare la precisione e la spiegabilità del risultato.",
    )
    effective_nc_detail = min(int(nc_count_detail), int(num_nc), 20)
    if num_nc > 20:
        st.caption("È possibile dettagliare fino a 20 NC; le restanti saranno trattate secondo la logica prudenziale già prevista dal modello.")
    if int(nc_count_detail) > int(num_nc):
        st.warning("Il numero di NC dettagliate non può superare il totale NC: il calcolo userà solo le prime NC coerenti con il totale indicato.")

    if effective_nc_detail:
        st.markdown("#### Non-conformity detail")
        for index in range(effective_nc_detail):
            col1, col2 = st.columns(2)
            with col1:
                theme = st.selectbox(
                    f"Tema NC {index + 1}",
                    list(NC_THEME_WEIGHTS.keys()),
                    key=f"nc_theme_{index}",
                )
            with col2:
                severity = st.selectbox(
                    f"Gravità NC {index + 1}",
                    ["bassa", "media", "alta"],
                    key=f"nc_severity_{index}",
                )
            nc_items.append({"theme": theme, "severity": severity})

    st.markdown("### Follow-up of critical issues")
    q1, q2, q3 = st.columns(3)
    with q1:
        crit_open = st.number_input("Criticità aperte", min_value=0, max_value=200, value=0, key="assessment_crit_open")
    with q2:
        crit_ontime = st.number_input("Risolte nei tempi", min_value=0, max_value=200, value=0, key="assessment_crit_ontime")
    with q3:
        crit_late = st.number_input("Risolte in ritardo", min_value=0, max_value=200, value=0, key="assessment_crit_late")

    st.markdown("### Awareness")
    w1, w2 = st.columns([1, 2])
    with w1:
        awareness_count = st.number_input(
            "Sensibilizzazioni effettuate",
            min_value=0,
            max_value=1000,
            value=0,
            key="assessment_awareness_count",
        )
    with w2:
        awareness_types = st.multiselect(
            "Tipologie sensibilizzazioni",
            list(AWARENESS_BONUS.keys()),
            key="assessment_awareness_types",
            help="Seleziona le tipologie realmente svolte. Il modello verifica anche la coerenza con attività e fase.",
        )

with step3:
    st.markdown("### Organisation & HSE coverage")
    o1, o2, o3 = st.columns(3)
    with o1:
        app = st.number_input(
            "Numero appaltatori",
            min_value=0,
            max_value=100,
            value=1,
            key="assessment_contractors",
        )
    with o2:
        sub = st.number_input(
            "Numero subappaltatori",
            min_value=0,
            max_value=100,
            value=0,
            key="assessment_subcontractors",
        )
    with o3:
        company_hse = st.number_input(
            "HSE delle imprese",
            min_value=0,
            max_value=100,
            value=1,
            key="assessment_company_hse",
            help="Il modello confronta questo valore con il numero di appaltatori.",
        )

    if app > 0:
        ratio = company_hse / app
        if ratio < 1:
            st.warning(f"Copertura HSE inferiore al rapporto 1:1 ({company_hse}/{app}). Il modello applicherà un malus proporzionato.")
        elif ratio == 1:
            st.success("Copertura HSE in rapporto 1:1 con gli appaltatori: fascia neutra del modello.")
        else:
            st.info(f"Copertura HSE superiore al rapporto 1:1 ({company_hse}/{app}). Il bonus resta comunque limitato.")

with step4:
    st.markdown("### Activities & interference")
    activities = st.multiselect(
        "Tipologie attività presenti *",
        list(ACTIVITY_WEIGHTS.keys()),
        default=["civili"],
        key="assessment_activities",
        help="Seleziona tutte le lavorazioni effettivamente in corso nel periodo analizzato.",
    )
    interference_level = st.slider(
        "Livello di interferenza delle attività",
        min_value=1,
        max_value=10,
        value=3,
        step=1,
        key="assessment_interference",
        help="Valutazione dell'interferenza effettiva tra lavorazioni, corretta dal modello per criticità e contemporaneità.",
    )
    i1, i2, i3 = st.columns(3)
    i1.metric("Attività contemporanee", len(activities))
    i2.metric("Interferenza dichiarata", f"{interference_level}/10")
    if activities:
        preview_interference, _, _ = calculate_interference_score(interference_level, normalize_list(activities))
        i3.metric("Interference score", f"{preview_interference:.1f}")
    else:
        i3.metric("Interference score", "—")

    if interference_level <= 3:
        st.success("Interferenza dichiarata bassa.")
    elif interference_level <= 5:
        st.info("Interferenza dichiarata moderata.")
    elif interference_level <= 7:
        st.warning("Interferenza dichiarata elevata: verificare coordinamento, sequenze e segregazioni.")
    else:
        st.error("Interferenza dichiarata molto elevata: verificare separazione temporale/spaziale e condizioni di prosecuzione.")

with step5:
    st.markdown("### Review your assessment")
    draft_data = {
        "nome": nome,
        "created_by": created_by,
        "fase": fase,
        "inspections": inspections,
        "num_nc": num_nc,
        "nc_items": nc_items,
        "stopworks": stopworks,
        "crit_open": crit_open,
        "crit_ontime": crit_ontime,
        "crit_late": crit_late,
        "app": app,
        "sub": sub,
        "company_hse": company_hse,
        "activities": normalize_list(activities),
        "interference_level": interference_level,
        "awareness_count": awareness_count,
        "awareness_types": normalize_list(awareness_types),
    }
    completeness, completeness_label, completeness_notes = calculate_data_completeness(draft_data)

    r1, r2, r3, r4 = st.columns(4)
    r1.metric("Cantiere", nome.strip() or "Da valorizzare")
    r2.metric("NC", int(num_nc))
    r3.metric("Interferenza", f"{interference_level}/10")
    r4.metric("Completezza dati", f"{completeness}%")

    review_rows = pd.DataFrame([
        {"Area": "Project", "Sintesi": f"{nome.strip() or '—'} · {fase} · {created_by.strip() or 'compilatore non indicato'}"},
        {"Area": "HSE performance", "Sintesi": f"{inspections} ispezioni · {num_nc} NC · {stopworks} Stop Work · {crit_open} criticità aperte"},
        {"Area": "Organisation", "Sintesi": f"{app} appaltatori · {sub} subappaltatori · {company_hse} HSE imprese"},
        {"Area": "Activities", "Sintesi": f"{len(activities)} attività · interferenza {interference_level}/10"},
        {"Area": "Awareness", "Sintesi": f"{awareness_count} sensibilizzazioni · {len(awareness_types)} tipologie"},
    ])
    st.dataframe(review_rows, use_container_width=True, hide_index=True)

    if completeness_notes:
        st.warning("Completezza dati: " + completeness_label + ". Verifica: " + "; ".join(completeness_notes) + ".")
    else:
        st.success("Completezza dati buona: non risultano incoerenze di compilazione rilevanti.")

    st.markdown(
        '<div class="responsible-note"><b>Before calculation.</b> Il risultato riflette esclusivamente i dati inseriti. Verifica che il perimetro temporale e operativo sia coerente e che le NC più rilevanti siano dettagliate.</div>',
        unsafe_allow_html=True,
    )

    calcola = st.button(
        "Calculate Risk Index",
        type="primary",
        use_container_width=True,
        key="calculate_risk_v22",
        disabled=not bool(nome.strip()) or not bool(activities),
    )
    if not nome.strip():
        st.caption("Inserisci il nome del cantiere per abilitare il calcolo.")
    elif not activities:
        st.caption("Seleziona almeno un'attività per abilitare il calcolo.")


# =========================================================
# AVVIO DEL CALCOLO
# =========================================================
if calcola:
    run_calculation(draft_data)


if not st.session_state.calcolo_effettuato:
    st.info(
        "Completa i cinque step dell’assessment e usa **Calculate Risk Index** nella sezione Review."
    )
    st.stop()


# =========================================================
# LETTURA REPORT SALVATO
# =========================================================
report = st.session_state.report

data = report["data"]
risk = report["risk"]
level = report["level"]
driver_df = report["driver_df"]
actions_df = report["actions_df"]
explanations_df = report["explanations_df"]
summary = report["summary"]
root_cause = report["root_cause"]
toolbox = report["toolbox"]
action_plan_df = report["action_plan_df"]
reduce_risk = report["reduce_risk"]
technical_data = report["technical_data"]
scorecard_df = report["scorecard_df"]
scorecard_overall = report["scorecard_overall"]



# =========================================================
# DECISION OUTPUT · RESULT + TRACEABILITY
# =========================================================
completeness, completeness_label, completeness_notes = calculate_data_completeness(data)
previous_risk = _previous_site_risk(data.get("nome", ""))
risk_color = risk_level_color(level)
risk_position = max(0, min(100, float(risk)))

if previous_risk is None:
    trend_text = "Prima valutazione disponibile per questo cantiere"
else:
    delta = float(risk) - float(previous_risk)
    if delta > 0:
        trend_text = f"↑ +{delta:.1f} rispetto all'ultima valutazione salvata"
    elif delta < 0:
        trend_text = f"↓ {delta:.1f} rispetto all'ultima valutazione salvata"
    else:
        trend_text = "→ invariato rispetto all'ultima valutazione salvata"

st.markdown("## Decision output")
st.markdown(
    f"""
    <div class="risk-result" style="--risk-color:{risk_color}; --risk-position:{risk_position}%">
        <div class="risk-result-grid">
            <div>
                <div class="hse-kicker">RISK INDEX</div>
                <div><span class="risk-number">{float(risk):.0f}</span> <span class="risk-denom">/ 100</span></div>
                <div class="risk-level">{safe_html(level).upper()}</div>
            </div>
            <div>
                <div class="hse-kicker">RISK POSITION</div>
                <div class="risk-bar"><div class="risk-marker"></div></div>
                <div style="display:flex;justify-content:space-between;color:#708591;font-size:.77rem;"><span>Low</span><span>Medium</span><span>High</span><span>Critical</span></div>
                <p style="margin:.75rem 0 0;color:#435F70;"><b>{safe_html(data['nome'])}</b> · {safe_html(data['fase'])} · {safe_html(trend_text)}</p>
            </div>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

k1, k2, k3, k4 = st.columns(4)
k1.metric("Data completeness", f"{completeness}%", completeness_label)
k2.metric("Detection", technical_data["nc_detection_ratio"])
k3.metric("Interferenza", f"{data['interference_level']} / 10")
k4.metric("Copertura HSE", f"{data['company_hse']} / {data['app']}")

positive_drivers = driver_df[driver_df["Valore"] > 0].sort_values("Valore", ascending=False).head(4)
protective_drivers = driver_df[driver_df["Valore"] < 0].sort_values("Valore").head(4)
left_driver, right_driver = st.columns(2)
with left_driver:
    st.markdown("### Main risk drivers")
    if positive_drivers.empty:
        st.info("Nessun driver negativo dominante.")
    else:
        items = "".join(
            f"<li><span>{safe_html(row['Driver'])}</span><b>+{float(row['Valore']):.1f}</b></li>"
            for _, row in positive_drivers.iterrows()
        )
        st.markdown(f'<div class="hse-panel"><ul class="driver-list">{items}</ul></div>', unsafe_allow_html=True)
with right_driver:
    st.markdown("### Protective factors")
    if protective_drivers.empty:
        st.info("Nessun bonus significativo applicato.")
    else:
        items = "".join(
            f"<li><span>{safe_html(row['Driver'])}</span><b>{float(row['Valore']):.1f}</b></li>"
            for _, row in protective_drivers.iterrows()
        )
        st.markdown(f'<div class="hse-panel"><ul class="driver-list">{items}</ul></div>', unsafe_allow_html=True)

st.markdown("### Executive Summary")
st.markdown(summary)

if completeness_notes:
    st.warning("Data completeness " + completeness_label.lower() + ": " + "; ".join(completeness_notes) + ".")
if data["company_hse"] < data["app"]:
    st.warning("La copertura HSE delle imprese è inferiore al numero degli appaltatori.")
if data["interference_level"] >= 8:
    st.error("Interferenza molto elevata: il punteggio non sostituisce la verifica immediata delle condizioni operative e delle misure di coordinamento.")

st.markdown(
    '<div class="responsible-note"><b>Interpretation.</b> Il Risk Index è un indicatore sintetico e spiegabile. Un valore basso non dimostra automaticamente l’assenza di condizioni pericolose; NC ad alta gravità o rischi non controllati richiedono comunque gestione specifica.</div>',
    unsafe_allow_html=True,
)

# =========================================================
# SALVATAGGIO VALUTAZIONE
# =========================================================
st.markdown("### Save assessment")
col_save_info, col_save_button = st.columns([2.4, 1])
with col_save_info:
    if database_ready():
        st.caption("Il salvataggio archivia input, risultato, driver, action plan, scorecard, data e compilatore. Le simulazioni non vengono archiviate automaticamente.")
    else:
        st.caption("Database non connesso: il report resta disponibile nella sessione corrente ma non sarà persistente.")
with col_save_button:
    if st.button("Save assessment", type="primary", use_container_width=True, key="save_assessment"):
        if not database_ready():
            st.warning("Database non configurato: impossibile archiviare in modo persistente.")
        else:
            try:
                saved = save_assessment(report, data.get("created_by", ""))
                saved_id = saved[0].get("id") if isinstance(saved, list) and saved else None
                st.success(f"Valutazione salvata correttamente{f' · ID {saved_id}' if saved_id else ''}.")
            except Exception as exc:
                st.error(f"Salvataggio non riuscito: {exc}")

st.divider()

# =========================================================
# TAB
# =========================================================
tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10 = st.tabs([
    "🎯 Azioni e alert",
    "📋 Action Plan",
    "🔍 Root Cause",
    "🧰 Toolbox Talk",
    "📉 Riduci rischio",
    "🧮 Dettaglio calcolo",
    "🧭 HSE Guidance",
    "⭐ HSE Scorecard",
    "📑 HSE Toolkit",
    "📤 Export"
])

with tab1:
    st.subheader("Azioni e alert generati dal motore")
    st.dataframe(actions_df, use_container_width=True, hide_index=True)

with tab2:
    st.subheader("Action Plan operativo")
    st.dataframe(action_plan_df, use_container_width=True, hide_index=True)

with tab3:
    st.subheader("Root Cause Analysis per categorie")
    st.caption("La motivazione del punteggio è organizzata per area e collegata ai contributi effettivi dei driver.")
    st.markdown(root_cause)

with tab4:
    st.subheader("Toolbox Talk")
    st.markdown(toolbox)

with tab5:
    st.subheader("Come ridurre il rischio")
    st.markdown(reduce_risk)

with tab6:
    st.subheader("Motivazione tecnica del calcolo")
    st.dataframe(explanations_df, use_container_width=True, hide_index=True)

    st.subheader("Indicatori tecnici")
    technical_df = pd.DataFrame([
        {"Indicatore": "Indice di detection", "Valore": technical_data["nc_detection_ratio"], "Descrizione": "Rapporto tra NC totali e ispezioni"},
        {"Indicatore": "Interference score", "Valore": technical_data["interference_score"], "Descrizione": "Interferenza corretta per tipo e numero attività"},
        {"Indicatore": "Fattore criticità attività", "Valore": technical_data["activity_criticality_factor"], "Descrizione": "Peso dell'attività più critica"},
        {"Indicatore": "Fattore numero attività", "Valore": technical_data["activity_count_factor"], "Descrizione": "Correttivo per contemporaneità"},
        {"Indicatore": "Correttivo copertura HSE", "Valore": technical_data["hse_coverage_adjustment"], "Descrizione": technical_data["hse_coverage_status"]},
        {"Indicatore": "Bonus Stop Work", "Valore": -technical_data["stop_work_bonus"], "Descrizione": "Solo bonus, con impatto massimo limitato a -6"},
        {"Indicatore": "Impatto diretto fase", "Valore": technical_data["phase_direct_impact"], "Descrizione": "La fase non incide direttamente sul Risk Index"},
        {"Indicatore": "Bonus coerenza fase/sensibilizzazioni", "Valore": -technical_data["phase_awareness_bonus"], "Descrizione": ", ".join(technical_data["phase_awareness_matches"]) if technical_data["phase_awareness_matches"] else "Nessuna corrispondenza"}
    ])
    st.dataframe(technical_df, use_container_width=True, hide_index=True)

with tab7:
    st.subheader("HSE Guidance · rules-based")
    st.caption("Questa funzione non utilizza AI generativa o API esterne: interpreta la valutazione tramite regole deterministiche definite nel codice.")
    question = st.text_input(
        "Fai una domanda sul report",
        placeholder="Esempio: la copertura HSE è sufficiente?",
        key="advisor_question"
    )
    if st.button("Generate guidance", key="advisor_button"):
        if question.strip():
            st.session_state.advisor_answer = local_hse_advisor(
                question, data, risk, level, actions_df, technical_data
            )
        else:
            st.session_state.advisor_answer = "Scrivi una domanda."
    if st.session_state.advisor_answer:
        st.markdown(st.session_state.advisor_answer)

with tab8:
    st.subheader("HSE Scorecard")
    st.metric("Valutazione complessiva", f"{scorecard_overall:.1f} / 10")
    for _, row in scorecard_df.iterrows():
        col_score, col_reason = st.columns([1, 2.3])
        with col_score:
            st.markdown(f"### {row['Area']}")
            st.markdown(f"## {row['Stelle']}")
            st.markdown(f"**{float(row['Punteggio']):.1f} / 10**")
        with col_reason:
            st.write(row["Motivazione"])
            st.progress(float(row["Punteggio"]) / 10)
        st.divider()

with tab9:
    st.subheader("Generatore Checklist HSE")
    st.info(
        "La checklist non viene compilata nel tool. Seleziona i capitoli e scarica il format Word da inviare all'HSE di cantiere."
    )

    suggested_sections = ["Controlli generali pre-attività"]
    suggested_sections.extend(
        ACTIVITY_TO_CHECKLIST[activity]
        for activity in data["activities"]
        if activity in ACTIVITY_TO_CHECKLIST
    )
    suggested_sections = list(dict.fromkeys(suggested_sections))

    selected_checklists = st.multiselect(
        "Capitoli da inserire nel documento",
        options=list(CHECKLIST_LIBRARY.keys()),
        default=suggested_sections,
        help="Le attività già selezionate nel calcolo vengono proposte automaticamente, ma puoi aggiungere o rimuovere capitoli."
    )

    if selected_checklists:
        total_questions = sum(len(CHECKLIST_LIBRARY[item]) for item in selected_checklists)
        st.caption(f"Documento composto da {len(selected_checklists)} capitoli e {total_questions} verifiche.")
        checklist_buffer = generate_checklist_word(data, selected_checklists)
        st.download_button(
            label="📥 Scarica Checklist Word",
            data=checklist_buffer,
            file_name=f"Checklist_HSE_{safe_filename(data['nome'])}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )
    else:
        st.warning("Seleziona almeno un capitolo da inserire nella checklist.")

with tab10:
    st.subheader("Download report")
    st.caption("La sezione Export contiene gli output complessivi del Risk Tool. La checklist Word rimane disponibile nell'HSE Toolkit.")

    ppt = generate_ppt(
        data,
        risk,
        level,
        driver_df,
        action_plan_df,
        summary,
        root_cause,
        toolbox,
        technical_data,
        scorecard_df,
        scorecard_overall
    )
    ppt_buffer = BytesIO()
    ppt.save(ppt_buffer)
    ppt_buffer.seek(0)

    excel_buffer = generate_excel(
        data,
        risk,
        level,
        driver_df,
        actions_df,
        action_plan_df,
        explanations_df,
        summary,
        root_cause,
        toolbox,
        reduce_risk,
        technical_data,
        scorecard_df,
        scorecard_overall
    )

    filename = safe_filename(data["nome"])
    col_a, col_b = st.columns(2)

    with col_a:
        st.download_button(
            label="📥 Scarica PowerPoint",
            data=ppt_buffer,
            file_name=f"HSE_Risk_Report_{filename}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

    with col_b:
        st.download_button(
            label="📥 Scarica Excel",
            data=excel_buffer,
            file_name=f"HSE_Risk_Report_{filename}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )
