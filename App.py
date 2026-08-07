from io import BytesIO
from datetime import datetime

import altair as alt
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Pt, Inches as DocxInches


# =========================================================
# CONFIGURAZIONE STREAMLIT
# =========================================================
st.set_page_config(
    page_title="HSE Risk Platform AI Locale",
    page_icon="🦺",
    layout="wide"
)


# =========================================================
# PARAMETRI DEL MODELLO
# =========================================================

# Gravità NC ridotte a tre livelli
NC_SEVERITY_WEIGHTS = {
    "bassa": 3,
    "media": 8,
    "alta": 16
}

NC_THEME_WEIGHTS = {
    "elettrico": 16,
    "scavi": 16,
    "quota": 14,
    "sollevamento": 14,
    "spazi confinati": 18,
    "incendio": 14,
    "viabilità": 10,
    "dpi": 7,
    "ordine e pulizia": 6,
    "documentale": 5,
    "ambientale": 8,
    "altro": 6
}

ACTIVITY_WEIGHTS = {
    "elettrici": 18,
    "scavi": 18,
    "lavori in quota": 16,
    "sollevamenti": 16,
    "spazi confinati": 20,
    "hot work": 15,
    "movimentazione mezzi": 12,
    "civili": 8,
    "meccanici": 10,
    "altro": 6
}

AWARENESS_BONUS = {
    "rischio elettrico": 8,
    "scavi e sottoservizi": 8,
    "lavori in quota": 7,
    "sollevamenti": 7,
    "spazi confinati": 9,
    "hot work / incendio": 7,
    "viabilità e mezzi": 6,
    "dpi": 4,
    "ordine e pulizia": 4,
    "toolbox generale": 3,
    "altro": 2
}

LINKS_ACTIVITY_NC = {
    "elettrici": ["elettrico"],
    "scavi": ["scavi"],
    "lavori in quota": ["quota"],
    "sollevamenti": ["sollevamento"],
    "spazi confinati": ["spazi confinati"],
    "hot work": ["incendio"],
    "movimentazione mezzi": ["viabilità"]
}

LINKS_AWARENESS_ACTIVITY = {
    "rischio elettrico": ["elettrici"],
    "scavi e sottoservizi": ["scavi"],
    "lavori in quota": ["lavori in quota"],
    "sollevamenti": ["sollevamenti"],
    "spazi confinati": ["spazi confinati"],
    "hot work / incendio": ["hot work"],
    "viabilità e mezzi": ["movimentazione mezzi"],
    "dpi": ["civili", "meccanici", "elettrici", "scavi"],
    "ordine e pulizia": ["civili", "meccanici"]
}


# La fase non genera più un malus diretto sul Risk Index.
# Viene però utilizzata per verificare la coerenza delle
# sensibilizzazioni con il momento operativo del cantiere.
LINKS_PHASE_AWARENESS = {
    "cantierizzazione": [
        "scavi e sottoservizi",
        "viabilità e mezzi",
        "ordine e pulizia",
        "toolbox generale"
    ],
    "costruzione": [
        "lavori in quota",
        "sollevamenti",
        "dpi",
        "ordine e pulizia",
        "toolbox generale"
    ],
    "commissioning": [
        "rischio elettrico",
        "hot work / incendio",
        "dpi",
        "toolbox generale"
    ],
    "punch list": [
        "rischio elettrico",
        "lavori in quota",
        "ordine e pulizia",
        "dpi",
        "toolbox generale"
    ]
}

STOP_WORK_BONUS_PER_EVENT = 1.5
STOP_WORK_BONUS_CAP = 6.0
PHASE_AWARENESS_BONUS_PER_MATCH = 2
PHASE_AWARENESS_BONUS_CAP = 6


# =========================================================
# FUNZIONI GENERALI
# =========================================================
def normalize_list(values):
    return [
        value.strip().lower()
        for value in values
        if value and value.strip()
    ]


def risk_level(risk):
    if risk <= 20:
        return "Basso"
    if risk <= 40:
        return "Medio basso"
    if risk <= 60:
        return "Medio"
    if risk <= 80:
        return "Alto"
    return "Critico"


def safe_filename(value):
    return (
        value.strip()
        .replace(" ", "_")
        .replace("/", "_")
        .replace("\\", "_")
    )


# =========================================================
# INTERFERENZA ATTIVITÀ
# =========================================================
def calculate_interference_score(
    interference_level,
    activities
):
    """
    Il livello di interferenza è impostato dall'utilizzatore da 1 a 10.

    Il risultato è corretto in funzione di:
    - pericolosità dell'attività più critica;
    - numero delle attività contemporanee.

    Livello 1 = impatto minimo.
    Livello 10 = impatto massimo.
    """

    if not activities:
        return 0, 0, 1

    highest_activity_weight = max(
        ACTIVITY_WEIGHTS.get(activity, 6)
        for activity in activities
    )

    # Fattore compreso indicativamente tra 0,30 e 1
    activity_criticality_factor = (
        highest_activity_weight /
        max(ACTIVITY_WEIGHTS.values())
    )

    # Aumento fino al 25% in presenza di più attività
    activity_count_factor = min(
        1 + max(0, len(activities) - 1) * 0.08,
        1.25
    )

    normalized_level = (interference_level - 1) / 9

    interference_score = (
        normalized_level
        * 100
        * activity_criticality_factor
        * activity_count_factor
    )

    interference_score = min(
        round(interference_score, 1),
        100
    )

    return (
        interference_score,
        round(activity_criticality_factor, 2),
        round(activity_count_factor, 2)
    )


# =========================================================
# COPERTURA HSE IMPRESE
# =========================================================
def calculate_hse_coverage(
    contractors,
    company_hse
):
    """
    Confronta il numero di HSE delle imprese con gli appaltatori.

    - HSE = appaltatori: nessun bonus/malus.
    - HSE < appaltatori: malus fino a +30.
    - HSE > appaltatori: bonus fino a -15.
    """

    if contractors <= 0:
        return 0, "Non valutabile"

    ratio = company_hse / contractors

    if company_hse == contractors:
        return 0, "Copertura equilibrata"

    if company_hse < contractors:
        shortage_ratio = (
            contractors - company_hse
        ) / contractors

        malus = min(
            shortage_ratio * 30,
            30
        )

        return round(malus, 1), "Copertura insufficiente"

    surplus_ratio = (
        company_hse - contractors
    ) / contractors

    bonus = min(
        surplus_ratio * 15,
        15
    )

    return -round(bonus, 1), "Copertura superiore al minimo"


# =========================================================
# CALCOLO RISCHIO
# =========================================================
def calculate_risk(data):
    details = []
    explanations = []

    inspections = data["inspections"]
    num_nc = data["num_nc"]
    nc_items = data["nc_items"]
    activities = data["activities"]
    awareness_types = data["awareness_types"]
    awareness_count = data["awareness_count"]

    # -----------------------------------------------------
    # 1. NON CONFORMITÀ
    # -----------------------------------------------------
    nc_score = 0

    for item in nc_items:
        theme = item["theme"]
        severity = item["severity"]

        theme_score = NC_THEME_WEIGHTS.get(theme, 6)
        severity_score = NC_SEVERITY_WEIGHTS.get(severity, 8)

        item_score = theme_score + severity_score
        nc_score += item_score

        explanations.append(
            f"NC su '{theme}' con gravità '{severity}': "
            f"+{item_score} punti grezzi."
        )

    # Le NC non dettagliate ricevono un peso medio prudenziale
    unclassified_nc = max(
        num_nc - len(nc_items),
        0
    )

    if unclassified_nc > 0:
        unclassified_score = min(
            unclassified_nc * 8,
            40
        )

        nc_score += unclassified_score

        explanations.append(
            f"{unclassified_nc} NC non dettagliate: "
            f"+{unclassified_score} punti prudenziali."
        )

    nc_score = min(nc_score, 100)

    # -----------------------------------------------------
    # 2. INDICE DI DETECTION
    # Rapporto tra NC e ispezioni
    # -----------------------------------------------------
    nc_detection_ratio = (
        num_nc / max(inspections, 1)
    )

    if nc_detection_ratio >= 1:
        detection_penalty = 25
        explanations.append(
            "Indice di detection molto elevato: "
            "almeno una NC per ispezione, +25 punti."
        )

    elif nc_detection_ratio >= 0.5:
        detection_penalty = 15
        explanations.append(
            "Indice di detection significativo: "
            "almeno una NC ogni due ispezioni, +15 punti."
        )

    elif nc_detection_ratio >= 0.25:
        detection_penalty = 8
        explanations.append(
            "Indice di detection moderato: "
            "almeno una NC ogni quattro ispezioni, +8 punti."
        )

    else:
        detection_penalty = 0
        explanations.append(
            "Indice di detection contenuto: "
            "nessun malus aggiuntivo."
        )

    # -----------------------------------------------------
    # 3. BONUS/MALUS RAPPORTO ISPEZIONI
    # -----------------------------------------------------
    if inspections == 0 and num_nc > 0:
        inspection_ratio_adjustment = 30

        explanations.append(
            "NC presenti senza ispezioni registrate: "
            "malus rapporto ispezioni +30."
        )

    elif inspections < 3 and num_nc >= 3:
        inspection_ratio_adjustment = 25

        explanations.append(
            "Numero di ispezioni insufficiente rispetto alle NC: "
            "malus rapporto ispezioni +25."
        )

    elif inspections >= 20 and num_nc <= 2:
        inspection_ratio_adjustment = -10

        explanations.append(
            "Elevato presidio ispettivo con poche NC: "
            "bonus rapporto ispezioni -10."
        )

    elif inspections >= 10 and num_nc <= 2:
        inspection_ratio_adjustment = -5

        explanations.append(
            "Buon presidio ispettivo con poche NC: "
            "bonus rapporto ispezioni -5."
        )

    else:
        inspection_ratio_adjustment = 0

        explanations.append(
            "Rapporto ispezioni/NC nella fascia neutra: "
            "nessun bonus o malus."
        )

    # -----------------------------------------------------
    # 4. STOP WORK - SOLO BONUS, CON IMPATTO LIMITATO
    # -----------------------------------------------------
    stop_work_bonus = min(
        data["stopworks"] * STOP_WORK_BONUS_PER_EVENT,
        STOP_WORK_BONUS_CAP
    )

    if data["stopworks"] > 0:
        explanations.append(
            f"{data['stopworks']} Stop Work registrati: "
            f"bonus cultura della sicurezza -{stop_work_bonus:.1f}. "
            "Il beneficio è limitato per evitare che compensi "
            "eccessivamente altri fattori di rischio."
        )
    else:
        explanations.append(
            "Nessuno Stop Work registrato: nessun bonus applicato."
        )

    # -----------------------------------------------------
    # 5. CRITICITÀ
    # -----------------------------------------------------
    crit_score = min(
        data["crit_open"] * 12
        + data["crit_late"] * 10
        + data["crit_ontime"] * 3,
        100
    )

    if data["crit_open"] > 0:
        explanations.append(
            f"{data['crit_open']} criticità aperte: "
            "incremento del rischio."
        )

    if data["crit_late"] > 0:
        explanations.append(
            f"{data['crit_late']} criticità risolte in ritardo: "
            "peggioramento del follow-up."
        )

    if data["crit_ontime"] > 0:
        explanations.append(
            f"{data['crit_ontime']} criticità risolte nei tempi: "
            "mantengono un peso residuo limitato."
        )

    # -----------------------------------------------------
    # 6. COMPLESSITÀ ORGANIZZATIVA
    # -----------------------------------------------------
    complexity_score = min(
        data["app"] * 5
        + data["sub"] * 8,
        70
    )

    if data["sub"] >= 5:
        complexity_score += 15

        explanations.append(
            "Numero elevato di subappaltatori: "
            "+15 punti grezzi aggiuntivi."
        )

    complexity_score = min(
        complexity_score,
        100
    )

    # -----------------------------------------------------
    # 7. COPERTURA HSE IMPRESE
    # -----------------------------------------------------
    hse_coverage_adjustment, hse_coverage_status = (
        calculate_hse_coverage(
            data["app"],
            data["company_hse"]
        )
    )

    if hse_coverage_adjustment > 0:
        explanations.append(
            f"Copertura HSE insufficiente: "
            f"{data['company_hse']} HSE per "
            f"{data['app']} appaltatori, "
            f"malus +{hse_coverage_adjustment}."
        )

    elif hse_coverage_adjustment < 0:
        explanations.append(
            f"Copertura HSE superiore al rapporto 1:1: "
            f"{data['company_hse']} HSE per "
            f"{data['app']} appaltatori, "
            f"bonus {hse_coverage_adjustment}."
        )

    else:
        explanations.append(
            "Copertura HSE equilibrata rispetto agli appaltatori: "
            "nessun bonus o malus."
        )

    # -----------------------------------------------------
    # 8. TIPOLOGIA E NUMERO ATTIVITÀ
    # -----------------------------------------------------
    activity_score = 0

    for activity in activities:
        score = ACTIVITY_WEIGHTS.get(
            activity,
            6
        )

        activity_score += score

        explanations.append(
            f"Attività '{activity}': "
            f"+{score} punti grezzi."
        )

    # Mantiene la logica già presente relativa al numero
    # di attività contemporanee.
    if len(activities) >= 3:
        activity_score += 20

        explanations.append(
            "Tre o più tipologie di attività contemporanee: "
            "+20 punti grezzi."
        )

    if len(activities) >= 5:
        activity_score += 15

        explanations.append(
            "Cinque o più tipologie di attività contemporanee: "
            "+15 punti grezzi aggiuntivi."
        )

    activity_score = min(
        activity_score,
        100
    )

    # -----------------------------------------------------
    # 9. LIVELLO INTERFERENZA ATTIVITÀ
    # -----------------------------------------------------
    (
        interference_score,
        activity_criticality_factor,
        activity_count_factor
    ) = calculate_interference_score(
        data["interference_level"],
        activities
    )

    if activities:
        explanations.append(
            f"Livello di interferenza "
            f"{data['interference_level']}/10, "
            f"corretto per pericolosità e numero attività: "
            f"{interference_score} punti grezzi."
        )
    else:
        explanations.append(
            "Nessuna attività selezionata: "
            "interferenza non valorizzata."
        )

    # -----------------------------------------------------
    # 10. CORRELAZIONE ATTIVITÀ / NC
    # -----------------------------------------------------
    correlation_penalty = 0

    nc_themes = [
        item["theme"]
        for item in nc_items
    ]

    for activity in activities:
        expected_nc_themes = LINKS_ACTIVITY_NC.get(
            activity,
            []
        )

        for theme in expected_nc_themes:
            if theme in nc_themes:
                correlation_penalty += 15

                explanations.append(
                    f"Correlazione critica tra attività "
                    f"'{activity}' e NC su '{theme}': "
                    "+15 punti."
                )

    correlation_penalty = min(
        correlation_penalty,
        45
    )

    # -----------------------------------------------------
    # 11. FASE DEL CANTIERE
    # -----------------------------------------------------
    # La fase resta un dato operativo e viene utilizzata per
    # valutare la coerenza delle sensibilizzazioni. Non genera
    # più alcun malus diretto sul Risk Index.
    phase_score = 0

    explanations.append(
        f"Fase '{data['fase']}': nessun impatto diretto "
        "sul Risk Index; utilizzata per la verifica di coerenza "
        "delle sensibilizzazioni."
    )

    # -----------------------------------------------------
    # 12. BONUS SENSIBILIZZAZIONI
    # -----------------------------------------------------
    awareness_bonus = 0

    for awareness in awareness_types:
        bonus = AWARENESS_BONUS.get(
            awareness,
            2
        )

        awareness_bonus += bonus

        explanations.append(
            f"Sensibilizzazione '{awareness}': "
            f"bonus -{bonus}."
        )

    if awareness_count >= 5:
        awareness_bonus += 5

        explanations.append(
            "Almeno cinque sensibilizzazioni: "
            "bonus aggiuntivo -5."
        )

    if awareness_count >= 10:
        awareness_bonus += 5

        explanations.append(
            "Almeno dieci sensibilizzazioni: "
            "bonus continuità -5."
        )

    for awareness in awareness_types:
        covered_activities = LINKS_AWARENESS_ACTIVITY.get(
            awareness,
            []
        )

        for activity in activities:
            if activity in covered_activities:
                awareness_bonus += 5

                explanations.append(
                    f"Sensibilizzazione '{awareness}' "
                    f"coerente con attività '{activity}': "
                    "bonus -5."
                )

    phase_awareness_matches = [
        awareness
        for awareness in awareness_types
        if awareness in LINKS_PHASE_AWARENESS.get(
            data["fase"],
            []
        )
    ]

    phase_awareness_bonus = min(
        len(phase_awareness_matches)
        * PHASE_AWARENESS_BONUS_PER_MATCH,
        PHASE_AWARENESS_BONUS_CAP
    )

    if phase_awareness_bonus > 0:
        awareness_bonus += phase_awareness_bonus

        explanations.append(
            f"Sensibilizzazioni coerenti con la fase "
            f"'{data['fase']}': "
            f"{', '.join(phase_awareness_matches)}. "
            f"Bonus aggiuntivo -{phase_awareness_bonus}."
        )
    else:
        explanations.append(
            f"Nessuna sensibilizzazione specificamente coerente "
            f"con la fase '{data['fase']}': "
            "nessun bonus di coerenza fase."
        )

    awareness_bonus = min(
        awareness_bonus,
        35
    )

    # -----------------------------------------------------
    # FORMULA FINALE
    # -----------------------------------------------------
    risk = (
        nc_score * 0.20
        + crit_score * 0.15
        + activity_score * 0.18
        + complexity_score * 0.10
        + interference_score * 0.12
        + detection_penalty
        + inspection_ratio_adjustment
        + correlation_penalty
        + hse_coverage_adjustment
        - awareness_bonus
        - stop_work_bonus
    )

    risk = max(
        0,
        min(100, risk)
    )

    level = risk_level(risk)

    # Nel grafico vengono mostrati gli impatti effettivi
    # già moltiplicati per il peso della formula.
    details.append((
        "NC",
        round(nc_score * 0.20, 1)
    ))

    details.append((
        "Criticità",
        round(crit_score * 0.15, 1)
    ))

    details.append((
        "Attività",
        round(activity_score * 0.18, 1)
    ))

    details.append((
        "Interferenza attività",
        round(interference_score * 0.12, 1)
    ))

    details.append((
        "Complessità organizzativa",
        round(complexity_score * 0.10, 1)
    ))

    details.append((
        "Copertura HSE imprese",
        round(hse_coverage_adjustment, 1)
    ))

    details.append((
        "Bonus Stop Work",
        round(-stop_work_bonus, 1)
    ))

    details.append((
        "Indice di detection",
        round(detection_penalty, 1)
    ))

    details.append((
        "Bonus/Malus rapporto ispezioni",
        round(inspection_ratio_adjustment, 1)
    ))

    details.append((
        "Correlazioni critiche",
        round(correlation_penalty, 1)
    ))

    details.append((
        "Bonus sensibilizzazioni",
        round(-awareness_bonus, 1)
    ))

    technical_data = {
        "nc_detection_ratio": round(nc_detection_ratio, 2),
        "interference_score": interference_score,
        "activity_criticality_factor": activity_criticality_factor,
        "activity_count_factor": activity_count_factor,
        "hse_coverage_adjustment": hse_coverage_adjustment,
        "hse_coverage_status": hse_coverage_status,
        "stop_work_bonus": round(stop_work_bonus, 1),
        "phase_direct_impact": 0,
        "phase_awareness_matches": phase_awareness_matches,
        "phase_awareness_bonus": phase_awareness_bonus
    }

    return (
        risk,
        level,
        details,
        explanations,
        technical_data
    )


# =========================================================
# AZIONI E ALERT
# =========================================================
def generate_base_actions(
    data,
    risk,
    level,
    technical_data
):
    actions = []

    nc_themes = [
        item["theme"]
        for item in data["nc_items"]
    ]

    high_nc = [
        item
        for item in data["nc_items"]
        if item["severity"] == "alta"
    ]

    if risk >= 80:
        actions.append((
            "MUST HAVE",
            "🔴 Rischio critico: convocare immediatamente "
            "il coordinamento operativo e validare le misure "
            "prima della prosecuzione."
        ))

    if high_nc:
        actions.append((
            "MUST HAVE",
            f"🔴 Presenti {len(high_nc)} NC ad alta gravità. "
            "Definire responsabile, scadenza, azione correttiva "
            "e verifica di efficacia."
        ))

    if technical_data["nc_detection_ratio"] >= 1:
        actions.append((
            "MUST HAVE",
            "🔴 Indice di detection molto elevato: viene rilevata "
            "almeno una NC per ogni ispezione. Analizzare le cause "
            "ricorrenti e rafforzare la prevenzione."
        ))

    elif technical_data["nc_detection_ratio"] >= 0.5:
        actions.append((
            "MUST HAVE",
            "🟠 Indice di detection significativo: aumentare "
            "il controllo preventivo e verificare la ripetitività "
            "delle NC."
        ))

    if data["interference_level"] >= 8:
        actions.append((
            "MUST HAVE",
            f"🔴 Interferenza attività molto elevata "
            f"({data['interference_level']}/10): valutare "
            "separazione temporale o spaziale delle lavorazioni."
        ))

    elif data["interference_level"] >= 6:
        actions.append((
            "MUST HAVE",
            f"🟠 Interferenza attività elevata "
            f"({data['interference_level']}/10): predisporre "
            "coordinamento giornaliero e matrice interferenze."
        ))

    if data["company_hse"] < data["app"]:
        actions.append((
            "MUST HAVE",
            f"🟠 Copertura HSE insufficiente: "
            f"{data['company_hse']} HSE rispetto a "
            f"{data['app']} appaltatori. Rafforzare il presidio "
            "o formalizzare una diversa organizzazione di controllo."
        ))

    elif data["company_hse"] > data["app"]:
        actions.append((
            "NICE TO HAVE",
            "🟢 Copertura HSE superiore al rapporto 1:1. "
            "Mantenere il presidio e orientarlo sulle attività "
            "a maggiore interferenza."
        ))

    for activity in data["activities"]:
        if activity == "elettrici":
            actions.append((
                "MUST HAVE",
                "⚡ Verificare sezionamenti, LOTO, autorizzazioni, "
                "DPI e competenze PES/PAV/PEI."
            ))

        elif activity == "scavi":
            actions.append((
                "MUST HAVE",
                "🚧 Verificare sottoservizi, stabilità dei fronti, "
                "accessi, delimitazioni e condizioni dello scavo."
            ))

        elif activity == "lavori in quota":
            actions.append((
                "MUST HAVE",
                "🪜 Verificare parapetti, ancoraggi, PLE, "
                "imbracature e piano di emergenza."
            ))

        elif activity == "sollevamenti":
            actions.append((
                "MUST HAVE",
                "🏗️ Verificare piano di sollevamento, portate, "
                "accessori, interferenze e segregazione dell'area."
            ))

        elif activity == "spazi confinati":
            actions.append((
                "MUST HAVE",
                "☠️ Verificare autorizzazione, atmosfera, "
                "recupero di emergenza e presidio esterno."
            ))

        elif activity == "hot work":
            actions.append((
                "MUST HAVE",
                "🔥 Verificare permesso hot work, estintori, "
                "materiali combustibili e fire watch."
            ))

    if len(data["activities"]) >= 3:
        actions.append((
            "MUST HAVE",
            "🔴 Predisporre matrice interferenziale giornaliera "
            "e briefing pre-job tra le imprese."
        ))

    if "elettrico" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC elettriche: sospendere le lavorazioni "
            "non conformi fino al ripristino delle condizioni sicure."
        ))

    if "scavi" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC sugli scavi: rivalutare stabilità, accessi, "
            "delimitazioni e sottoservizi."
        ))

    if "quota" in nc_themes:
        actions.append((
            "MUST HAVE",
            "🔴 NC sui lavori in quota: verificare immediatamente "
            "protezioni collettive e sistemi anticaduta."
        ))

    if data["awareness_count"] == 0:
        actions.append((
            "MUST HAVE",
            "🟡 Nessuna sensibilizzazione registrata: "
            "pianificare toolbox mirati prima delle attività critiche."
        ))

    missing_awareness = []

    if (
        "elettrici" in data["activities"]
        and "rischio elettrico" not in data["awareness_types"]
    ):
        missing_awareness.append("rischio elettrico")

    if (
        "scavi" in data["activities"]
        and "scavi e sottoservizi" not in data["awareness_types"]
    ):
        missing_awareness.append("scavi e sottoservizi")

    if (
        "lavori in quota" in data["activities"]
        and "lavori in quota" not in data["awareness_types"]
    ):
        missing_awareness.append("lavori in quota")

    if (
        "sollevamenti" in data["activities"]
        and "sollevamenti" not in data["awareness_types"]
    ):
        missing_awareness.append("sollevamenti")

    if missing_awareness:
        actions.append((
            "MUST HAVE",
            "🟠 Mancano sensibilizzazioni mirate su: "
            + ", ".join(missing_awareness)
            + "."
        ))

    if data["crit_open"] > 0:
        actions.append((
            "MUST HAVE",
            f"🟠 Criticità aperte: {data['crit_open']}. "
            "Definire priorità, owner e data di chiusura."
        ))

    if data["crit_late"] > 0:
        actions.append((
            "MUST HAVE",
            f"🟠 Criticità chiuse in ritardo: "
            f"{data['crit_late']}. Analizzare le cause del ritardo."
        ))

    if level in ["Basso", "Medio basso"]:
        actions.append((
            "NICE TO HAVE",
            "🟢 Mantenere il trend positivo con ispezioni mirate, "
            "tracciamento NC e toolbox sulle attività giornaliere."
        ))

    if len(actions) < 4:
        actions.append((
            "IDEA",
            "Mantenere il presidio operativo e aggiornare "
            "la valutazione al variare delle attività."
        ))

    return actions


# =========================================================
# EXECUTIVE SUMMARY
# =========================================================
def local_executive_summary(
    data,
    risk,
    level,
    driver_df,
    technical_data
):
    malus_drivers = (
        driver_df[driver_df["Valore"] > 0]
        .sort_values("Valore", ascending=False)
        .head(3)
    )

    if malus_drivers.empty:
        top_text = "nessun driver negativo dominante"
    else:
        top_text = ", ".join(
            f"{row['Driver']} ({round(row['Valore'], 1)})"
            for _, row in malus_drivers.iterrows()
        )

    activities = (
        ", ".join(data["activities"])
        if data["activities"]
        else "nessuna attività critica selezionata"
    )

    awareness = (
        ", ".join(data["awareness_types"])
        if data["awareness_types"]
        else "nessuna sensibilizzazione mirata"
    )

    if level in ["Alto", "Critico"]:
        decision = (
            "Si raccomanda un coordinamento operativo immediato, "
            "con verifica delle attività critiche, delle interferenze, "
            "della copertura HSE e delle NC prioritarie."
        )

    elif level == "Medio":
        decision = (
            "Si raccomanda di rafforzare il presidio operativo, "
            "ridurre le interferenze e programmare azioni mirate."
        )

    else:
        decision = (
            "Il livello risulta sotto controllo, ma deve essere "
            "mantenuto il monitoraggio al variare delle attività."
        )

    return f"""
Il cantiere **{data['nome']}** presenta un Risk Index pari a
**{round(risk)} / 100**, classificato come **{level}**.

I principali driver negativi sono: **{top_text}**.

Le attività in corso sono: **{activities}**.

Il livello di interferenza dichiarato è **{data['interference_level']}/10**,
con un impatto calcolato pari a
**{technical_data['interference_score']} punti grezzi**.

La copertura HSE è pari a **{data['company_hse']} HSE**
rispetto a **{data['app']} appaltatori**:
**{technical_data['hse_coverage_status']}**.

L'indice di detection è pari a
**{technical_data['nc_detection_ratio']} NC per ispezione**.

Le sensibilizzazioni registrate sono: **{awareness}**.

{decision}
"""


# =========================================================
# ROOT CAUSE
# =========================================================
def local_root_cause(
    data,
    risk,
    level,
    technical_data,
    driver_df=None
):
    categories = {
        "Non conformità": [],
        "Attività": [],
        "Interferenze": [],
        "Organizzazione": [],
        "Controlli": [],
        "Cultura HSE": []
    }

    high_nc = [item for item in data["nc_items"] if item["severity"] == "alta"]
    nc_themes = [item["theme"] for item in data["nc_items"]]

    if data["num_nc"] > 0:
        categories["Non conformità"].append(
            f"Sono presenti {data['num_nc']} NC, di cui {len(high_nc)} ad alta gravità."
        )
    else:
        categories["Non conformità"].append("Non risultano NC registrate nel periodo analizzato.")

    if technical_data["nc_detection_ratio"] >= 0.5:
        categories["Non conformità"].append(
            f"L'indice di detection è significativo ({technical_data['nc_detection_ratio']} NC/ispezione)."
        )

    if data["activities"]:
        categories["Attività"].append(
            "Attività presenti: " + ", ".join(data["activities"]) + "."
        )
    if len(data["activities"]) >= 3:
        categories["Attività"].append(
            "La contemporaneità di almeno tre tipologie di lavoro aumenta la complessità operativa."
        )

    categories["Interferenze"].append(
        f"Il livello dichiarato è {data['interference_level']}/10, con impatto tecnico pari a {technical_data['interference_score']} punti grezzi."
    )
    if data["interference_level"] >= 6:
        categories["Interferenze"].append(
            "È necessario rafforzare la separazione spaziale o temporale e il coordinamento giornaliero."
        )

    categories["Organizzazione"].append(
        f"Sono presenti {data['app']} appaltatori, {data['sub']} subappaltatori e {data['company_hse']} HSE delle imprese."
    )
    categories["Organizzazione"].append(
        f"Valutazione della copertura HSE: {technical_data['hse_coverage_status']}."
    )
    if data["sub"] >= 5:
        categories["Organizzazione"].append(
            "L'elevato numero di subappaltatori richiede una governance più strutturata."
        )

    if data["crit_open"] > 0:
        categories["Controlli"].append(
            f"Sono ancora aperte {data['crit_open']} criticità."
        )
    if data["crit_late"] > 0:
        categories["Controlli"].append(
            f"Sono state chiuse in ritardo {data['crit_late']} criticità, segnale di follow-up da rafforzare."
        )
    if data["inspections"] == 0:
        categories["Controlli"].append("Non risultano ispezioni registrate.")
    else:
        categories["Controlli"].append(f"Sono state registrate {data['inspections']} ispezioni HSE.")

    if data["awareness_count"] > 0:
        categories["Cultura HSE"].append(
            f"Sono state eseguite {data['awareness_count']} sensibilizzazioni."
        )
    else:
        categories["Cultura HSE"].append(
            "Non risultano sensibilizzazioni registrate: manca un elemento preventivo rilevante."
        )
    if data["stopworks"] > 0:
        categories["Cultura HSE"].append(
            f"Sono stati registrati {data['stopworks']} Stop Work, considerati un segnale positivo di partecipazione."
        )

    contribution_map = {
        "Non conformità": ["NC", "Indice di detection", "Correlazioni critiche"],
        "Attività": ["Attività"],
        "Interferenze": ["Interferenza attività"],
        "Organizzazione": ["Complessità organizzativa", "Copertura HSE imprese"],
        "Controlli": ["Criticità", "Bonus/Malus rapporto ispezioni"],
        "Cultura HSE": ["Bonus Stop Work", "Bonus sensibilizzazioni"],
    }

    contributions = {}
    if driver_df is not None:
        for category, drivers in contribution_map.items():
            value = driver_df[driver_df["Driver"].isin(drivers)]["Valore"].sum()
            contributions[category] = round(float(value), 1)
    else:
        contributions = {category: 0.0 for category in categories}

    positive = {key: max(value, 0) for key, value in contributions.items()}
    total_positive = sum(positive.values())
    percentages = {
        key: round((value / total_positive) * 100) if total_positive > 0 else 0
        for key, value in positive.items()
    }

    blocks = []
    for category, items in categories.items():
        impact = contributions.get(category, 0)
        percentage = percentages.get(category, 0)
        blocks.append(f"### {category}\n\n**Contributo netto:** {impact:+.1f} punti | **Quota dei driver negativi:** {percentage}%")
        blocks.extend(f"- {item}" for item in items)

    dominant = max(positive, key=positive.get) if positive else "nessuna categoria"
    conclusion = (
        f"### Conclusione complessiva\n\nIl Risk Index è **{round(risk)}/100 ({level})**. "
        f"La categoria che contribuisce maggiormente ai driver negativi è **{dominant}**. "
        "Le priorità devono essere definite sulla base dei contributi più elevati, senza trascurare "
        "eventuali NC ad alta gravità o condizioni non controllate anche quando il loro peso numerico è limitato."
    )

    return "\n\n".join(blocks + [conclusion])


# =========================================================
# TOOLBOX TALK
# =========================================================
def local_toolbox_talk(
    data,
    technical_data
):
    activities = data["activities"]

    nc_themes = [
        item["theme"]
        for item in data["nc_items"]
    ]

    title_parts = []

    if activities:
        title_parts.extend(activities[:2])

    if nc_themes:
        title_parts.extend(nc_themes[:2])

    title = (
        " / ".join(title_parts)
        if title_parts
        else "Sicurezza operativa giornaliera"
    )

    rules = [
        "Non iniziare attività senza autorizzazione e briefing pre-job.",
        "Verificare le interferenze con le altre imprese prima di iniziare.",
        "Segnalare immediatamente variazioni del programma di lavoro.",
        "Mantenere l'area ordinata, accessibile e segregata.",
        "Utilizzare i DPI previsti per il rischio specifico.",
        "Applicare lo Stop Work in caso di condizione non controllata."
    ]

    if data["interference_level"] >= 6:
        rules.append(
            "Il livello di interferenza è elevato: rispettare "
            "aree, orari e sequenze operative concordate."
        )

    if data["company_hse"] < data["app"]:
        rules.append(
            "In assenza di presidio HSE dedicato, il preposto deve "
            "segnalare immediatamente ogni variazione o criticità."
        )

    if "elettrici" in activities or "elettrico" in nc_themes:
        rules.append(
            "Per le attività elettriche verificare sezionamento, "
            "assenza tensione e procedura LOTO."
        )

    if "scavi" in activities or "scavi" in nc_themes:
        rules.append(
            "Per gli scavi verificare sottoservizi, stabilità "
            "dei fronti e accessi sicuri."
        )

    if "lavori in quota" in activities or "quota" in nc_themes:
        rules.append(
            "Per i lavori in quota verificare le protezioni "
            "anticaduta prima dell'accesso."
        )

    if "sollevamenti" in activities:
        rules.append(
            "Nessuno deve sostare sotto carichi sospesi "
            "o nelle aree di manovra."
        )

    rules_text = "\n".join(
        f"- {rule}"
        for rule in rules
    )

    return f"""
### Titolo

**Toolbox Talk - {title}**

### Dati operativi

- Fase: **{data['fase']}**
- Interferenza: **{data['interference_level']}/10**
- HSE imprese: **{data['company_hse']}**
- Appaltatori: **{data['app']}**
- Indice di detection:
  **{technical_data['nc_detection_ratio']} NC/ispezione**

### Obiettivo

Allineare la squadra sui rischi della giornata,
sulle interferenze e sulle condizioni minime di sicurezza.

### Regole operative

{rules_text}

### Comportamenti vietati

- Iniziare lavorazioni fuori programma senza coordinamento.
- Rimuovere segregazioni o protezioni.
- Accedere alle aree di un'altra impresa senza autorizzazione.
- Proseguire in presenza di un rischio non controllato.

### Domande finali

1. Quali attività interferiscono con la nostra?
2. Chi coordina le lavorazioni di oggi?
3. Qual è il rischio principale?
4. Chi deve essere avvisato in caso di variazione?
5. Quando deve essere applicato lo Stop Work?
"""


# =========================================================
# ACTION PLAN
# =========================================================
def local_action_plan(
    data,
    risk,
    level,
    technical_data
):
    rows = []

    if risk >= 80:
        rows.append([
            "Coordinamento HSE immediato",
            "Critica",
            "Site Manager / HSE",
            "Immediata",
            "Verbale e piano azioni"
        ])

    if technical_data["nc_detection_ratio"] >= 0.5:
        rows.append([
            "Analizzare indice di detection e ricorrenza NC",
            "Alta",
            "HSE Manager",
            "48 ore",
            "Analisi NC per tema e causa"
        ])

    if data["interference_level"] >= 6:
        rows.append([
            "Formalizzare matrice interferenze giornaliera",
            "Alta",
            "Construction Manager",
            "Prima delle attività",
            "Matrice e briefing firmato"
        ])

    if data["interference_level"] >= 8:
        rows.append([
            "Valutare separazione temporale o spaziale delle attività",
            "Critica",
            "Site Manager",
            "Immediata",
            "Programma lavori aggiornato"
        ])

    if data["company_hse"] < data["app"]:
        rows.append([
            "Rafforzare il presidio HSE delle imprese",
            "Alta",
            "Responsabili appaltatori",
            "48 ore",
            "Organigramma e presenze HSE"
        ])

    if data["crit_open"] > 0:
        rows.append([
            "Chiudere le criticità aperte prioritarie",
            "Alta",
            "HSE / Responsabile impresa",
            "24-48 ore",
            "Registro criticità aggiornato"
        ])

    if data["crit_late"] > 0:
        rows.append([
            "Analizzare i ritardi nella chiusura",
            "Media",
            "HSE Manager",
            "7 giorni",
            "Analisi cause e azioni correttive"
        ])

    if "elettrici" in data["activities"]:
        rows.append([
            "Verificare attività elettriche e LOTO",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist e autorizzazioni"
        ])

    if "scavi" in data["activities"]:
        rows.append([
            "Verificare sicurezza degli scavi",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist scavi"
        ])

    if "lavori in quota" in data["activities"]:
        rows.append([
            "Verificare lavori in quota",
            "Alta",
            "Preposto / HSE",
            "24 ore",
            "Checklist quota e PLE"
        ])

    if data["awareness_count"] == 0:
        rows.append([
            "Eseguire toolbox mirato",
            "Alta",
            "HSE / Preposto",
            "Prima dell'attività",
            "Registro presenze"
        ])

    if not rows:
        rows.append([
            "Mantenere il monitoraggio HSE",
            "Media",
            "HSE",
            "7 giorni",
            "Report ispezione"
        ])

    return pd.DataFrame(
        rows,
        columns=[
            "Azione",
            "Priorità",
            "Owner suggerito",
            "Scadenza",
            "Evidenza richiesta"
        ]
    )


# =========================================================
# RIDUZIONE RISCHIO
# =========================================================
def local_reduce_risk(
    data,
    risk,
    technical_data
):
    if risk > 80:
        target = 80
    elif risk > 60:
        target = 60
    elif risk > 40:
        target = 40
    else:
        target = 20

    levers = []

    if technical_data["nc_detection_ratio"] >= 0.25:
        levers.append(
            "Ridurre la frequenza delle NC attraverso analisi "
            "delle cause ricorrenti e azioni preventive."
        )

    if data["interference_level"] >= 6:
        levers.append(
            "Ridurre il livello di interferenza separando "
            "le attività per area, orario o sequenza."
        )

    if data["company_hse"] < data["app"]:
        levers.append(
            "Aumentare la copertura HSE delle imprese "
            "fino ad almeno un rapporto 1:1 con gli appaltatori."
        )

    if data["crit_open"] > 0:
        levers.append(
            "Chiudere le criticità aperte più rilevanti."
        )

    if data["awareness_count"] == 0:
        levers.append(
            "Eseguire toolbox mirati sulle attività critiche."
        )

    if data["num_nc"] > 0:
        levers.append(
            "Chiudere le NC ad alta gravità con verifica di efficacia."
        )

    if data["sub"] >= 5:
        levers.append(
            "Rafforzare il coordinamento dei subappaltatori."
        )

    if not levers:
        levers.append(
            "Mantenere controlli periodici e aggiornare "
            "il rischio al cambio delle attività."
        )

    levers_text = "\n".join(
        f"- {lever}"
        for lever in levers
    )

    return f"""
### Obiettivo suggerito

Portare il Risk Index sotto **{target}**.

### Leve principali

{levers_text}

### Strategia consigliata

Agire prioritariamente su:

- indice di detection;
- interferenza tra attività;
- copertura HSE delle imprese;
- NC ad alta gravità;
- criticità aperte;
- sensibilizzazioni mirate.
"""


# =========================================================
# HSE ADVISOR LOCALE
# =========================================================
def local_hse_advisor(
    question,
    data,
    risk,
    level,
    actions_df,
    technical_data
):
    query = question.lower()

    if (
        "interfer" in query
        or "contemporan" in query
    ):
        return f"""
Il livello di interferenza dichiarato è
**{data['interference_level']}/10**.

L'impatto tecnico calcolato è pari a
**{technical_data['interference_score']} punti grezzi**.

Il valore è corretto in funzione della tipologia
e del numero delle attività in corso.

Le azioni prioritarie sono:

1. predisporre una matrice interferenziale;
2. separare le attività per area o orario;
3. effettuare un briefing tra imprese;
4. aggiornare il programma giornaliero;
5. definire chi coordina ogni interferenza.
"""

    if (
        "hse" in query
        or "appaltator" in query
        or "copertura" in query
    ):
        return f"""
Sono presenti **{data['company_hse']} HSE delle imprese**
rispetto a **{data['app']} appaltatori**.

La valutazione della copertura è:
**{technical_data['hse_coverage_status']}**.

Il correttivo applicato all'indice è pari a
**{technical_data['hse_coverage_adjustment']} punti**.

Un valore positivo è un malus.
Un valore negativo è un bonus.
"""

    if (
        "detection" in query
        or "ispezion" in query
        or "rapporto nc" in query
    ):
        return f"""
L'indice di detection è pari a
**{technical_data['nc_detection_ratio']} NC per ispezione**.

Il parametro rappresenta il rapporto tra il numero
totale delle NC e il numero delle ispezioni.

Più il rapporto cresce, maggiore è il malus,
perché aumenta la frequenza con cui vengono
individuate anomalie durante i controlli.
"""

    if (
        "24" in query
        or "domani" in query
        or "subito" in query
    ):
        top_actions = actions_df.head(5)["Azione"].tolist()

        actions_text = "\n".join(
            f"{index + 1}. {action}"
            for index, action in enumerate(top_actions)
        )

        return f"""
Nelle prossime 24 ore suggerisco:

{actions_text}
"""

    if (
        "principale" in query
        or "rischio" in query
    ):
        top_actions = actions_df.head(3)["Azione"].tolist()

        actions_text = "\n".join(
            f"- {action}"
            for action in top_actions
        )

        return f"""
Il rischio principale deriva dalla combinazione di:

- attività critiche;
- livello di interferenza {data['interference_level']}/10;
- indice di detection
  {technical_data['nc_detection_ratio']};
- copertura HSE
  {data['company_hse']} su {data['app']} appaltatori;
- criticità aperte e NC.

Azioni prioritarie:

{actions_text}
"""

    if (
        "ridurre" in query
        or "abbassare" in query
    ):
        return local_reduce_risk(
            data,
            risk,
            technical_data
        )

    if (
        "bonus stop" in query
        or "stop work" in query
        or "stopwork" in query
    ):
        return f"""
Sono stati registrati **{data['stopworks']} Stop Work**.

Nel modello gli Stop Work sono considerati esclusivamente
come un segnale positivo di cultura della sicurezza:
il bonus applicato è **-{technical_data['stop_work_bonus']} punti**.

Il beneficio è volutamente limitato a un massimo di
**-{STOP_WORK_BONUS_CAP} punti**, così non può compensare
NC, criticità o interferenze rilevanti.
"""

    if (
        "fermare" in query
        or "stop" in query
    ):
        if level in ["Alto", "Critico"]:
            return (
                "Con livello Alto o Critico valuterei la sospensione "
                "selettiva delle attività non controllate, in particolare "
                "in presenza di NC ad alta gravità o interferenze elevate."
            )

        return (
            "Non emerge automaticamente la necessità di fermare "
            "le attività, ma deve essere applicato lo Stop Work "
            "in presenza di una condizione non controllata."
        )

    return f"""
Il report presenta un livello **{level}**
con Risk Index **{round(risk)}/100**.

Per interpretarlo considera prioritariamente:

- NC e relativa gravità;
- indice di detection;
- rapporto tra ispezioni e NC;
- attività e livello di interferenza;
- copertura HSE delle imprese;
- criticità aperte;
- sensibilizzazioni mirate.

Prova a chiedere, ad esempio:

- Qual è il rischio principale?
- Come posso ridurre il rischio?
- La copertura HSE è sufficiente?
- Come viene calcolata l'interferenza?
- Cosa devo fare nelle prossime 24 ore?
"""


# =========================================================
# GRAFICO DRIVER BONUS / MALUS
# =========================================================
def create_driver_chart(driver_df):
    chart_df = driver_df.copy()

    chart_df["Tipo"] = chart_df["Valore"].apply(
        lambda value:
        "Malus"
        if value > 0
        else "Bonus"
        if value < 0
        else "Neutro"
    )

    chart_df["Etichetta"] = chart_df["Valore"].apply(
        lambda value: f"{value:+.1f}"
    )

    bars = (
        alt.Chart(chart_df)
        .mark_bar(cornerRadiusEnd=4)
        .encode(
            x=alt.X(
                "Valore:Q",
                title="Impatto sull'indice di rischio"
            ),
            y=alt.Y(
                "Driver:N",
                sort="-x",
                title=None
            ),
            color=alt.Color(
                "Tipo:N",
                scale=alt.Scale(
                    domain=[
                        "Malus",
                        "Bonus",
                        "Neutro"
                    ],
                    range=[
                        "#D62728",
                        "#2CA02C",
                        "#A0A0A0"
                    ]
                ),
                legend=alt.Legend(
                    title="Impatto"
                )
            ),
            tooltip=[
                alt.Tooltip(
                    "Driver:N",
                    title="Driver"
                ),
                alt.Tooltip(
                    "Valore:Q",
                    title="Impatto",
                    format=".1f"
                ),
                alt.Tooltip(
                    "Tipo:N",
                    title="Tipo"
                )
            ]
        )
    )

    zero_line = (
        alt.Chart(
            pd.DataFrame({"zero": [0]})
        )
        .mark_rule(
            color="#444444",
            strokeWidth=1
        )
        .encode(
            x="zero:Q"
        )
    )

    labels = (
        alt.Chart(chart_df)
        .mark_text(
            align="left",
            baseline="middle",
            dx=5
        )
        .encode(
            x="Valore:Q",
            y=alt.Y(
                "Driver:N",
                sort="-x"
            ),
            text="Etichetta:N"
        )
    )

    return (
        bars + zero_line + labels
    ).properties(
        height=420
    )



# =========================================================
# HSE SCORECARD
# =========================================================
def clamp(value, minimum=0.0, maximum=10.0):
    return max(minimum, min(maximum, value))


def generate_hse_scorecard(data, risk, technical_data):
    """Genera una scorecard sintetica su scala 0-10."""
    organization = 10.0
    organization -= min(data["sub"] * 0.35, 2.0)
    organization -= max(technical_data["hse_coverage_adjustment"], 0) / 6
    organization += min(max(-technical_data["hse_coverage_adjustment"], 0) / 8, 1.0)

    planning = 10.0
    planning -= max(data["interference_level"] - 3, 0) * 0.55
    planning -= 1.2 if len(data["activities"]) >= 3 else 0
    planning -= 0.8 if data["crit_late"] > 0 else 0

    operational = 10.0
    operational -= min(technical_data["nc_detection_ratio"] * 3.2, 4.0)
    operational -= min(data["crit_open"] * 0.45, 2.5)
    operational -= min(sum(1 for item in data["nc_items"] if item["severity"] == "alta") * 0.8, 2.0)

    culture = 6.0
    culture += min(data["awareness_count"] * 0.25, 2.0)
    culture += min(data["stopworks"] * 0.30, 1.5)
    culture += min(len(data["awareness_types"]) * 0.18, 0.8)
    culture -= 1.0 if data["awareness_count"] == 0 else 0

    coordination = 10.0
    coordination -= max(data["interference_level"] - 2, 0) * 0.55
    coordination -= min(max(len(data["activities"]) - 2, 0) * 0.45, 1.8)
    coordination -= 0.8 if data["sub"] >= 5 else 0

    rows = [
        ("Organizzazione", organization, "Adeguatezza della struttura, complessità imprese e copertura HSE."),
        ("Pianificazione", planning, "Programmazione delle attività, gestione delle interferenze e puntualità del follow-up."),
        ("Controllo operativo", operational, "Esito dei controlli, frequenza delle NC e criticità ancora aperte."),
        ("Cultura HSE", culture, "Sensibilizzazioni, partecipazione e utilizzo responsabile dello Stop Work."),
        ("Coordinamento", coordination, "Gestione della contemporaneità, delle imprese e delle lavorazioni interferenti."),
    ]

    scorecard_df = pd.DataFrame(rows, columns=["Area", "Punteggio", "Motivazione"])
    scorecard_df["Punteggio"] = scorecard_df["Punteggio"].apply(lambda value: round(clamp(value), 1))
    scorecard_df["Stelle"] = scorecard_df["Punteggio"].apply(
        lambda value: "★" * max(1, min(5, round(value / 2))) + "☆" * (5 - max(1, min(5, round(value / 2))))
    )
    overall = round(scorecard_df["Punteggio"].mean(), 1)
    return scorecard_df, overall


# =========================================================
# HSE TOOLKIT - CHECKLIST WORD
# =========================================================
CHECKLIST_LIBRARY = {
    "Controlli generali pre-attività": [
        "L'attività è inserita nel programma lavori / sinottico aggiornato?",
        "È stato eseguito il coordinamento con le imprese coinvolte e interferenti?",
        "Sono stati individuati il preposto e i referenti operativi dell'attività?",
        "La documentazione operativa è disponibile, aggiornata e coerente con l'attività?",
        "Sono state verificate formazione, informazione, addestramento e idoneità dei lavoratori?",
        "Sono disponibili attrezzature, mezzi e DPI idonei e verificati?",
        "L'area di lavoro è stata delimitata, segnalata e resa accessibile in sicurezza?",
        "Sono state valutate le condizioni ambientali e meteorologiche?",
        "Sono state definite le modalità di comunicazione e gestione delle emergenze?",
        "È stato eseguito il briefing / pre-job check prima dell'avvio?",
    ],
    "Sollevamenti": [
        "È stato predisposto e condiviso il piano di sollevamento?",
        "Sono stati verificati peso, baricentro e punti di presa del carico?",
        "La gru o l'apparecchio di sollevamento è idoneo e correttamente posizionato?",
        "Sono disponibili verifiche, controlli e documentazione dell'attrezzatura?",
        "Gruista, imbracatore e segnalatore possiedono formazione e abilitazioni richieste?",
        "Gli accessori di sollevamento sono identificati, integri e adeguati alla portata?",
        "L'area di manovra e la zona sottostante sono segregate?",
        "Sono state verificate interferenze con linee elettriche, strutture e altre attività?",
        "Sono definite modalità di comunicazione univoche tra gli operatori?",
        "Sono rispettati i limiti meteo e di vento previsti?",
    ],
    "Lavori elettrici": [
        "Sono definiti ruoli, responsabilità e autorizzazioni per il lavoro elettrico?",
        "Sono verificate le competenze PES, PAV, PEI o equivalenti richieste?",
        "È disponibile il permesso / piano di lavoro applicabile?",
        "Sono state eseguite sezionamento, blocco, identificazione e prevenzione delle richiusure?",
        "È stata verificata l'assenza di tensione con strumento idoneo?",
        "Sono state applicate eventuali terre e cortocircuiti previste?",
        "Sono definite e rispettate le distanze di sicurezza dalle parti attive?",
        "DPI, attrezzi isolati e strumenti sono idonei e verificati?",
        "Sono gestite le interferenze con impianti o stalli adiacenti in esercizio?",
        "Sono definite le procedure di emergenza e soccorso elettrico?",
    ],
    "Scavi": [
        "Sono stati individuati e tracciati i sottoservizi presenti?",
        "È disponibile l'autorizzazione allo scavo e la documentazione prevista?",
        "Sono state valutate stabilità dei fronti, profondità e caratteristiche del terreno?",
        "Sono previste armature, pendenze o sistemi di protezione adeguati?",
        "Sono presenti accessi e uscite sicure dallo scavo?",
        "Materiali e mezzi sono mantenuti a distanza sicura dal ciglio?",
        "Lo scavo è delimitato e segnalato anche nelle ore notturne?",
        "Sono gestite acqua, cedimenti, atmosfere pericolose e condizioni meteo?",
        "Sono impediti accessi e soste nel raggio operativo dei mezzi?",
        "È prevista una verifica prima di ogni turno e dopo eventi che possano alterare lo scavo?",
    ],
    "Lavori in quota": [
        "È stata privilegiata una protezione collettiva rispetto ai DPI anticaduta?",
        "Ponteggi, parapetti, trabattelli o opere provvisionali sono conformi e verificati?",
        "Sono individuati e verificati i punti di ancoraggio?",
        "Imbracature, cordini e dispositivi anticaduta sono integri e idonei?",
        "I lavoratori hanno formazione e addestramento specifici?",
        "Sono gestiti caduta materiali e segregazione dell'area sottostante?",
        "Scale e accessi sono utilizzati esclusivamente nelle condizioni consentite?",
        "Sono considerate condizioni meteo e vento?",
        "È disponibile un piano di recupero e soccorso dell'operatore sospeso?",
        "Sono state verificate interferenze con impianti elettrici e altre lavorazioni?",
    ],
    "PLE": [
        "La PLE è idonea alla quota, allo sbraccio e alle condizioni del terreno?",
        "Sono disponibili verifiche periodiche, controlli e manuale d'uso?",
        "L'operatore è abilitato e autorizzato all'utilizzo?",
        "È stata eseguita la verifica pre-uso?",
        "Stabilizzatori e piastre sono correttamente posizionati?",
        "Sono rispettati i limiti di portata e numero di occupanti?",
        "Sono gestite distanze da linee elettriche e ostacoli?",
        "L'area sottostante e il raggio di manovra sono segregati?",
        "È previsto l'utilizzo dei DPI anticaduta indicati dal costruttore?",
        "È noto il sistema di discesa di emergenza e presente personale addestrato a terra?",
    ],
    "Hot Work": [
        "È stato emesso il permesso di lavoro a caldo?",
        "Sono stati rimossi o protetti i materiali combustibili?",
        "Sono disponibili estintori idonei e facilmente raggiungibili?",
        "È stato nominato il fire watch quando necessario?",
        "Bombole, tubazioni e riduttori sono integri e correttamente posizionati?",
        "È garantita adeguata ventilazione?",
        "Sono controllate scintille, scorie e propagazione del calore verso aree adiacenti?",
        "Sono gestite atmosfere potenzialmente esplosive o infiammabili?",
        "I lavoratori utilizzano DPI idonei a calore, radiazioni e fumi?",
        "È previsto il controllo dell'area al termine dell'attività?",
    ],
    "Spazi confinati": [
        "L'ambiente è stato classificato e autorizzato come spazio confinato / sospetto di inquinamento?",
        "È disponibile una procedura specifica e un permesso di ingresso?",
        "Sono stati isolati impianti, energie e possibili ingressi di sostanze?",
        "Sono state eseguite misure atmosferiche prima e durante l'accesso?",
        "Sono garantite ventilazione e condizioni respirabili?",
        "Gli addetti possiedono formazione e addestramento specifici?",
        "È presente un sorvegliante esterno dedicato?",
        "Sono disponibili comunicazioni efficaci tra interno ed esterno?",
        "È predisposto un piano di emergenza con attrezzature di recupero?",
        "È impedito l'accesso a personale non autorizzato?",
    ],
    "Movimentazione mezzi": [
        "È definito e segnalato il piano di viabilità del cantiere?",
        "Sono separati, ove possibile, percorsi pedonali e carrabili?",
        "Mezzi e operatori sono autorizzati e in possesso dei requisiti previsti?",
        "Sono stati eseguiti i controlli pre-uso dei mezzi?",
        "Sono funzionanti segnalatori acustici, luminosi e dispositivi di visibilità?",
        "Sono individuati punti ciechi e necessità di moviere?",
        "Velocità e sensi di marcia sono definiti e rispettati?",
        "Le aree di carico, scarico e manovra sono segregate?",
        "Sono gestite interferenze con scavi, sollevamenti e lavorazioni a terra?",
        "Le condizioni del fondo e delle piste sono adeguate?",
    ],
    "Attività civili": [
        "Le aree di lavoro sono ordinate, pulite e prive di ostacoli?",
        "Sono verificati accessi, percorsi e protezioni dei bordi?",
        "Attrezzature manuali ed elettriche sono idonee e controllate?",
        "Sono gestite polveri, rumore e proiezione di materiali?",
        "Sono previste modalità sicure per taglio, perforazione e demolizioni minori?",
        "Sono stoccati correttamente materiali e prodotti?",
        "Sono controllate le interferenze con impianti esistenti e sottoservizi?",
        "Sono disponibili i DPI specifici per le attività?",
    ],
    "Attività meccaniche": [
        "Macchine e componenti sono stabilizzati prima delle lavorazioni?",
        "Sono isolate e scaricate le energie residue?",
        "Attrezzi e attrezzature sono adatti e in buono stato?",
        "Sono gestiti schiacciamento, cesoiamento e intrappolamento?",
        "Sono previsti sistemi sicuri per movimentazione e posizionamento dei componenti?",
        "Sono controllati fluidi in pressione e temperature residue?",
        "Sono definite le modalità di prova e rimessa in servizio?",
        "L'area è segregata durante prove, avviamenti e movimentazioni?",
    ],
}

ACTIVITY_TO_CHECKLIST = {
    "elettrici": "Lavori elettrici",
    "scavi": "Scavi",
    "lavori in quota": "Lavori in quota",
    "sollevamenti": "Sollevamenti",
    "spazi confinati": "Spazi confinati",
    "hot work": "Hot Work",
    "movimentazione mezzi": "Movimentazione mezzi",
    "civili": "Attività civili",
    "meccanici": "Attività meccaniche",
}


def generate_checklist_word(data, selected_sections):
    document = Document()
    section = document.sections[0]
    section.top_margin = DocxInches(0.65)
    section.bottom_margin = DocxInches(0.65)
    section.left_margin = DocxInches(0.7)
    section.right_margin = DocxInches(0.7)

    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(9)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CHECKLIST HSE PRE-ATTIVITÀ")
    run.bold = True
    run.font.size = Pt(16)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Verifica preventiva da compilare a cura dell'HSE di cantiere").italic = True

    info_table = document.add_table(rows=3, cols=4)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    info_table.style = "Table Grid"
    info_values = [
        ("Cantiere", data.get("nome", ""), "Data", datetime.now().strftime("%d/%m/%Y")),
        ("Fase", data.get("fase", ""), "Impresa", ""),
        ("Compilata da", "", "Attività / area", ""),
    ]
    for row, values in zip(info_table.rows, info_values):
        for index, value in enumerate(values):
            row.cells[index].text = str(value)
            row.cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            if index in (0, 2):
                row.cells[index].paragraphs[0].runs[0].bold = True

    document.add_paragraph()
    legend = document.add_paragraph()
    legend.add_run("Modalità di compilazione: ").bold = True
    legend.add_run("barrare SÌ, NO oppure N.A. e riportare nelle note le azioni necessarie.")

    for section_name in selected_sections:
        questions = CHECKLIST_LIBRARY.get(section_name, [])
        if not questions:
            continue

        heading = document.add_paragraph()
        heading.paragraph_format.space_before = Pt(10)
        heading.paragraph_format.space_after = Pt(4)
        run = heading.add_run(section_name.upper())
        run.bold = True
        run.font.size = Pt(11)

        table = document.add_table(rows=1, cols=6)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = "Table Grid"
        headers = ["N.", "Verifica", "SÌ", "NO", "N.A.", "Note / azioni"]
        for index, header in enumerate(headers):
            cell = table.rows[0].cells[index]
            cell.text = header
            cell.paragraphs[0].runs[0].bold = True
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        for number, question in enumerate(questions, start=1):
            cells = table.add_row().cells
            cells[0].text = str(number)
            cells[1].text = question
            cells[2].text = "☐"
            cells[3].text = "☐"
            cells[4].text = "☐"
            cells[5].text = ""
            for cell in cells:
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        widths = [0.35, 4.55, 0.45, 0.45, 0.5, 1.45]
        for row in table.rows:
            for index, width in enumerate(widths):
                row.cells[index].width = DocxInches(width)

    document.add_paragraph()
    notes = document.add_paragraph()
    notes.add_run("NOTE GENERALI / AZIONI DA INTRAPRENDERE").bold = True
    for _ in range(4):
        document.add_paragraph("________________________________________________________________________________")

    signatures = document.add_table(rows=2, cols=3)
    signatures.alignment = WD_TABLE_ALIGNMENT.CENTER
    signatures.style = "Table Grid"
    for index, label in enumerate(["HSE di cantiere", "Preposto", "Responsabile impresa"]):
        signatures.cell(0, index).text = label
        signatures.cell(0, index).paragraphs[0].runs[0].bold = True
        signatures.cell(1, index).text = "\nFirma: __________________________\n"

    footer = document.sections[0].footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("Checklist HSE generata tramite HSE Risk Platform")

    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer


# =========================================================
# POWERPOINT
# =========================================================
def generate_ppt(
    data,
    risk,
    level,
    driver_df,
    action_plan_df,
    summary,
    root_cause,
    toolbox,
    technical_data,
    scorecard_df,
    scorecard_overall
):
    prs = Presentation()

    slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    slide.shapes.title.text = "HSE Risk Report"

    slide.placeholders[1].text = (
        f"Cantiere: {data['nome']}\n"
        f"Data: {datetime.now().strftime('%d/%m/%Y')}"
    )

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Executive Summary"

    slide.placeholders[1].text = (
        summary
        .replace("**", "")
        .replace("###", "")
    )[:1700]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Sintesi del rischio"

    slide.placeholders[1].text = (
        f"Risk Index: {round(risk)} / 100\n"
        f"Livello: {level}\n\n"
        f"Fase: {data['fase']}\n"
        f"Indice di detection: "
        f"{technical_data['nc_detection_ratio']}\n"
        f"Interferenza: {data['interference_level']}/10\n"
        f"HSE imprese: {data['company_hse']}\n"
        f"Appaltatori: {data['app']}\n"
        f"Copertura HSE: "
        f"{technical_data['hse_coverage_status']}\n"
        f"Attività: "
        f"{', '.join(data['activities']) if data['activities'] else 'Nessuna'}"
    )

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = "Driver del rischio"

    table = slide.shapes.add_table(
        len(driver_df) + 1,
        3,
        Inches(0.6),
        Inches(1.3),
        Inches(8.5),
        Inches(5)
    ).table

    table.cell(0, 0).text = "Driver"
    table.cell(0, 1).text = "Impatto"
    table.cell(0, 2).text = "Tipo"

    for index, row in driver_df.iterrows():
        value = float(row["Valore"])

        impact_type = (
            "Malus"
            if value > 0
            else "Bonus"
            if value < 0
            else "Neutro"
        )

        table.cell(
            index + 1,
            0
        ).text = str(row["Driver"])

        table.cell(
            index + 1,
            1
        ).text = f"{value:+.1f}"

        table.cell(
            index + 1,
            2
        ).text = impact_type

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Action Plan"

    action_text = ""

    for _, row in action_plan_df.head(8).iterrows():
        action_text += (
            f"{row['Priorità']} - "
            f"{row['Azione']} - "
            f"{row['Scadenza']}\n"
        )

    slide.placeholders[1].text = action_text[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Root Cause Analysis"

    slide.placeholders[1].text = (
        root_cause
        .replace("###", "")
        .replace("**", "")
    )[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = "Toolbox Talk"

    slide.placeholders[1].text = (
        toolbox
        .replace("###", "")
        .replace("**", "")
    )[:1800]

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = "HSE Scorecard"

    score_table = slide.shapes.add_table(
        len(scorecard_df) + 2,
        4,
        Inches(0.6),
        Inches(1.25),
        Inches(8.6),
        Inches(4.8)
    ).table

    score_table.cell(0, 0).text = "Area"
    score_table.cell(0, 1).text = "Stelle"
    score_table.cell(0, 2).text = "Punteggio"
    score_table.cell(0, 3).text = "Motivazione"

    for index, row in scorecard_df.iterrows():
        score_table.cell(index + 1, 0).text = str(row["Area"])
        score_table.cell(index + 1, 1).text = str(row["Stelle"])
        score_table.cell(index + 1, 2).text = f"{float(row['Punteggio']):.1f} / 10"
        score_table.cell(index + 1, 3).text = str(row["Motivazione"])

    last_row = len(scorecard_df) + 1
    score_table.cell(last_row, 0).text = "VALUTAZIONE COMPLESSIVA"
    score_table.cell(last_row, 1).text = ""
    score_table.cell(last_row, 2).text = f"{scorecard_overall:.1f} / 10"
    score_table.cell(last_row, 3).text = "Media delle cinque aree della scorecard."

    return prs


# =========================================================
# EXCEL
# =========================================================
def generate_excel(
    data,
    risk,
    level,
    driver_df,
    actions_df,
    action_plan_df,
    explanations_df,
    summary,
    root_cause,
    toolbox,
    reduce_risk,
    technical_data,
    scorecard_df,
    scorecard_overall
):
    buffer = BytesIO()

    summary_df = pd.DataFrame([{
        "Data": datetime.now().strftime("%d/%m/%Y %H:%M"),
        "Cantiere": data["nome"],
        "Risk Index": round(risk, 1),
        "Livello": level,
        "Fase": data["fase"],
        "Ispezioni": data["inspections"],
        "NC": data["num_nc"],
        "Indice di detection":
            technical_data["nc_detection_ratio"],
        "Interferenza attività":
            data["interference_level"],
        "Interference score":
            technical_data["interference_score"],
        "Appaltatori": data["app"],
        "Subappaltatori": data["sub"],
        "HSE imprese": data["company_hse"],
        "Copertura HSE":
            technical_data["hse_coverage_status"],
        "Correttivo copertura HSE":
            technical_data["hse_coverage_adjustment"],
        "Stop Work": data["stopworks"],
        "Criticità aperte": data["crit_open"],
        "Criticità in tempo": data["crit_ontime"],
        "Criticità in ritardo": data["crit_late"],
        "Sensibilizzazioni":
            data["awareness_count"],
        "Tipologie attività":
            ", ".join(data["activities"]),
        "Tipologie sensibilizzazioni":
            ", ".join(data["awareness_types"])
    }])

    local_outputs_df = pd.DataFrame([
        {
            "Sezione": "Executive Summary",
            "Testo": summary
        },
        {
            "Sezione": "Root Cause Analysis",
            "Testo": root_cause
        },
        {
            "Sezione": "Toolbox Talk",
            "Testo": toolbox
        },
        {
            "Sezione": "Come ridurre il rischio",
            "Testo": reduce_risk
        }
    ])

    nc_df = pd.DataFrame(
        data["nc_items"]
    )

    with pd.ExcelWriter(
        buffer,
        engine="openpyxl"
    ) as writer:
        summary_df.to_excel(
            writer,
            sheet_name="Sintesi",
            index=False
        )

        driver_df.to_excel(
            writer,
            sheet_name="Driver rischio",
            index=False
        )

        actions_df.to_excel(
            writer,
            sheet_name="Azioni base",
            index=False
        )

        action_plan_df.to_excel(
            writer,
            sheet_name="Action Plan",
            index=False
        )

        explanations_df.to_excel(
            writer,
            sheet_name="Motivazione",
            index=False
        )

        local_outputs_df.to_excel(
            writer,
            sheet_name="Output locali",
            index=False
        )

        nc_df.to_excel(
            writer,
            sheet_name="NC dettaglio",
            index=False
        )

        scorecard_export_df = scorecard_df.copy()
        scorecard_export_df.loc[len(scorecard_export_df)] = {
            "Area": "Valutazione complessiva",
            "Punteggio": scorecard_overall,
            "Motivazione": "Media delle cinque aree della scorecard",
            "Stelle": ""
        }
        scorecard_export_df.to_excel(
            writer,
            sheet_name="HSE Scorecard",
            index=False
        )

    buffer.seek(0)

    return buffer


# =========================================================
# ESECUZIONE CALCOLO E SESSION STATE
# =========================================================
def run_calculation(data):
    (
        risk,
        level,
        details,
        explanations,
        technical_data
    ) = calculate_risk(data)

    base_actions = generate_base_actions(
        data,
        risk,
        level,
        technical_data
    )

    driver_df = pd.DataFrame(
        details,
        columns=[
            "Driver",
            "Valore"
        ]
    )

    actions_df = pd.DataFrame(
        base_actions,
        columns=[
            "Priorità",
            "Azione"
        ]
    )

    explanations_df = pd.DataFrame(
        explanations,
        columns=[
            "Motivazione"
        ]
    )

    summary = local_executive_summary(
        data,
        risk,
        level,
        driver_df,
        technical_data
    )

    root_cause = local_root_cause(
        data,
        risk,
        level,
        technical_data,
        driver_df
    )

    toolbox = local_toolbox_talk(
        data,
        technical_data
    )

    action_plan_df = local_action_plan(
        data,
        risk,
        level,
        technical_data
    )

    reduce_risk = local_reduce_risk(
        data,
        risk,
        technical_data
    )

    scorecard_df, scorecard_overall = generate_hse_scorecard(
        data,
        risk,
        technical_data
    )

    st.session_state.report = {
        "data": data,
        "risk": risk,
        "level": level,
        "driver_df": driver_df,
        "actions_df": actions_df,
        "explanations_df": explanations_df,
        "summary": summary,
        "root_cause": root_cause,
        "toolbox": toolbox,
        "action_plan_df": action_plan_df,
        "reduce_risk": reduce_risk,
        "technical_data": technical_data,
        "scorecard_df": scorecard_df,
        "scorecard_overall": scorecard_overall
    }

    st.session_state.calcolo_effettuato = True
    st.session_state.advisor_answer = ""


if "calcolo_effettuato" not in st.session_state:
    st.session_state.calcolo_effettuato = False

if "advisor_answer" not in st.session_state:
    st.session_state.advisor_answer = ""


# =========================================================
# INTERFACCIA
# =========================================================
st.title("🦺 HSE Risk Platform AI Locale")

st.caption(
    "Risk assessment HSE con Executive Summary, Root Cause, "
    "Toolbox Talk, Action Plan e Advisor locale gratuito."
)

st.success(
    "Modalità gratuita attiva: nessuna API key richiesta."
)


with st.sidebar:
    st.header("📋 Dati cantiere")

    nome = st.text_input(
        "Nome cantiere",
        value="Cantiere Demo",
        help=(
            "Identifica il sito o il progetto. Il nome viene riportato "
            "nei report PowerPoint, Excel e nella checklist Word."
        )
    )

    fase = st.selectbox(
        "Fase",
        [
            "cantierizzazione",
            "costruzione",
            "commissioning",
            "punch list"
        ],
        help=(
            "La fase non aumenta direttamente il Risk Index. "
            "Serve per contestualizzare gli output e premiare "
            "le sensibilizzazioni coerenti con il momento operativo."
        )
    )

    st.divider()

    # -----------------------------------------------------
    # CONTROLLI E NC
    # -----------------------------------------------------
    st.header("🔍 Controlli e NC")

    inspections = st.number_input(
        "Numero ispezioni",
        min_value=0,
        max_value=1000,
        value=10,
        help=(
            "Numero di ispezioni HSE eseguite nel periodo analizzato. "
            "Influisce sull'indice di detection e sul rapporto ispezioni/NC."
        )
    )

    num_nc = st.number_input(
        "Numero totale NC",
        min_value=0,
        max_value=300,
        value=0,
        help=(
            "Numero complessivo di non conformità rilevate. "
            "Le NC dettagliate ricevono un peso per tema e gravità; "
            "quelle non dettagliate ricevono un peso prudenziale medio."
        )
    )

    nc_items = []

    nc_count_detail = st.number_input(
        "Quante NC vuoi dettagliare?",
        min_value=0,
        max_value=20,
        value=min(num_nc, 3),
        help=(
            "Permette di specificare tema e gravità delle NC principali. "
            "Il numero non può superare il totale delle NC inserite."
        )
    )

    for index in range(nc_count_detail):
        col1, col2 = st.columns(2)

        with col1:
            theme = st.selectbox(
                f"Tema NC {index + 1}",
                list(NC_THEME_WEIGHTS.keys()),
                key=f"nc_theme_{index}",
                help=(
                    "Il tema determina il peso tecnico della NC. "
                    "Temi ad alto potenziale, come elettrico o spazi "
                    "confinati, hanno un impatto maggiore."
                )
            )

        with col2:
            severity = st.selectbox(
                f"Gravità NC {index + 1}",
                [
                    "bassa",
                    "media",
                    "alta"
                ],
                key=f"nc_severity_{index}",
                help=(
                    "La gravità integra il peso del tema: bassa +3, "
                    "media +8, alta +16 punti grezzi."
                )
            )

        nc_items.append({
            "theme": theme,
            "severity": severity
        })

    st.divider()

    # -----------------------------------------------------
    # STOP WORK E CRITICITÀ
    # -----------------------------------------------------
    st.header("⛔ Stop Work e criticità")

    stopworks = st.number_input(
        "Stop Work",
        min_value=0,
        max_value=100,
        value=0,
        help=(
            "Gli Stop Work sono considerati esclusivamente come bonus "
            "di cultura della sicurezza: -1,5 punti ciascuno, con un "
            "massimo di -6 punti complessivi."
        )
    )

    crit_open = st.number_input(
        "Criticità aperte",
        min_value=0,
        max_value=200,
        value=0,
        help=(
            "Criticità ancora da risolvere. Hanno il peso maggiore "
            "nel driver criticità perché rappresentano esposizioni attive."
        )
    )

    crit_ontime = st.number_input(
        "Criticità risolte in tempo",
        min_value=0,
        max_value=200,
        value=0,
        help=(
            "Criticità chiuse entro la scadenza. Mantengono un peso "
            "residuo limitato per rappresentare l'evento rilevato."
        )
    )

    crit_late = st.number_input(
        "Criticità risolte in ritardo",
        min_value=0,
        max_value=200,
        value=0,
        help=(
            "Criticità chiuse oltre la scadenza. Generano un malus "
            "per evidenziare debolezze nel follow-up."
        )
    )

    st.divider()

    # -----------------------------------------------------
    # ORGANIZZAZIONE
    # -----------------------------------------------------
    st.header("🏗️ Organizzazione")

    app = st.number_input(
        "Numero appaltatori",
        min_value=0,
        max_value=100,
        value=1,
        help=(
            "Numero di imprese appaltatrici presenti. Contribuisce "
            "alla complessità organizzativa e al calcolo della copertura HSE."
        )
    )

    sub = st.number_input(
        "Numero subappaltatori",
        min_value=0,
        max_value=100,
        value=0,
        help=(
            "Numero di imprese subappaltatrici presenti. Ha un peso "
            "maggiore degli appaltatori nella complessità organizzativa."
        )
    )

    company_hse = st.number_input(
        "Numero HSE delle imprese",
        min_value=0,
        max_value=100,
        value=1,
        help=(
            "Il valore viene confrontato con il numero "
            "degli appaltatori. Rapporto 1:1 neutro; "
            "copertura inferiore genera malus; "
            "copertura superiore genera bonus."
        )
    )

    st.divider()

    # -----------------------------------------------------
    # ATTIVITÀ
    # -----------------------------------------------------
    st.header("⚙️ Attività in corso")

    activities = st.multiselect(
        "Tipologie attività presenti",
        list(ACTIVITY_WEIGHTS.keys()),
        default=["civili"],
        help=(
            "Seleziona tutte le lavorazioni effettivamente in corso. "
            "Tipologia, numero e contemporaneità incidono sul rischio "
            "e sulle correlazioni con NC e sensibilizzazioni."
        )
    )

    interference_level = st.slider(
        "Livello di interferenza delle attività",
        min_value=1,
        max_value=10,
        value=3,
        step=1,
        help=(
            "1 indica interferenze minime; "
            "10 indica interferenze molto elevate. "
            "L'impatto è corretto in base alla tipologia "
            "e al numero delle attività selezionate."
        )
    )

    if interference_level <= 3:
        st.caption("🟢 Interferenza bassa")
    elif interference_level <= 5:
        st.caption("🟡 Interferenza moderata")
    elif interference_level <= 7:
        st.caption("🟠 Interferenza elevata")
    else:
        st.caption("🔴 Interferenza molto elevata")

    st.divider()

    # -----------------------------------------------------
    # SENSIBILIZZAZIONI
    # -----------------------------------------------------
    st.header("🧠 Sensibilizzazioni")

    awareness_count = st.number_input(
        "Numero sensibilizzazioni effettuate",
        min_value=0,
        max_value=1000,
        value=0,
        help=(
            "Numero complessivo di toolbox o sensibilizzazioni eseguite "
            "nel periodo. La continuità genera un bonus limitato."
        )
    )

    awareness_types = st.multiselect(
        "Tipologie sensibilizzazioni effettuate",
        list(AWARENESS_BONUS.keys()),
        help=(
            "La tipologia determina il bonus e permette di verificare "
            "la coerenza con le attività presenti e con la fase del cantiere."
        )
    )

    st.divider()

    calcola = st.button(
        "Calcola rischio",
        type="primary",
        use_container_width=True
    )


# =========================================================
# AVVIO DEL CALCOLO
# =========================================================
if calcola:
    data = {
        "nome": nome,
        "fase": fase,
        "inspections": inspections,
        "num_nc": num_nc,
        "nc_items": nc_items,
        "stopworks": stopworks,
        "crit_open": crit_open,
        "crit_ontime": crit_ontime,
        "crit_late": crit_late,
        "app": app,
        "sub": sub,
        "company_hse": company_hse,
        "activities": normalize_list(
            activities
        ),
        "interference_level": interference_level,
        "awareness_count": awareness_count,
        "awareness_types": normalize_list(
            awareness_types
        )
    }

    run_calculation(data)


if not st.session_state.calcolo_effettuato:
    st.info(
        "Compila i dati nella sidebar e clicca su "
        "**Calcola rischio**."
    )
    st.stop()


# =========================================================
# LETTURA REPORT SALVATO
# =========================================================
report = st.session_state.report

data = report["data"]
risk = report["risk"]
level = report["level"]
driver_df = report["driver_df"]
actions_df = report["actions_df"]
explanations_df = report["explanations_df"]
summary = report["summary"]
root_cause = report["root_cause"]
toolbox = report["toolbox"]
action_plan_df = report["action_plan_df"]
reduce_risk = report["reduce_risk"]
technical_data = report["technical_data"]
scorecard_df = report["scorecard_df"]
scorecard_overall = report["scorecard_overall"]


# =========================================================
# KPI
# =========================================================
col1, col2, col3, col4, col5 = st.columns(5)

col1.metric(
    "Risk Index",
    f"{round(risk)} / 100"
)

col2.metric(
    "Livello",
    level
)

col3.metric(
    "Indice detection",
    technical_data["nc_detection_ratio"]
)

col4.metric(
    "Interferenza",
    f"{data['interference_level']} / 10"
)

col5.metric(
    "Copertura HSE",
    f"{data['company_hse']} / {data['app']}"
)

st.divider()


# =========================================================
# GRAFICO E SUMMARY
# =========================================================
left, right = st.columns(
    [1.3, 1]
)

with left:
    st.subheader("📊 Impatto dei driver sul Risk Index")

    st.caption(
        "Le colonne rosse aumentano il rischio. "
        "Le colonne verdi rappresentano bonus che riducono l'indice."
    )

    chart = create_driver_chart(
        driver_df
    )

    st.altair_chart(
        chart,
        use_container_width=True
    )

with right:
    st.subheader("🤖 Executive Summary locale")

    st.markdown(summary)

    if (
        data["awareness_count"] > 0
        and not data["awareness_types"]
    ):
        st.warning(
            "Hai inserito sensibilizzazioni senza tipologia. "
            "Seleziona le tipologie per valorizzarle nel calcolo."
        )

    if (
        data["num_nc"]
        > len(data["nc_items"])
    ):
        st.warning(
            f"Hai indicato {data['num_nc']} NC totali, "
            f"ma ne hai dettagliate solo "
            f"{len(data['nc_items'])}. "
            "Alle NC non dettagliate viene applicato "
            "un peso prudenziale medio."
        )

    if data["company_hse"] < data["app"]:
        st.error(
            "La copertura HSE delle imprese è inferiore "
            "al numero degli appaltatori."
        )

    if data["interference_level"] >= 8:
        st.error(
            "Il livello di interferenza dichiarato "
            "è molto elevato."
        )


st.divider()


# =========================================================
# TAB
# =========================================================
tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10 = st.tabs([
    "🎯 Azioni e alert",
    "📋 Action Plan",
    "🔍 Root Cause",
    "🧰 Toolbox Talk",
    "📉 Riduci rischio",
    "🧮 Dettaglio calcolo",
    "💬 HSE Advisor",
    "⭐ HSE Scorecard",
    "📑 HSE Toolkit",
    "📤 Export"
])

with tab1:
    st.subheader("Azioni e alert generati dal motore")
    st.dataframe(actions_df, use_container_width=True, hide_index=True)

with tab2:
    st.subheader("Action Plan operativo")
    st.dataframe(action_plan_df, use_container_width=True, hide_index=True)

with tab3:
    st.subheader("Root Cause Analysis per categorie")
    st.caption("La motivazione del punteggio è organizzata per area e collegata ai contributi effettivi dei driver.")
    st.markdown(root_cause)

with tab4:
    st.subheader("Toolbox Talk locale")
    st.markdown(toolbox)

with tab5:
    st.subheader("Come ridurre il rischio")
    st.markdown(reduce_risk)

with tab6:
    st.subheader("Motivazione tecnica del calcolo")
    st.dataframe(explanations_df, use_container_width=True, hide_index=True)

    st.subheader("Indicatori tecnici")
    technical_df = pd.DataFrame([
        {"Indicatore": "Indice di detection", "Valore": technical_data["nc_detection_ratio"], "Descrizione": "Rapporto tra NC totali e ispezioni"},
        {"Indicatore": "Interference score", "Valore": technical_data["interference_score"], "Descrizione": "Interferenza corretta per tipo e numero attività"},
        {"Indicatore": "Fattore criticità attività", "Valore": technical_data["activity_criticality_factor"], "Descrizione": "Peso dell'attività più critica"},
        {"Indicatore": "Fattore numero attività", "Valore": technical_data["activity_count_factor"], "Descrizione": "Correttivo per contemporaneità"},
        {"Indicatore": "Correttivo copertura HSE", "Valore": technical_data["hse_coverage_adjustment"], "Descrizione": technical_data["hse_coverage_status"]},
        {"Indicatore": "Bonus Stop Work", "Valore": -technical_data["stop_work_bonus"], "Descrizione": "Solo bonus, con impatto massimo limitato a -6"},
        {"Indicatore": "Impatto diretto fase", "Valore": technical_data["phase_direct_impact"], "Descrizione": "La fase non incide direttamente sul Risk Index"},
        {"Indicatore": "Bonus coerenza fase/sensibilizzazioni", "Valore": -technical_data["phase_awareness_bonus"], "Descrizione": ", ".join(technical_data["phase_awareness_matches"]) if technical_data["phase_awareness_matches"] else "Nessuna corrispondenza"}
    ])
    st.dataframe(technical_df, use_container_width=True, hide_index=True)

with tab7:
    st.subheader("Chiedi all'HSE Advisor locale")
    question = st.text_input(
        "Fai una domanda sul report",
        placeholder="Esempio: la copertura HSE è sufficiente?",
        key="advisor_question"
    )
    if st.button("Chiedi all'Advisor", key="advisor_button"):
        if question.strip():
            st.session_state.advisor_answer = local_hse_advisor(
                question, data, risk, level, actions_df, technical_data
            )
        else:
            st.session_state.advisor_answer = "Scrivi una domanda."
    if st.session_state.advisor_answer:
        st.markdown(st.session_state.advisor_answer)

with tab8:
    st.subheader("HSE Scorecard")
    st.metric("Valutazione complessiva", f"{scorecard_overall:.1f} / 10")
    for _, row in scorecard_df.iterrows():
        col_score, col_reason = st.columns([1, 2.3])
        with col_score:
            st.markdown(f"### {row['Area']}")
            st.markdown(f"## {row['Stelle']}")
            st.markdown(f"**{float(row['Punteggio']):.1f} / 10**")
        with col_reason:
            st.write(row["Motivazione"])
            st.progress(float(row["Punteggio"]) / 10)
        st.divider()

with tab9:
    st.subheader("Generatore Checklist HSE")
    st.info(
        "La checklist non viene compilata nel tool. Seleziona i capitoli e scarica il format Word da inviare all'HSE di cantiere."
    )

    suggested_sections = ["Controlli generali pre-attività"]
    suggested_sections.extend(
        ACTIVITY_TO_CHECKLIST[activity]
        for activity in data["activities"]
        if activity in ACTIVITY_TO_CHECKLIST
    )
    suggested_sections = list(dict.fromkeys(suggested_sections))

    selected_checklists = st.multiselect(
        "Capitoli da inserire nel documento",
        options=list(CHECKLIST_LIBRARY.keys()),
        default=suggested_sections,
        help="Le attività già selezionate nel calcolo vengono proposte automaticamente, ma puoi aggiungere o rimuovere capitoli."
    )

    if selected_checklists:
        total_questions = sum(len(CHECKLIST_LIBRARY[item]) for item in selected_checklists)
        st.caption(f"Documento composto da {len(selected_checklists)} capitoli e {total_questions} verifiche.")
        checklist_buffer = generate_checklist_word(data, selected_checklists)
        st.download_button(
            label="📥 Scarica Checklist Word",
            data=checklist_buffer,
            file_name=f"Checklist_HSE_{safe_filename(data['nome'])}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )
    else:
        st.warning("Seleziona almeno un capitolo da inserire nella checklist.")

with tab10:
    st.subheader("Download report")
    st.caption("La sezione Export contiene gli output complessivi del Risk Tool. La checklist Word rimane disponibile nell'HSE Toolkit.")

    ppt = generate_ppt(
        data,
        risk,
        level,
        driver_df,
        action_plan_df,
        summary,
        root_cause,
        toolbox,
        technical_data,
        scorecard_df,
        scorecard_overall
    )
    ppt_buffer = BytesIO()
    ppt.save(ppt_buffer)
    ppt_buffer.seek(0)

    excel_buffer = generate_excel(
        data,
        risk,
        level,
        driver_df,
        actions_df,
        action_plan_df,
        explanations_df,
        summary,
        root_cause,
        toolbox,
        reduce_risk,
        technical_data,
        scorecard_df,
        scorecard_overall
    )

    filename = safe_filename(data["nome"])
    col_a, col_b = st.columns(2)

    with col_a:
        st.download_button(
            label="📥 Scarica PowerPoint",
            data=ppt_buffer,
            file_name=f"HSE_Risk_Report_{filename}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

    with col_b:
        st.download_button(
            label="📥 Scarica Excel",
            data=excel_buffer,
            file_name=f"HSE_Risk_Report_{filename}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )
